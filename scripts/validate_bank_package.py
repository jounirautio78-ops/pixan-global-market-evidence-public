#!/usr/bin/env python3
"""Fail-closed validation for the public, generated lender package."""

from __future__ import annotations

import copy
import csv
import hashlib
import json
import re
import sys
import zipfile
from datetime import datetime
from decimal import Decimal, InvalidOperation, ROUND_HALF_UP
from pathlib import Path
from typing import Any
from urllib.parse import parse_qsl, urlparse
from xml.etree import ElementTree as ET
from zoneinfo import ZoneInfo

from openpyxl import load_workbook
from pptx import Presentation

try:
    from bank_register_parity import validate_register_parity
except ModuleNotFoundError:  # Support importing this file as scripts.validate_bank_package.
    from scripts.bank_register_parity import validate_register_parity


ROOT = Path(__file__).resolve().parents[1]
MANIFEST_PATH = ROOT / "site" / "data" / "bank-package-manifest.json"
LOCK_PATH = ROOT / "source" / "bank-package-en-lock.json"
CHANGELOG_PATH = ROOT / "site" / "data" / "changelog.json"
REGISTER_CSV_PATH = ROOT / "site" / "data" / "bank-evidence-register.csv"
EN_REGISTER_CSV_PATH = ROOT / "site" / "data" / "bank-evidence-register-en.csv"
MARKET_VALUES_PATH = ROOT / "site" / "data" / "market-values.json"
GLOBAL_BASE_PATH = ROOT / "site" / "data" / "global-base-layer.json"
VENDOR_RESPONSE_CONTROL_PATH = ROOT / "site" / "data" / "vendor-response-control.json"
SOURCE_VENDOR_RESPONSE_CONTROL_PATH = ROOT / "source" / "vendor-response-control.json"
COUNTRY_SCENARIOS_PATH = ROOT / "site" / "data" / "country-scenarios.json"
PUBLIC_FX_PATH = ROOT / "site" / "data" / "fx-rates.json"
SOURCE_FX_PATH = ROOT / "source" / "fx-rates.json"
PUBLIC_FX_SCHEMA_PATH = ROOT / "site" / "schemas" / "fx-rates.schema.json"
SOURCE_FX_SCHEMA_PATH = ROOT / "source" / "schemas" / "fx-rates.schema.json"
ARTIFACT_BUILDER_PATH = ROOT / "scripts" / "artifact-build" / "build_bank_package_artifacts.mjs"
RELEASE_ID = "2026-08-03-germany-vendor-audit-v43"
RELEASE_VERSION = "2026.08.03-43"
RELEASE_DATE = "2026-08-03"
LOCK_RELATIVE_PATH = "source/bank-package-en-lock.json"
# The reviewed v43 lock is written only by the once-daily artifact build. Keep
# the last reviewed hash here until that build has completed, then replace it
# with the new lock SHA-256 before the release validator is run.
EXPECTED_LOCK_SHA256 = "7264671ca143797d2db3a3d4f667918774dfe1eee4303478abe30c824c4bd35d"
PACKAGE_TIME_ZONE = "Asia/Nicosia"
EXPECTED_PACKAGE_CADENCE = {
    "frequency": "once_daily",
    "timeZone": PACKAGE_TIME_ZONE,
    "dashboardMayUpdateIntraday": True,
}
FHM_SOURCE_ID = "SE-FHM-PUBLIC-RECORD-RESPONSE-2026-07-24"
FHM_SOURCE_URL = (
    "https://www.folkhalsomyndigheten.se/regler-och-tillsyn/"
    "tobak-och-nikotinprodukter-regler-for-tillverkning-handel-och-hantering/"
    "elektroniska-cigaretter-och-pafyllningsbehallare-sa-foljer-du-reglerna/"
)
SWEDEN_STRUCTURE_BASIS = "official_registration_structure_count_not_sales_or_market_value"
SWEDEN_STRUCTURE_METRICS = {
    "reporting_entities_count": "REPORTING-ENTITIES",
    "notified_products_count": "NOTIFIED-PRODUCTS",
    "active_products_count": "ACTIVE-PRODUCTS",
    "withdrawn_products_count": "WITHDRAWN-PRODUCTS",
}
EXPECTED_MARKET_OBSERVATIONS = 174
EXPECTED_MARKET_SOURCES = 54
EXPECTED_OFFICIAL_OBSERVATIONS = 152
EXPECTED_OFFICIAL_MARKET_MEASURES = 116
EXPECTED_SWEDEN_STRUCTURE_COUNTS = 36

REGISTER_HEADERS = [
    "Väite",
    "Dia/osio",
    "Todiste",
    "Lähde",
    "Päivämäärä",
    "Laskentatapa",
    "Oletukset",
    "Luottamustaso",
    "Puutteet / tarvittava lisänäyttö",
]
ALLOWED_STATUSES = {"Vahvistettu", "Tuettu", "Oletus", "Puuttuu"}
EN_REGISTER_HEADERS = [
    "Claim",
    "Slide/section",
    "Evidence",
    "Source",
    "Date",
    "Calculation method",
    "Assumptions",
    "Confidence",
    "Gaps / additional evidence needed",
]
EN_ALLOWED_STATUSES = {"Confirmed", "Supported", "Assumption", "Missing"}
EUR_EQUIVALENT_HEADERS = {
    "fi": [
        "Tietuetyyppi",
        "Tunniste",
        "Erä / komponentti",
        "Maa / maantiede",
        "Vuosi",
        "Periodi",
        "Alkuperäinen määrä",
        "Valuutta",
        "ECB-kurssi (valuuttayksikköä / EUR)",
        "EUR-vasta-arvo (täysi tarkkuus)",
        "Rate ID",
        "ECB-lähde URL",
        "Tila",
        "Syy / menetelmä",
    ],
    "en": [
        "Record type",
        "Record ID",
        "Item / component",
        "Country / geography",
        "Year",
        "Period",
        "Original amount",
        "Currency",
        "ECB rate (currency units / EUR)",
        "EUR equivalent (full precision)",
        "Rate ID",
        "ECB source URL",
        "Status",
        "Reason / method",
    ],
}
EUR_EQUIVALENT_SHEET_NAMES = {"fi": "Eurovastineet", "en": "EUR equivalents"}
EXPECTED_LOCKED_EUR_EQUIVALENT_ROWS = 84
EXPECTED_LOCKED_EUR_STATUS_COUNTS = {
    "computed": 48,
    "already_eur": 29,
    "not_computed": 7,
}
EXPECTED_TEMPLATE_INPUTS = {
    "scripts/artifact-build/seeds/v17/pixan-bank-deck-short-en.pptx",
    "scripts/artifact-build/seeds/v17/pixan-bank-deck-large-en.pptx",
    "scripts/artifact-build/seeds/v17/pixan-bank-evidence-register-en.xlsx",
    "scripts/artifact-build/seeds/v17/pixan-bank-deck-short-fi.pptx",
    "scripts/artifact-build/seeds/v17/pixan-bank-deck-large-fi.pptx",
    "scripts/artifact-build/seeds/v17/pixan-bank-evidence-register-fi.xlsx",
}
EXPECTED_INPUTS = {
    "scripts/artifact-build/build_bank_package_artifacts.mjs",
    "scripts/analyze_nz_2024_returns.py",
    "scripts/build_global_base.py",
    "scripts/refresh_global_base.py",
    "scripts/build_vendor_response_control.py",
    *EXPECTED_TEMPLATE_INPUTS,
    "site/data/atlas.json",
    "site/data/changelog.json",
    "site/data/country-scenarios.json",
    "site/data/donor-cockpit.json",
    "site/data/third-donor-screen.json",
    "site/data/evidence-lanes.json",
    "site/data/fx-rates.json",
    "site/data/global-base-layer.json",
    "site/data/global-base-layer.csv",
    "site/data/market-values.json",
    "site/data/patent-history.json",
    "site/data/vendor-response-control.json",
    "site/data/vendor-response-control.csv",
    "site/schemas/fx-rates.schema.json",
    "site/schemas/global-base-layer.schema.json",
    "site/schemas/third-donor-screen.schema.json",
    "source/bank-evidence-register-en.json",
    "source/bank-package-en-lock.json",
    "source/fx-rates.json",
    "source/global-base-config.json",
    "source/global-base-observations.json",
    "source/country-method-route-config.json",
    "source/COUNTRY_METHOD_ROUTE_MAP.md",
    "source/FIVE_COUNTRY_METHOD_SPRINT_2026-07-27.md",
    "source/ITALY_ADM_RESPONSE_BOUNDARY_2026-07-24.md",
    "source/POLAND_EUCEG_ANNUAL_SALES_REQUEST_2026-07-28.md",
    "source/top20-data-request-routes.json",
    "source/paid-data-procurement.json",
    "source/third-donor-screen.json",
    "source/vendor-response-control.json",
    "source/schemas/fx-rates.schema.json",
    "source/schemas/global-base-layer.schema.json",
    "source/schemas/third-donor-screen.schema.json",
    "source/NZ_2023_ANNUAL_RETURNS_FAIL_CLOSED.md",
    "source/NZ_2024_ANNUAL_RETURNS_RECONCILIATION.md",
    "source/NZ_2024_DONOR_CLOSURE_PACK.md",
    "source/NZ_2024_D8_D10_OFFICIAL_SOURCE_AUDIT.md",
    "source/NZ_DONOR_FOLLOWUP_PACK_2026-08-07.md",
    "source/NZ_2024_PRODUCT_SCOPE_AUDIT.json",
    "source/NZ_2024_RPS_RETAIL_VALUE_SENSITIVITY.md",
    "source/NZ_2024_WORKBOOK_MANIFEST.json",
    "source/CANADA_RCS_2019_2025_RETAIL_SALES.md",
    "source/CANADA_RCS_TAX_BASIS_CLARIFICATION_2026-07-29.md",
    "source/CANADA_RCS_SCOPE_QUALITY_CLARIFICATION_2026-07-30.md",
    "source/CANADA_2024_DONOR_CLOSURE_PACK.md",
    "source/CANADA_2024_D5_D7_D10_OFFICIAL_SOURCE_AUDIT.md",
    "source/CANADA_INDEPENDENT_D5_D7_D10_ROUTE_MAP_2026-07-31.md",
    "source/NZ_CA_DE_DONOR_CONTROL_SPRINT_2026-08-02.md",
    "source/GERMANY_VENDOR_AUDIT_BOUNDARY_2026-08-03.md",
    "source/THIRD_DONOR_SCREEN_2026-07-27.md",
    "source/POLAND_2020_2025_RECONSTRUCTION.md",
    "source/POLAND_D1_D10_PREASSESSMENT_2026-07-31.md",
    "source/FOLLOW_UP_DRAFTS_2026-07-28.md",
    "source/US_FTC_2015_2021_REPORTED_SALES.md",
    "source/SWEDEN_FHM_REGISTRATION_STRUCTURE_2018_2026.md",
}
EXPECTED_ARTIFACTS = {
    "short-deck-fi": {
        "kind": "pptx",
        "language": "fi",
        "path": "downloads/pixan-bank-deck-short-fi.pptx",
        "slideCount": 6,
    },
    "large-deck-fi": {
        "kind": "pptx",
        "language": "fi",
        "path": "downloads/pixan-bank-deck-large-fi.pptx",
        "slideCount": 30,
    },
    "evidence-register-fi": {
        "kind": "xlsx",
        "language": "fi",
        "path": "downloads/pixan-bank-evidence-register-fi.xlsx",
    },
    "short-deck-en": {
        "kind": "pptx",
        "language": "en",
        "path": "downloads/pixan-bank-deck-short-en.pptx",
        "slideCount": 6,
    },
    "large-deck-en": {
        "kind": "pptx",
        "language": "en",
        "path": "downloads/pixan-bank-deck-large-en.pptx",
        "slideCount": 30,
    },
    "evidence-register-en": {
        "kind": "xlsx",
        "language": "en",
        "path": "downloads/pixan-bank-evidence-register-en.xlsx",
    },
}
FORBIDDEN_TEXT = (
    "/users/",
    "\\users\\",
    "file://",
    "tmp/pdfs",
)
PRIVATE_IDENTIFIER_FINGERPRINTS = frozenset(
    {
        (7, "46d7415f6182ece9e933e8e9f780957e449361e0dbe10e34f46c186cad3382a1"),
        (7, "f910f0bbe95037851d18ca33b91ee7fc9f334c6cfcd02deaf66af4501c8a884c"),
        (9, "7e6578c2e34b53136741c6efe7799a2dce739651c22404a7894b48d42aa88b41"),
        (13, "933536a17b00f1b39ba9d3585427bd7232d44960ab35754318c1da8e4cf6c5be"),
        (25, "40f45830e7e3e21d88245728fe87f76b2e8919543a502aad248a465487cacee3"),
    }
)
FORBIDDEN_ARCHIVE_PARTS = (
    "vbaproject",
    "/embeddings/",
    "/externalLinks/",
    "/oleObject",
    "/comments",
    "connections.xml",
)
SENSITIVE_QUERY_KEYS = {
    "access_token",
    "api_key",
    "apikey",
    "authorization",
    "key",
    "password",
    "secret",
    "sig",
    "signature",
    "token",
    "x-amz-credential",
    "x-amz-signature",
}
SHA256_RE = re.compile(r"^[0-9a-f]{64}$")


def parse_release_timestamp(value: Any, label: str, errors: list[str]) -> datetime | None:
    if not isinstance(value, str):
        errors.append(f"{label} must be an ISO-8601 timestamp with an explicit offset")
        return None
    try:
        parsed = datetime.fromisoformat(value.replace("Z", "+00:00"))
    except ValueError:
        errors.append(f"{label} must be a valid ISO-8601 timestamp")
        return None
    if parsed.tzinfo is None or parsed.utcoffset() is None:
        errors.append(f"{label} must include an explicit UTC offset")
        return None
    return parsed


def validate_daily_package_snapshot(
    manifest: dict[str, Any],
    changelog: dict[str, Any],
    errors: list[str],
) -> bool:
    """Return whether reviewed-input hash drift is allowed for this package snapshot."""

    cadence = manifest.get("cadence")
    if cadence != EXPECTED_PACKAGE_CADENCE:
        errors.append(
            "manifest cadence must be exactly once_daily in Asia/Nicosia and allow intraday dashboard updates"
        )

    releases = changelog.get("releases")
    if not isinstance(releases, list) or not releases or not all(
        isinstance(release, dict) for release in releases
    ):
        errors.append("public changelog must contain at least one release object")
        return False

    manifest_release = manifest.get("release")
    if not isinstance(manifest_release, dict) or set(manifest_release) != {
        "id",
        "version",
        "publishedAt",
    }:
        errors.append("manifest release must contain exactly id, version and publishedAt")
        return False

    matching_indexes = [
        index
        for index, release in enumerate(releases)
        if {
            key: release.get(key) for key in ("id", "version", "publishedAt")
        }
        == manifest_release
    ]
    if len(matching_indexes) != 1:
        errors.append("manifest release must match exactly one public changelog release")
        return False

    latest_release = {
        key: releases[0].get(key) for key in ("id", "version", "publishedAt")
    }
    manifest_timestamp = parse_release_timestamp(
        manifest_release.get("publishedAt"),
        "manifest release publishedAt",
        errors,
    )
    latest_timestamp = parse_release_timestamp(
        latest_release.get("publishedAt"),
        "latest changelog release publishedAt",
        errors,
    )
    if manifest_timestamp is None or latest_timestamp is None:
        return False

    package_date = manifest_timestamp.astimezone(ZoneInfo(PACKAGE_TIME_ZONE)).date()
    dashboard_date = latest_timestamp.astimezone(ZoneInfo(PACKAGE_TIME_ZONE)).date()
    if package_date != dashboard_date:
        errors.append(
            "bank package is older than the current dashboard calendar date in Asia/Nicosia"
        )
        return False
    if manifest.get("asOf") != changelog.get("asOf") or manifest.get("asOf") != package_date.isoformat():
        errors.append(
            "manifest asOf, changelog asOf and the Asia/Nicosia package calendar date must match"
        )
        return False

    release_index = matching_indexes[0]
    if release_index == 0:
        return False
    if manifest_timestamp >= latest_timestamp:
        errors.append(
            "an earlier changelog package snapshot must precede the latest same-day release timestamp"
        )
        return False
    return cadence == EXPECTED_PACKAGE_CADENCE


def validate_release_lock(
    manifest: dict[str, Any],
    lock: dict[str, Any],
    errors: list[str],
) -> bool:
    """Bind an earlier same-day package to its immutable reviewed snapshot."""

    valid = True
    expected_lock_keys = {
        "schemaVersion",
        "release",
        "asOf",
        "reviewedInputs",
        "artifacts",
        "generatedBy",
    }
    if not isinstance(lock, dict) or set(lock) != expected_lock_keys:
        errors.append("bank-package lock has an unexpected schema")
        return False
    if sha256(LOCK_PATH) != EXPECTED_LOCK_SHA256:
        errors.append("bank-package lock SHA-256 differs from the reviewed package snapshot")
        valid = False
    if (
        lock.get("schemaVersion") != 2
        or lock.get("release") != manifest.get("release")
        or lock.get("asOf") != manifest.get("asOf")
    ):
        errors.append("bank-package lock release/asOf differs from the manifest snapshot")
        valid = False

    manifest_inputs = manifest.get("inputs")
    locked_inputs = lock.get("reviewedInputs")
    if not isinstance(manifest_inputs, list) or not isinstance(locked_inputs, list):
        errors.append("bank-package manifest/lock reviewed inputs must be arrays")
        valid = False
    else:
        self_lock_entries = [
            item
            for item in manifest_inputs
            if isinstance(item, dict) and item.get("path") == LOCK_RELATIVE_PATH
        ]
        if self_lock_entries != [{"path": LOCK_RELATIVE_PATH, "sha256": EXPECTED_LOCK_SHA256}]:
            errors.append("manifest must bind exactly one reviewed package lock SHA-256")
            valid = False
        expected_locked_inputs = [
            item
            for item in manifest_inputs
            if not (isinstance(item, dict) and item.get("path") == LOCK_RELATIVE_PATH)
        ]
        if locked_inputs != expected_locked_inputs:
            errors.append("bank-package lock reviewedInputs differ from the manifest snapshot")
            valid = False

    manifest_artifacts = manifest.get("artifacts")
    locked_artifacts = lock.get("artifacts")
    if not isinstance(manifest_artifacts, list) or not isinstance(locked_artifacts, list):
        errors.append("bank-package manifest/lock artifacts must be arrays")
        valid = False
    else:
        expected_locked_artifacts: list[dict[str, Any]] = []
        for item in manifest_artifacts:
            if not isinstance(item, dict):
                expected_locked_artifacts.append({})
                continue
            count_key = "slideCount" if item.get("kind") == "pptx" else "rowCount"
            projection = {
                key: item.get(key)
                for key in ("id", "kind", "language", "sha256", "bytes", count_key)
            }
            projection["path"] = f"site/{item.get('path')}"
            expected_locked_artifacts.append(projection)
        if locked_artifacts != expected_locked_artifacts:
            errors.append("bank-package lock artifacts differ from the manifest snapshot")
            valid = False
    return valid


def validate_v22_market_bindings(errors: list[str]) -> None:
    try:
        market = load_json(MARKET_VALUES_PATH)
    except ValueError as error:
        errors.append(str(error))
        return
    observations = market.get("observations")
    sources = market.get("sources")
    if not isinstance(observations, list) or len(observations) != EXPECTED_MARKET_OBSERVATIONS:
        errors.append(
            f"current bank-package input requires exactly {EXPECTED_MARKET_OBSERVATIONS} market observations"
        )
        return
    if not isinstance(sources, list) or len(sources) != EXPECTED_MARKET_SOURCES:
        errors.append(f"current bank-package input requires exactly {EXPECTED_MARKET_SOURCES} market sources")
        sources = [] if not isinstance(sources, list) else sources
    source_by_id = {
        item.get("sourceId"): item
        for item in sources
        if isinstance(item, dict) and isinstance(item.get("sourceId"), str)
    }
    if len(source_by_id) != len(sources):
        errors.append("market sources must have unique string sourceId values")
    fhm_source = source_by_id.get(FHM_SOURCE_ID)
    if not isinstance(fhm_source, dict) or fhm_source.get("pageUrl") != FHM_SOURCE_URL:
        errors.append("v27 market sources must retain the reviewed public FHM reference")
    for source_id in {
        "NZ-MOH-ANNUAL-RETURNS-2024",
        "NZ-MOH-ANNUAL-RETURN-REQUIREMENTS",
        "NZ-MOH-ANNUAL-RETURNS-2024-GUIDE",
    }:
        source = source_by_id.get(source_id)
        if not isinstance(source, dict) or source.get("publisher") != "New Zealand Ministry of Health":
            errors.append(f"v27 market sources lack reviewed New Zealand source {source_id}")
    observation_by_id = {
        item.get("observationId"): item
        for item in observations
        if isinstance(item, dict) and isinstance(item.get("observationId"), str)
    }
    if len(observation_by_id) != len(observations):
        errors.append("market observations must have unique string observationId values")
        return
    official = [
        item
        for item in observations
        if str(item.get("evidenceStatus", "")).startswith("official")
    ]
    sweden_structure = [
        item
        for item in observations
        if item.get("marketValueBasis") == SWEDEN_STRUCTURE_BASIS
        or FHM_SOURCE_ID in item.get("sourceIds", [])
    ]
    sweden_structure_ids = {
        item.get("observationId") for item in sweden_structure if isinstance(item, dict)
    }
    official_market_measures = [
        item for item in official if item.get("observationId") not in sweden_structure_ids
    ]
    if len(official) != EXPECTED_OFFICIAL_OBSERVATIONS:
        errors.append(
            f"current bank-package input requires exactly {EXPECTED_OFFICIAL_OBSERVATIONS} official observations"
        )
    if (
        len(official_market_measures) != EXPECTED_OFFICIAL_MARKET_MEASURES
        or {item.get("countryIso2") for item in official_market_measures} != {
        "CA", "DE", "ES", "FI", "JP", "NZ", "PL", "SE", "US"
        }
    ):
        errors.append(
            f"current bank-package input requires {EXPECTED_OFFICIAL_MARKET_MEASURES} official market measures across nine reviewed countries"
        )

    expected_structure_ids = {
        f"SE-{year}-FHM-{suffix}"
        for year in range(2018, 2027)
        for suffix in SWEDEN_STRUCTURE_METRICS.values()
    }
    if (
        len(sweden_structure) != EXPECTED_SWEDEN_STRUCTURE_COUNTS
        or sweden_structure_ids != expected_structure_ids
    ):
        errors.append(
            "v27 bank package requires the exact 36-record Swedish FHM structure series"
        )
    for item in sweden_structure:
        value = item.get("value")
        metric = item.get("metric")
        snapshot = item.get("year") == 2026
        expected_unit = "reporting_entity" if metric == "reporting_entities_count" else "product"
        if (
            item.get("countryIso2") != "SE"
            or item.get("geography") != "Sweden"
            or metric not in SWEDEN_STRUCTURE_METRICS
            or isinstance(value, bool)
            or not isinstance(value, (int, float))
            or not float(value).is_integer()
            or value < 0
            or item.get("unit") != expected_unit
            or item.get("currency") is not None
            or item.get("period")
            != (
                "current_snapshot_as_of_2026_07_24"
                if snapshot
                else "authority_supplied_year_label"
            )
            or item.get("finality")
            != (
                "official_current_snapshot"
                if snapshot
                else "official_response_year_label"
            )
            or item.get("marketValueBasis") != SWEDEN_STRUCTURE_BASIS
            or item.get("comparableMarketValue") is not False
            or item.get("atlasEstimate") is not False
            or not str(item.get("evidenceStatus", "")).startswith("official")
            or item.get("sourceIds") != [FHM_SOURCE_ID]
        ):
            errors.append(
                "Swedish FHM structure record is not a non-sales count: "
                f"{item.get('observationId')!r}"
            )
        if snapshot:
            snapshot_text = " ".join(
                str(item.get(field, ""))
                for field in ("period", "finality", "limitationEn", "limitationFi")
            ).casefold()
            if (
                item.get("period") == "calendar_year"
                or ("snapshot" not in snapshot_text and "tilannekuva" not in snapshot_text)
            ):
                errors.append(
                    f"{item.get('observationId')}: 2026 must be disclosed as a snapshot, not a full year"
                )

    exact_observations = {
        "NZ-2024-SPECIALIST-RETAIL-SALES-LOWER-BOUND": (
            280_000_000,
            "official_provisional",
            "official_lower_bound_with_quality_warning",
        ),
        "NZ-2024-SPECIALIST-RETAIL-PRODUCT-SALES-RAW-FILE-SUM": (
            280_684_512.81,
            "derived_official_files",
            "reproduced_raw_file_sum_with_quality_warning",
        ),
        "NZ-2024-IDENTIFIED-VAPING-PRODUCT-SALES-RAW-SUM": (
            274_180_410.21,
            "derived_official_files",
            "keyword_classified_raw_file_sum_with_quality_warning",
        ),
        "EU-2023-EC-E-CIGARETTE-MARKET-BENCHMARK": (
            4_990_000_000,
            "institutional_supported",
            "published_secondary_benchmark",
        ),
        "US-2021-FTC-CARTRIDGE-DISPOSABLE-REPORTED-SALES": (
            2_763_284_338,
            "official_table_derived",
            "official_table_sum",
        ),
    }
    for observation_id, (value, status, finality) in exact_observations.items():
        item = observation_by_id.get(observation_id)
        if (
            not isinstance(item, dict)
            or item.get("value") != value
            or item.get("evidenceStatus") != status
            or item.get("finality") != finality
            or item.get("comparableMarketValue") is not False
            or item.get("atlasEstimate") is not False
        ):
            errors.append(f"v27 reviewed observation binding differs: {observation_id}")

    protocol = market.get("donorProtocol")
    criteria = protocol.get("criteria") if isinstance(protocol, dict) else None
    candidates = market.get("donorCandidates")
    readiness = market.get("meta", {}).get("modelReadiness", {})
    if (
        not isinstance(criteria, list)
        or [item.get("criterionId") for item in criteria] != [f"D{index}" for index in range(1, 11)]
    ):
        errors.append("donor protocol must contain ordered criteria D1-D10")
    if (
        not isinstance(candidates, list)
        or len(candidates) != 5
        or any(item.get("decision") != "not_accepted" for item in candidates)
    ):
        errors.append("all five reviewed donor candidates must remain not accepted")
    candidate_by_id = {
        item.get("candidateId"): item
        for item in candidates or []
        if isinstance(item, dict) and isinstance(item.get("candidateId"), str)
    }
    nz_candidate = candidate_by_id.get("NZ-2024-IDENTIFIED-VAPING-RETAIL-SUBTOTAL")
    if (
        not isinstance(nz_candidate, dict)
        or nz_candidate.get("referenceType") != "observation"
        or nz_candidate.get("referenceId")
        != "NZ-2024-IDENTIFIED-VAPING-PRODUCT-SALES-RAW-SUM"
        or nz_candidate.get("passedCriteria")
        != ["D1", "D2", "D3", "D4", "D6", "D7", "D9"]
        or nz_candidate.get("failedCriteria") != ["D5"]
        or nz_candidate.get("openCriteria") != ["D8", "D10"]
    ):
        errors.append("New Zealand v27 donor candidate must remain 7/10 and not accepted")
    if (
        readiness.get("comparableFullYearMarketValueDonors") != 0
        or readiness.get("minimumRequiredDonors") != 3
    ):
        errors.append("accepted-donor gate must remain 0/3")


def validate_v43_vendor_and_global_boundary(errors: list[str]) -> None:
    """Validate the privacy-safe Germany result without reading licensed data."""

    try:
        vendor_control = load_json(VENDOR_RESPONSE_CONTROL_PATH)
        source_vendor_control = load_json(SOURCE_VENDOR_RESPONSE_CONTROL_PATH)
        market = load_json(MARKET_VALUES_PATH)
        global_base = load_json(GLOBAL_BASE_PATH)
    except ValueError as error:
        errors.append(str(error))
        return

    public_source_projection = copy.deepcopy(vendor_control)
    public_source_projection.pop("summary", None)
    for vendor in public_source_projection.get("vendors", []):
        if not isinstance(vendor, dict):
            continue
        vendor.pop("evidenceReceivedCount", None)
        vendor.pop("evaluatedGateCount", None)
        vendor.pop("mandatoryGatePassCount", None)
    if public_source_projection != source_vendor_control:
        errors.append("public vendor-response control differs from the reviewed source")

    benchmark = vendor_control.get("germanyBenchmark")
    vendors = vendor_control.get("vendors")
    euromonitor = next(
        (
            item
            for item in vendors or []
            if isinstance(item, dict)
            and item.get("vendorId") == "euromonitor-passport-nicotine"
        ),
        None,
    )
    expected_gate_results = {
        "G1": "pass",
        "G2": "fail",
        "G3": "fail",
        "G4": "not_testable",
        "G5": "fail",
        "G6": "fail",
    }
    gate_results = (euromonitor or {}).get("gateResults")
    public_status_en = str((euromonitor or {}).get("publicStatusEn", ""))
    public_status_fi = str((euromonitor or {}).get("publicStatusFi", ""))
    evaluated_gate_count = sum(
        (gate_results.get(gate_id) or {}).get("status") != "missing"
        for gate_id in expected_gate_results
    ) if isinstance(gate_results, dict) else -1
    mandatory_gate_pass_count = sum(
        (gate_results.get(gate_id) or {}).get("status") == "pass"
        for gate_id in expected_gate_results
    ) if isinstance(gate_results, dict) else -1
    if (
        vendor_control.get("schemaVersion") != 3
        or vendor_control.get("asOf") != RELEASE_DATE
        or vendor_control.get("version") != RELEASE_VERSION
        or vendor_control.get("status")
        != "public_status_only_germany_extract_received_wider_package_not_authorised"
        or not isinstance(benchmark, dict)
        or benchmark.get("benchmarkId") != "DE-BLIND-1.0.0"
        or benchmark.get("status") != "numeric_pass_scope_open"
        or benchmark.get("vendorPassDoesNotEstablishDonorPass") is not True
        or benchmark.get("donorGateEffect") != "none"
        or not isinstance(euromonitor, dict)
        or euromonitor.get("quoteReceived") is not True
        or euromonitor.get("responseState")
        != "evaluation_extract_received_private_audit_complete"
        or euromonitor.get("mandatoryGatePassCount", mandatory_gate_pass_count) != 1
        or euromonitor.get("evaluatedGateCount", evaluated_gate_count) != 6
        or mandatory_gate_pass_count != 1
        or evaluated_gate_count != 6
        or euromonitor.get("scoringState") != "not_scored"
        or euromonitor.get("weightedScore") is not None
        or euromonitor.get("evaluationExtractAuthorised") is not True
        or euromonitor.get("evaluationExtractReceived") is not True
        or euromonitor.get("widerPackagePurchaseAuthorised") is not False
        or (euromonitor.get("receivedEvidence") or {}).get(
            "officialAnchorReconciliation"
        )
        is not True
        or not isinstance(gate_results, dict)
        or any(
            (gate_results.get(gate_id) or {}).get("status") != expected_status
            for gate_id, expected_status in expected_gate_results.items()
        )
        or "full 19-tab Germany evaluation extract" not in public_status_en
        or "numerical liquid-volume proximity tests passed" not in public_status_en
        or "no wider 25/50/78-country subscription is authorised" not in public_status_en
        or "donor gate remains 0/3" not in public_status_en
        or "global value remains not_computed" not in public_status_en
        or "NOT SCORED" not in public_status_en
        or "täysi 19 välilehden arviointiote" not in public_status_fi
        or "numeerinen läheisyystesti läpäistiin" not in public_status_fi
        or "laajempaa 25/50/78 maan tilausta ei ole valtuutettu" not in public_status_fi
        or "donor-portti pysyy 0/3:ssa" not in public_status_fi
        or "maailmanarvo not_computed-tilassa" not in public_status_fi
        or "EI PISTEYTETTY" not in public_status_fi
    ):
        errors.append("v43 Germany extract and 1/6 vendor-gate boundary differs")

    germany_candidate = next(
        (
            item
            for item in market.get("donorCandidates", [])
            if isinstance(item, dict)
            and item.get("candidateId") == "DE-2025-LIQUID-RETAIL-MODEL"
        ),
        None,
    )
    germany_decision_reason_en = str(
        (germany_candidate or {}).get(
            "decisionReasonEn",
            (germany_candidate or {}).get("blockerEn", ""),
        )
    )
    readiness = market.get("meta", {}).get("modelReadiness", {})
    global_retail = global_base.get("globalRetailSales", {})
    method_summary = global_base.get("methodRouteControl", {}).get("summary", {})
    if (
        not isinstance(germany_candidate, dict)
        or germany_candidate.get("decision") != "not_accepted"
        or "Germany remains NOT ACCEPTED" not in germany_decision_reason_en
        or readiness.get("comparableFullYearMarketValueDonors") != 0
        or readiness.get("minimumRequiredDonors") != 3
        or global_retail.get("value") is not None
        or global_retail.get("currency") is not None
        or global_retail.get("eligibleObservationCount") != 0
        or method_summary.get("eligibleForGlobalRollupCount") != 0
        or method_summary.get("donorAcceptedCount") != 0
    ):
        errors.append(
            "v43 must keep Germany outside the donor set, donor gate 0/3 and global value not_computed"
        )


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_json(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as error:
        raise ValueError(f"cannot read {path.relative_to(ROOT)}: {error}") from error


def decimal_value(value: Any) -> Decimal | None:
    if isinstance(value, bool):
        return None
    try:
        result = Decimal(str(value))
    except (InvalidOperation, ValueError):
        return None
    return result if result.is_finite() else None


def validate_artifact_builder_fx_contract(builder_text: str, errors: list[str]) -> None:
    required_tokens = {
        '"site/data/fx-rates.json"': "public FX reviewed input",
        '"site/schemas/fx-rates.schema.json"': "public FX schema reviewed input",
        '"source/fx-rates.json"': "source FX reviewed input",
        '"source/schemas/fx-rates.schema.json"': "source FX schema reviewed input",
        "buildEurEquivalentRows": "data-driven EUR row builder",
        "scenario_component": "country-scenario component rows",
        "market_observation": "market-observation rows",
        '"model"': "market-model rows",
        "compatible_ecb_rate_missing": "missing-rate fail-closed reason",
        '"not_computed"': "missing-rate fail-closed status",
        "eur_equivalent = original_amount / currency_units_per_eur": "reviewed FX formula",
        "`=G${sheetRow}/I${sheetRow}`": "full-precision worksheet formula",
        '"EUR equivalents"': "English EUR-equivalent sheet",
        '"Eurovastineet"': "Finnish EUR-equivalent sheet",
        "[FX methodology]": "deck FX methodology notes",
        "fxSourcesInDeckNotes": "deck-source QA lock",
        "eurEquivalentRowsAfterReopen": "workbook-row QA lock",
        "validateV27MarketEvidence": "v27 market-role validation",
        FHM_SOURCE_ID: "Swedish FHM source binding",
        SWEDEN_STRUCTURE_BASIS: "Swedish non-sales structure-role marker",
        "officialMarketMeasures: 116": "116 official market-measure lock",
        "swedenRegisterStructure: 36": "36 Swedish register-count lock",
        "validateVendorGateBoundary": "v43 Germany vendor-gate validation",
    }
    for token, description in required_tokens.items():
        if token not in builder_text:
            errors.append(f"artifact builder lacks {description}: {token}")
    if re.search(r"=ROUND\(\s*G\$\{sheetRow", builder_text, flags=re.IGNORECASE):
        errors.append("artifact builder must not round EUR-equivalent worksheet formulas")


def validate_fx_artifact_inputs(
    public_fx: dict[str, Any],
    source_fx: dict[str, Any],
    errors: list[str],
) -> None:
    if public_fx != source_fx:
        errors.append("artifact FX input differs between source and site/data")
    if not PUBLIC_FX_SCHEMA_PATH.is_file() or not SOURCE_FX_SCHEMA_PATH.is_file():
        errors.append("artifact FX source and public schema files are both required")
    elif PUBLIC_FX_SCHEMA_PATH.read_bytes() != SOURCE_FX_SCHEMA_PATH.read_bytes():
        errors.append("artifact FX source and public schemas differ")
    policy = public_fx.get("calculationPolicy")
    if (
        public_fx.get("targetCurrency") != "EUR"
        or public_fx.get("provider", {}).get("name") != "European Central Bank"
        or not isinstance(policy, dict)
        or policy.get("formulaMachine")
        != "eur_equivalent = original_amount / currency_units_per_eur"
        or policy.get("missingRateStatus") != "not_computed"
        or policy.get("eligibleUnitRule") != "currency_must_equal_unit"
    ):
        errors.append("artifact FX input does not retain the reviewed ECB conversion policy")
    seen: set[tuple[str, int]] = set()
    for rate in public_fx.get("rates", []):
        if not isinstance(rate, dict):
            errors.append("artifact FX input contains a non-object rate")
            continue
        currency = rate.get("currency")
        year = rate.get("year")
        key = (currency, year)
        expected_id = f"ECB-EXR-A-{currency}-EUR-SP00-A-{year}"
        parsed = urlparse(str(rate.get("sourceUrl", "")))
        if (
            key in seen
            or rate.get("rateId") != expected_id
            or decimal_value(rate.get("currencyUnitsPerEur")) is None
            or decimal_value(rate.get("currencyUnitsPerEur")) <= 0
            or parsed.scheme != "https"
            or parsed.hostname != "data-api.ecb.europa.eu"
        ):
            errors.append(f"artifact FX rate is invalid: {currency}/{year}")
        seen.add(key)


def build_expected_eur_equivalent_rows(
    market: dict[str, Any],
    scenarios: dict[str, Any],
    fx: dict[str, Any],
) -> list[dict[str, Any]]:
    rates = {
        (rate.get("currency"), rate.get("year")): rate
        for rate in fx.get("rates", [])
        if isinstance(rate, dict)
    }
    eligible_periods = set(
        fx.get("calculationPolicy", {}).get("eligibleRecordPeriods", [])
    )
    rows: list[dict[str, Any]] = []

    def append(
        record_type: str,
        record_id: Any,
        item: Any,
        geography: Any,
        record: dict[str, Any],
    ) -> None:
        amount = decimal_value(record.get("value"))
        currency = record.get("currency")
        unit = record.get("unit")
        year = record.get("year")
        period = record.get("period")
        if (
            amount is None
            or amount <= 0
            or not isinstance(currency, str)
            or re.fullmatch(r"[A-Z]{3}", currency) is None
            or unit != currency
        ):
            return
        if currency == "EUR":
            status = "already_eur"
            reason = "original_currency_already_eur"
            rate_value: Decimal | None = Decimal("1")
            rate_id: str | None = "EUR-IDENTITY"
            source_url = fx.get("provider", {}).get("methodologyUrl")
        elif not isinstance(year, int) or isinstance(year, bool) or period not in eligible_periods:
            status = "not_computed"
            reason = "period_not_compatible_with_annual_average"
            rate_value = None
            rate_id = None
            source_url = fx.get("provider", {}).get("datasetUrl")
        else:
            rate = rates.get((currency, year))
            rate_value = decimal_value(rate.get("currencyUnitsPerEur")) if rate else None
            if rate_value is None or rate_value <= 0:
                status = "not_computed"
                reason = "compatible_ecb_rate_missing"
                rate_value = None
                rate_id = None
                source_url = fx.get("provider", {}).get("datasetUrl")
            else:
                status = "computed"
                reason = "original_amount_divided_by_ecb_annual_average"
                rate_id = rate.get("rateId")
                source_url = rate.get("sourceUrl")
        rows.append(
            {
                "recordType": record_type,
                "recordId": record_id,
                "item": item,
                "geography": geography,
                "year": year,
                "period": period,
                "originalAmount": amount,
                "currency": currency,
                "rateValue": rate_value,
                "rateId": rate_id,
                "sourceUrl": source_url,
                "status": status,
                "reason": reason,
            }
        )

    for observation in market.get("observations", []):
        if isinstance(observation, dict):
            append(
                "market_observation",
                observation.get("observationId"),
                observation.get("metric"),
                observation.get("geography"),
                observation,
            )

    for scenario in scenarios.get("countryYearScenarios", []):
        if not isinstance(scenario, dict):
            continue
        components = scenario.get("componentBreakdown")
        if not isinstance(components, dict):
            continue
        for range_key, component in components.items():
            if not isinstance(component, dict):
                continue
            for component_key, value in component.items():
                numeric = decimal_value(value)
                if numeric is None or numeric <= 0:
                    continue
                append(
                    "scenario_component",
                    scenario.get("scenarioId"),
                    f"{range_key}.{component_key}",
                    scenario.get("geography"),
                    {
                        "value": value,
                        "currency": scenario.get("currency"),
                        "unit": scenario.get("currency"),
                        "year": scenario.get("year"),
                        "period": "calendar_year",
                    },
                )

    for model in market.get("models", []):
        if not isinstance(model, dict):
            continue
        for bound in ("low", "base", "central", "high"):
            numeric = decimal_value(model.get(bound))
            if numeric is None or numeric <= 0:
                continue
            append(
                "model",
                model.get("modelId"),
                bound,
                model.get("geography"),
                {
                    "value": model.get(bound),
                    "currency": model.get("currency"),
                    "unit": model.get("currency"),
                    "year": model.get("year"),
                    "period": "calendar_year",
                },
            )
    return rows


def load_expected_eur_equivalent_rows(errors: list[str]) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    try:
        market = load_json(MARKET_VALUES_PATH)
        scenarios = load_json(COUNTRY_SCENARIOS_PATH)
        public_fx = load_json(PUBLIC_FX_PATH)
        source_fx = load_json(SOURCE_FX_PATH)
    except ValueError as error:
        errors.append(str(error))
        return [], {}
    validate_fx_artifact_inputs(public_fx, source_fx, errors)
    rows = build_expected_eur_equivalent_rows(market, scenarios, public_fx)
    if not rows:
        errors.append("artifact EUR-equivalent ledger has no eligible rows")
    return rows, public_fx


def deck_fx_markers(
    rows: list[dict[str, Any]],
    language: str,
) -> tuple[str, str, str, str, str]:
    by_key = {
        (row["recordType"], row["recordId"], row["item"]): row
        for row in rows
    }
    nz_low = by_key.get(
        ("scenario_component", "NZ-2024-RETAIL-RANGE", "low.combinedNzd")
    )
    nz_high = by_key.get(
        ("scenario_component", "NZ-2024-RETAIL-RANGE", "high.combinedNzd")
    )
    nz_observed = by_key.get(
        (
            "market_observation",
            "NZ-2024-IDENTIFIED-VAPING-PRODUCT-SALES-RAW-SUM",
            "derived_identified_vaping_product_sales_raw_sum",
        )
    )
    ftc = by_key.get(
        (
            "market_observation",
            "US-2021-FTC-CARTRIDGE-DISPOSABLE-REPORTED-SALES",
            "ftc_reported_cartridge_and_disposable_sales",
        )
    )
    canada_retail = by_key.get(
        (
            "market_observation",
            "CA-2024-STATCAN-RCS-VAPING-RETAIL-SALES",
            "statcan_rcs_vaping_retail_sales",
        )
    )
    canada_shipments = by_key.get(
        (
            "market_observation",
            "CA-2024-MANUFACTURER-IMPORTER-SHIPMENTS-VALUE",
            "manufacturer_importer_shipments_value",
        )
    )
    if not all(
        item
        and item["status"] == "computed"
        and isinstance(item["rateValue"], Decimal)
        and item["rateValue"] > 0
        for item in (nz_observed, nz_low, nz_high, ftc, canada_retail, canada_shipments)
    ):
        return (
            "eur not_computed",
            "eur not_computed",
            "eur not_computed",
            "eur not_computed",
            "eur not_computed",
        )
    nz_observed_eur = nz_observed["originalAmount"] / nz_observed["rateValue"]
    nz_low_eur = nz_low["originalAmount"] / nz_low["rateValue"]
    nz_high_eur = nz_high["originalAmount"] / nz_high["rateValue"]
    ftc_eur = ftc["originalAmount"] / ftc["rateValue"]
    canada_retail_eur = canada_retail["originalAmount"] / canada_retail["rateValue"]
    canada_shipments_eur = canada_shipments["originalAmount"] / canada_shipments["rateValue"]
    nz_observed_display = (nz_observed_eur / Decimal("1000000")).quantize(
        Decimal("0.1"), rounding=ROUND_HALF_UP
    )
    nz_low_display = (nz_low_eur / Decimal("1000000")).quantize(
        Decimal("0.1"), rounding=ROUND_HALF_UP
    )
    nz_high_display = (nz_high_eur / Decimal("1000000")).quantize(
        Decimal("0.1"), rounding=ROUND_HALF_UP
    )
    ftc_display = (ftc_eur / Decimal("1000000000")).quantize(
        Decimal("0.001"), rounding=ROUND_HALF_UP
    )
    canada_retail_display = (canada_retail_eur / Decimal("1000000")).quantize(
        Decimal("0.1"), rounding=ROUND_HALF_UP
    )
    canada_shipments_display = (canada_shipments_eur / Decimal("1000000")).quantize(
        Decimal("0.1"), rounding=ROUND_HALF_UP
    )
    if language == "fi":
        return (
            f"≈{str(nz_observed_display).replace('.', ',')} milj. eur",
            f"≈{str(nz_low_display).replace('.', ',')}–"
            f"{str(nz_high_display).replace('.', ',')} milj. eur",
            f"≈{str(ftc_display).replace('.', ',')} mrd eur",
            f"≈{str(canada_retail_display).replace('.', ',')} milj. eur",
            f"≈{str(canada_shipments_display).replace('.', ',')} milj. eur",
        )
    return (
        f"≈eur {nz_observed_display}m",
        f"≈eur {nz_low_display}–{nz_high_display}m",
        f"≈eur {ftc_display}bn",
        f"≈eur {canada_retail_display}m",
        f"≈eur {canada_shipments_display}m",
    )


def strings(value: Any):
    if isinstance(value, str):
        yield value
    elif isinstance(value, dict):
        for key, item in value.items():
            yield str(key)
            yield from strings(item)
    elif isinstance(value, list):
        for item in value:
            yield from strings(item)


def validate_forbidden_terms(label: str, text: str, errors: list[str]) -> None:
    lowered = text.casefold()
    for phrase in FORBIDDEN_TEXT:
        if phrase in lowered:
            errors.append(f"{label}: forbidden private/local term {phrase!r}")
    normalised = re.sub(r"[^a-z0-9]+", "", lowered)
    for length, expected in PRIVATE_IDENTIFIER_FINGERPRINTS:
        if any(
            hashlib.sha256(normalised[index:index + length].encode("utf-8")).hexdigest() == expected
            for index in range(max(0, len(normalised) - length + 1))
        ):
            errors.append(f"{label}: forbidden private identifier fingerprint")
            break


def validate_text(label: str, text: str, errors: list[str]) -> None:
    validate_forbidden_terms(label, text, errors)
    for match in re.finditer(r"https?://[^\s<>\"']+", text):
        parsed = urlparse(match.group(0).rstrip(".,);]"))
        if parsed.scheme != "https" or not parsed.netloc:
            errors.append(f"{label}: only public HTTPS links are allowed")
            continue
        query_keys = {key.casefold() for key, _ in parse_qsl(parsed.query, keep_blank_values=True)}
        if query_keys & SENSITIVE_QUERY_KEYS:
            errors.append(f"{label}: URL contains a sensitive query key")


def validate_external_https_target(label: str, target: str, errors: list[str]) -> None:
    """Require the entire external relationship target to be a safe HTTPS URL."""

    if target != target.strip() or any(character.isspace() or ord(character) < 32 for character in target):
        errors.append(f"{label}: external hyperlink target contains whitespace/control characters")
        return
    if re.search(r"%(?![0-9A-Fa-f]{2})", target):
        errors.append(f"{label}: external hyperlink target contains malformed percent-encoding")
        return
    try:
        parsed = urlparse(target)
        # Accessing port forces urllib to reject malformed or out-of-range ports.
        _ = parsed.port
    except ValueError:
        errors.append(f"{label}: external hyperlink target is malformed")
        return
    if parsed.scheme.casefold() != "https" or not parsed.netloc or not parsed.hostname:
        errors.append(f"{label}: external hyperlink target must be an absolute HTTPS URL")
        return
    if parsed.username is not None or parsed.password is not None:
        errors.append(f"{label}: external hyperlink target must not contain credentials")
    if "\\" in target:
        errors.append(f"{label}: external hyperlink target must not contain backslashes or UNC syntax")
    query_keys = {key.casefold() for key, _ in parse_qsl(parsed.query, keep_blank_values=True)}
    if query_keys & SENSITIVE_QUERY_KEYS:
        errors.append(f"{label}: external hyperlink target contains a sensitive query key")
    validate_forbidden_terms(label, target, errors)


def validate_ooxml(
    path: Path,
    errors: list[str],
    *,
    require_deterministic_zip: bool,
    allow_notes: bool,
) -> str:
    label = str(path.relative_to(ROOT))
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            names = [info.filename for info in infos]
            if not names or "[Content_Types].xml" not in names:
                errors.append(f"{label}: invalid OOXML package")
                return ""
            if len(names) != len(set(names)):
                errors.append(f"{label}: duplicate ZIP entries")
            if any(name.startswith("/") or ".." in Path(name).parts for name in names):
                errors.append(f"{label}: unsafe ZIP path")
            if require_deterministic_zip:
                if names != sorted(names):
                    errors.append(f"{label}: ZIP entries are not deterministically ordered")
                timestamps = {info.date_time for info in infos}
                if len(timestamps) != 1:
                    errors.append(f"{label}: ZIP timestamps are not normalized")

            extracted_text: list[str] = []
            for info in infos:
                lowered_name = f"/{info.filename}".casefold()
                if any(part.casefold() in lowered_name for part in FORBIDDEN_ARCHIVE_PARTS):
                    errors.append(f"{label}: forbidden OOXML part {info.filename}")
                if info.file_size > 20 * 1024 * 1024:
                    errors.append(f"{label}: oversized OOXML part {info.filename}")
                if info.filename.endswith((".xml", ".rels")):
                    payload = archive.read(info).decode("utf-8", errors="replace")
                    extracted_text.append(payload)
                    is_notes_part = "/notesslides/" in lowered_name or "/notesmasters/" in lowered_name
                    if is_notes_part:
                        if not allow_notes:
                            errors.append(f"{label}: notes parts are forbidden ({info.filename})")
                        elif info.filename.endswith(".xml"):
                            try:
                                ET.fromstring(payload)
                            except ET.ParseError:
                                errors.append(f"{label}: malformed notes part {info.filename}")
                    if info.filename.endswith(".rels"):
                        try:
                            root = ET.fromstring(payload)
                        except ET.ParseError:
                            errors.append(f"{label}: malformed relationship part {info.filename}")
                            continue
                        for relation in root:
                            if relation.attrib.get("TargetMode") != "External":
                                continue
                            target = relation.attrib.get("Target", "")
                            relation_type = relation.attrib.get("Type", "").casefold()
                            if "hyperlink" not in relation_type:
                                errors.append(f"{label}: external non-hyperlink relationship is forbidden")
                            else:
                                validate_external_https_target(f"{label} relationship", target, errors)
            combined = "\n".join(extracted_text)
            # Namespace declarations and relationship type identifiers use HTTP
            # URIs by OOXML design; URL policy applies only to visible content
            # and explicit TargetMode=External relationships.
            validate_forbidden_terms(label, combined, errors)
            return combined
    except (OSError, zipfile.BadZipFile) as error:
        errors.append(f"{label}: unreadable OOXML package: {error}")
        return ""


def slide_texts(path: Path, errors: list[str]) -> list[str]:
    try:
        presentation = Presentation(path)
    except Exception as error:  # python-pptx exposes parser-specific exceptions
        errors.append(f"{path.relative_to(ROOT)}: cannot parse presentation: {error}")
        return []
    output: list[str] = []
    for slide in presentation.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                chunks.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
        output.append("\n".join(chunks))
    return output


def validate_slide_source_notes(
    path: Path,
    fx: dict[str, Any],
    errors: list[str],
) -> None:
    label = str(path.relative_to(ROOT))
    rates = {
        (item.get("currency"), item.get("year")): item
        for item in fx.get("rates", [])
        if isinstance(item, dict)
    }
    required_fx_urls = {
        fx.get("provider", {}).get("methodologyUrl"),
        rates.get(("NZD", 2024), {}).get("sourceUrl"),
        rates.get(("USD", 2021), {}).get("sourceUrl"),
        rates.get(("CAD", 2024), {}).get("sourceUrl"),
    } - {None}
    required_source_urls = {*required_fx_urls, FHM_SOURCE_URL}
    formula = fx.get("calculationPolicy", {}).get("formulaEn")
    try:
        presentation = Presentation(path)
    except Exception as error:
        errors.append(f"{label}: cannot parse presentation notes: {error}")
        return
    for index, slide in enumerate(presentation.slides, start=1):
        if not slide.has_notes_slide or slide.notes_slide.notes_text_frame is None:
            errors.append(f"{label}: slide {index} is missing speaker notes")
            continue
        notes = str(slide.notes_slide.notes_text_frame.text or "").strip()
        validate_text(f"{label} slide {index} notes", notes, errors)
        if "[Sources]" not in notes:
            errors.append(f"{label}: slide {index} notes lack a [Sources] block")
        if not re.search(r"https://[^\s]+", notes):
            errors.append(f"{label}: slide {index} [Sources] block lacks a public HTTPS source")
        if "[FX methodology]" not in notes or not formula or formula not in notes:
            errors.append(f"{label}: slide {index} notes lack the reviewed FX methodology")
        for required_url in required_source_urls:
            if required_url not in notes:
                errors.append(
                    f"{label}: slide {index} notes lack required reviewed source {required_url}"
                )


def read_register_csv(
    path: Path,
    headers: list[str],
    allowed_statuses: set[str],
    errors: list[str],
) -> list[list[str]]:
    label = str(path.relative_to(ROOT))
    try:
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.reader(handle))
    except OSError as error:
        errors.append(f"{label}: {error}")
        return []
    if not rows or rows[0] != headers:
        errors.append(f"{label} has incorrect headers")
        return []
    data_rows = [[str(value) for value in row] for row in rows[1:] if any(str(value).strip() for value in row)]
    for index, row in enumerate(data_rows, start=2):
        if len(row) != len(headers):
            errors.append(f"{label} row {index} has {len(row)} columns")
            continue
        if row[7] not in allowed_statuses:
            errors.append(f"{label} row {index} has invalid status {row[7]!r}")
        validate_text(f"{label} row {index}", "\n".join(row), errors)
    if not data_rows:
        errors.append(f"{label} must contain evidence rows")
    elif {row[7] for row in data_rows if len(row) == len(headers)} != allowed_statuses:
        errors.append(f"{label} must visibly use all four evidence classifications")
    return data_rows


def validate_eur_equivalent_sheet(
    workbook: Any,
    language: str,
    expected_rows: list[dict[str, Any]],
    compare_with_current_inputs: bool,
    label: str,
    errors: list[str],
) -> None:
    sheet_name = EUR_EQUIVALENT_SHEET_NAMES[language]
    if sheet_name not in workbook.sheetnames:
        errors.append(f"{label}: missing {sheet_name} sheet")
        return
    sheet = workbook[sheet_name]
    headers = [str(sheet.cell(1, column).value or "") for column in range(1, 15)]
    if headers != EUR_EQUIVALENT_HEADERS[language]:
        errors.append(f"{label}: {sheet_name} headers are incorrect")
    actual_rows = [
        row
        for row in sheet.iter_rows(min_row=2, max_col=14)
        if any(cell.value not in (None, "") for cell in row)
    ]
    expected_row_count = (
        len(expected_rows)
        if compare_with_current_inputs
        else EXPECTED_LOCKED_EUR_EQUIVALENT_ROWS
    )
    if len(actual_rows) != expected_row_count:
        errors.append(
            f"{label}: {sheet_name} row coverage differs "
            f"({len(actual_rows)} != {expected_row_count})"
        )
    seen_keys: set[tuple[str, str, str]] = set()
    status_counts: dict[str, int] = {}
    for index, cells in enumerate(actual_rows, start=2):
        key = tuple(str(cells[column - 1].value or "") for column in (1, 2, 3))
        if any(not value for value in key):
            errors.append(f"{label}: {sheet_name} row {index} lacks a complete record key")
        elif key in seen_keys:
            errors.append(f"{label}: {sheet_name} row {index} duplicates record key {key!r}")
        seen_keys.add(key)

        status = str(cells[12].value or "")
        status_counts[status] = status_counts.get(status, 0) + 1
        original_amount = decimal_value(cells[6].value)
        rate_value = decimal_value(cells[8].value)
        currency = str(cells[7].value or "")
        rate_id = str(cells[10].value or "")
        source_url = str(cells[11].value or "")
        parsed = urlparse(source_url)
        if original_amount is None or original_amount <= 0:
            errors.append(f"{label}: {sheet_name} row {index} lacks a positive original amount")
        if not re.fullmatch(r"[A-Z]{3}", currency):
            errors.append(f"{label}: {sheet_name} row {index} has an invalid currency code")

        if status == "computed":
            expected_formula = f"=G{index}/I{index}"
            if rate_value is None or rate_value <= 0:
                errors.append(f"{label}: {sheet_name} row {index} lacks a positive ECB rate")
            if (
                not rate_id.startswith("ECB-EXR-A-")
                or parsed.scheme != "https"
                or parsed.hostname != "data-api.ecb.europa.eu"
            ):
                errors.append(
                    f"{label}: {sheet_name} row {index} lacks direct ECB rateId/source URL"
                )
        elif status == "already_eur":
            expected_formula = f"=G{index}"
            if currency != "EUR" or rate_value != Decimal("1") or rate_id != "EUR-IDENTITY":
                errors.append(
                    f"{label}: {sheet_name} row {index} lacks the EUR identity conversion"
                )
            if parsed.scheme != "https" or parsed.hostname != "www.ecb.europa.eu":
                errors.append(
                    f"{label}: {sheet_name} row {index} lacks the official ECB identity source"
                )
        elif status == "not_computed":
            expected_formula = None
            if cells[8].value not in (None, "") or rate_id:
                errors.append(
                    f"{label}: {sheet_name} row {index} must fail closed without an ECB rate/rateId"
                )
        else:
            expected_formula = None
            errors.append(f"{label}: {sheet_name} row {index} has invalid status {status!r}")
        if cells[9].value != expected_formula:
            errors.append(
                f"{label}: {sheet_name}!{cells[9].coordinate} must preserve "
                f"full-precision formula {expected_formula!r}"
            )

    if not compare_with_current_inputs:
        if status_counts != EXPECTED_LOCKED_EUR_STATUS_COUNTS:
            errors.append(
                f"{label}: {sheet_name} locked-snapshot status coverage differs "
                f"({status_counts!r} != {EXPECTED_LOCKED_EUR_STATUS_COUNTS!r})"
            )
        return

    for index, expected in enumerate(expected_rows, start=2):
        if index - 2 >= len(actual_rows):
            break
        cells = actual_rows[index - 2]
        expected_text = {
            1: expected["recordType"],
            2: expected["recordId"],
            3: expected["item"],
            4: expected["geography"],
            5: expected["year"],
            6: expected["period"],
            8: expected["currency"],
            11: expected["rateId"],
            12: expected["sourceUrl"],
            13: expected["status"],
            14: expected["reason"],
        }
        for column, value in expected_text.items():
            actual = cells[column - 1].value
            if ("" if actual is None else str(actual)) != ("" if value is None else str(value)):
                errors.append(
                    f"{label}: {sheet_name}!{cells[column - 1].coordinate} "
                    "differs from the reviewed FX row"
                )
        if decimal_value(cells[6].value) != expected["originalAmount"]:
            errors.append(
                f"{label}: {sheet_name}!{cells[6].coordinate} original amount differs"
            )
        if decimal_value(cells[8].value) != expected["rateValue"]:
            errors.append(
                f"{label}: {sheet_name}!{cells[8].coordinate} ECB rate differs"
            )


def validate_workbook(
    path: Path,
    csv_rows: list[list[str]],
    expected_headers: list[str],
    expected_eur_rows: list[dict[str, Any]],
    compare_eur_with_current_inputs: bool,
    errors: list[str],
) -> int:
    label = str(path.relative_to(ROOT))
    try:
        workbook = load_workbook(path, data_only=False, read_only=False)
    except Exception as error:
        errors.append(f"{label}: cannot parse workbook: {error}")
        return 0
    if "Evidence Register" not in workbook.sheetnames:
        errors.append(f"{label}: missing Evidence Register sheet")
        return 0
    for sheet in workbook.worksheets:
        if sheet.sheet_state != "visible":
            errors.append(f"{label}: hidden worksheets are forbidden ({sheet.title})")
        for row in sheet.iter_rows():
            for cell in row:
                if cell.comment is not None:
                    errors.append(f"{label}: comments are forbidden ({sheet.title}!{cell.coordinate})")
                if isinstance(cell.value, str):
                    validate_text(f"{label} {sheet.title}!{cell.coordinate}", cell.value, errors)
                    if cell.value.startswith("=") and "[" in cell.value:
                        errors.append(f"{label}: external workbook formula is forbidden")
                if cell.hyperlink is not None:
                    validate_text(f"{label} hyperlink", str(cell.hyperlink.target), errors)
    sheet = workbook["Evidence Register"]
    headers = [str(sheet.cell(1, column).value or "") for column in range(1, 10)]
    if headers != expected_headers:
        errors.append(f"{label}: Evidence Register headers are incorrect")
    workbook_rows: list[list[str]] = []
    for values in sheet.iter_rows(min_row=2, max_col=9, values_only=True):
        row = ["" if value is None else str(value) for value in values]
        if any(value.strip() for value in row):
            workbook_rows.append(row)
    if workbook_rows != csv_rows:
        errors.append(f"{label}: Evidence Register rows differ from the public CSV")
    is_finnish = expected_headers == REGISTER_HEADERS
    summary_name = "Yhteenveto" if is_finnish else "Summary"
    if summary_name not in workbook.sheetnames:
        errors.append(f"{label}: missing {summary_name} sheet")
    else:
        summary = workbook[summary_name]
        evidence_end = len(csv_rows) + 1
        expected_formulas = {
            "B8": f"=COUNTA('Evidence Register'!$A$2:$A${evidence_end})",
            **{
                f"B{row}": f"=COUNTIF('Evidence Register'!$H$2:$H${evidence_end},A{row})"
                for row in range(11, 15)
            },
        }
        for coordinate, formula in expected_formulas.items():
            if summary[coordinate].value != formula:
                errors.append(f"{label}: {summary_name}!{coordinate} must preserve formula {formula}")

    sources_name = "Lähteet" if is_finnish else "Sources"
    if sources_name not in workbook.sheetnames:
        errors.append(f"{label}: missing {sources_name} sheet")
    else:
        sources_sheet = workbook[sources_name]
        source_urls = {
            str(sources_sheet.cell(row, 4).value or "").strip()
            for row in range(2, sources_sheet.max_row + 1)
            if str(sources_sheet.cell(row, 4).value or "").strip()
        }
        register_urls = {
            match.group(0).rstrip(".,);]")
            for row in csv_rows
            for match in re.finditer(r"https://[^\s;]+", row[3])
        }
        missing_urls = sorted(register_urls - source_urls)
        if missing_urls:
            errors.append(f"{label}: {sources_name} omits register URLs: {missing_urls}")
        for required_url in {
            "https://www.ftc.gov/reports/e-cigarette-report-2015-2018",
            "https://www.ftc.gov/reports/e-cigarette-report-2021",
            "https://www.un.org/en/about-us/member-states",
            "https://www.un.org/en/about-us/non-member-states",
        }:
            if required_url not in source_urls:
                errors.append(f"{label}: {sources_name} lacks required source {required_url}")

    nz_prefix = (
        "Uuden-Seelannin vuoden 2024 AIS/AVP"
        if is_finnish
        else "New Zealand's 2024 AIS/AVP"
    )
    nz_row = next(
        (index for index, row in enumerate(csv_rows, start=2) if row[0].startswith(nz_prefix)),
        None,
    )
    if nz_row is None:
        errors.append(f"{label}: supported New Zealand model row is missing")
    else:
        row_height = sheet.row_dimensions[nz_row].height
        if row_height is None or row_height < 80:
            errors.append(f"{label}: supported New Zealand model row lacks the expanded review treatment")
        if sheet[f"F{nz_row}"].fill.fill_type is None:
            errors.append(f"{label}: supported New Zealand calculation cell lacks the review highlight")
    validate_eur_equivalent_sheet(
        workbook,
        "fi" if is_finnish else "en",
        expected_eur_rows,
        compare_eur_with_current_inputs,
        label,
        errors,
    )
    return len(workbook_rows)


def validate_manifest(errors: list[str]) -> None:
    expected_eur_rows, fx = load_expected_eur_equivalent_rows(errors)
    try:
        builder_text = ARTIFACT_BUILDER_PATH.read_text(encoding="utf-8")
    except OSError as error:
        errors.append(f"cannot read artifact builder: {error}")
        builder_text = ""
    validate_artifact_builder_fx_contract(builder_text, errors)
    try:
        manifest = load_json(MANIFEST_PATH)
        changelog = load_json(CHANGELOG_PATH)
        lock = load_json(LOCK_PATH)
    except ValueError as error:
        errors.append(str(error))
        return
    expected_keys = {
        "schemaVersion",
        "generatedFromPublicDataOnly",
        "release",
        "asOf",
        "cadence",
        "languages",
        "publicBoundary",
        "templateInputs",
        "inputs",
        "artifacts",
    }
    if not isinstance(manifest, dict) or set(manifest) != expected_keys:
        errors.append("bank-package-manifest.json has an unexpected schema")
        return
    if manifest.get("schemaVersion") != 2 or manifest.get("generatedFromPublicDataOnly") is not True:
        errors.append("manifest must declare schemaVersion 2 and public-data-only generation")
    if manifest.get("languages") != ["en", "fi"]:
        errors.append("manifest languages must be exactly en and fi")
    earlier_same_day_snapshot = validate_daily_package_snapshot(
        manifest,
        changelog,
        errors,
    )
    release_lock_valid = validate_release_lock(manifest, lock, errors)
    allow_reviewed_input_drift = earlier_same_day_snapshot and release_lock_valid
    package_release = manifest.get("release", {})
    if (
        package_release.get("id") != RELEASE_ID
        or package_release.get("version") != RELEASE_VERSION
        or manifest.get("asOf") != RELEASE_DATE
    ):
        errors.append(
            f"bank package must be locked to release {RELEASE_VERSION} as of {RELEASE_DATE}"
        )
    boundary = manifest.get("publicBoundary")
    if not isinstance(boundary, dict) or set(boundary) != {"en", "fi"}:
        errors.append("manifest publicBoundary must contain exactly en and fi")
    else:
        boundary_text = " ".join(str(value) for value in boundary.values())
        validate_text("manifest public boundary", boundary_text, errors)
        if "public" not in str(boundary.get("en", "")).casefold() or "julk" not in str(boundary.get("fi", "")).casefold():
            errors.append("manifest must state the public-data boundary in both languages")

    template_inputs = manifest.get("templateInputs")
    if not isinstance(template_inputs, list):
        errors.append("manifest templateInputs must be an array")
        template_inputs = []
    template_by_path = {
        item.get("path"): item
        for item in template_inputs
        if isinstance(item, dict) and set(item) == {"path", "sha256"}
    }
    if set(template_by_path) != EXPECTED_TEMPLATE_INPUTS or len(template_by_path) != len(template_inputs):
        errors.append("manifest templateInputs must contain the exact reviewed seed artifacts")
    for relative, item in template_by_path.items():
        seed_path = ROOT / relative
        if not seed_path.is_file():
            errors.append(f"manifest seed artifact is missing: {relative}")
        elif item.get("sha256") != sha256(seed_path):
            errors.append(f"manifest seed artifact hash differs: {relative}")

    generated_by = lock.get("generatedBy") if isinstance(lock, dict) else None
    if not isinstance(generated_by, dict):
        errors.append("bank-package lock lacks generatedBy lineage")
    else:
        if generated_by.get("tool") != "@oai/artifact-tool":
            errors.append("bank-package lock must identify @oai/artifact-tool")
        if not re.fullmatch(r"\d+\.\d+\.\d+(?:[-+][0-9A-Za-z.-]+)?", str(generated_by.get("toolVersion", ""))):
            errors.append("bank-package lock must record the runtime-resolved artifact-tool version")
        lock_templates = generated_by.get("sourceTemplates")
        if lock_templates != template_inputs:
            errors.append("bank-package lock sourceTemplates must match manifest templateInputs")
        quality = generated_by.get("qualityAssurance")
        if not isinstance(quality, dict) or any(
            quality.get(key) is not True
            for key in (
                "summaryFormulasAfterReopen",
                "allSlidesRendered",
                "allWorkbookSheetsRendered",
                "sourcesNotesOnEverySlide",
                "eurEquivalentRowsAfterReopen",
                "fxSourcesInDeckNotes",
            )
        ):
            errors.append("bank-package lock lacks required artifact QA lineage")

    if 'toolVersion: artifactToolVersion' not in builder_text or 'toolVersion: "2.8.' in builder_text:
        errors.append("artifact builder must derive the artifact-tool version at runtime")

    inputs = manifest.get("inputs")
    if not isinstance(inputs, list):
        errors.append("manifest inputs must be an array")
        inputs = []
    input_by_path = {
        item.get("path"): item for item in inputs if isinstance(item, dict) and set(item) == {"path", "sha256"}
    }
    expected_input_paths = EXPECTED_INPUTS
    if set(input_by_path) != expected_input_paths or len(input_by_path) != len(inputs):
        errors.append("manifest inputs must be the exact reviewed public-data allowlist")
    for relative, item in input_by_path.items():
        path = ROOT / relative
        if not path.is_file():
            errors.append(f"manifest input is missing: {relative}")
        elif not SHA256_RE.fullmatch(str(item.get("sha256", ""))):
            errors.append(f"manifest input hash is invalid: {relative}")
        elif item.get("sha256") != sha256(path):
            if not allow_reviewed_input_drift or relative == LOCK_RELATIVE_PATH:
                errors.append(f"manifest input hash differs: {relative}")

    artifacts = manifest.get("artifacts")
    if not isinstance(artifacts, list):
        errors.append("manifest artifacts must be an array")
        artifacts = []
    artifact_by_id = {
        item.get("id"): item
        for item in artifacts
        if isinstance(item, dict) and isinstance(item.get("id"), str)
    }
    if set(artifact_by_id) != set(EXPECTED_ARTIFACTS) or len(artifact_by_id) != len(artifacts):
        errors.append("manifest artifacts must contain exactly the six approved downloads")

    csv_rows_by_language = {
        "fi": read_register_csv(REGISTER_CSV_PATH, REGISTER_HEADERS, ALLOWED_STATUSES, errors),
        "en": read_register_csv(EN_REGISTER_CSV_PATH, EN_REGISTER_HEADERS, EN_ALLOWED_STATUSES, errors),
    }
    if any(len(rows) != 60 for rows in csv_rows_by_language.values()):
        errors.append("both v43 Evidence Registers must contain exactly 60 reviewed rows")
    register_markers = {
        "fi": (
            "280 684 512,81",
            "274 180 410,21",
            "189 402 451,96",
            "84 709 409,85",
            "68 548,40",
            "2 137 085,24",
            "4 367 017,37",
            "189 640 890",
            "203 340 531",
            "6 270 209",
            "183 370 681",
            "197 070 322",
            "1,495224911",
            "1,391282094",
            "2 763 284 338",
            "4,99 mrd",
            "1 219 160 000",
            "5,031748 %",
            "D1–D10",
            "174 havaintoa 54 lähteestä",
            "152 virallista havaintoa",
            "116 markkinamittariin",
            "36 Ruotsin FHM-rekisterirakenteen lukuun",
            "578 havaittua World Bank",
            "194/195",
            "190/195",
            "28 reviewed_method_plan",
            "0 reviewed_source_lead",
            "15 regional_tpd_pattern_only",
            "152 proxy_only_unscoped",
            "Itävallalla on tarkistettu",
            "Belgian virallinen osavuosiveroluku",
            "83 333 333 ml",
            "83 333 litraa",
            "9 kuukauden verovolyymi-indikaattorin",
            "Sveitsillä on 2024-10-01",
            "Luxemburgilla on 2024-10-01",
            "kulutukseen luovutettuja määriä",
            "hintalistaa ei kerrota vahvistamattomalla pakkausmäärällä",
            "Norjalla ei ole nykyistä Article 20(7)",
            "55 910 871,89 EUR",
            "84 309 841,41 EUR",
            "kyse ei ole vähittäismarkkina-arvosta",
            "1 451 529 litraa vuonna 2020",
            "4 382 500",
            "62 500",
            "laajan höyrystyslaiteryhmän",
            "eivät ole sähkötupakkakohtaisia",
            "eivät ole myyntiä",
            "26 000 litraa",
            "80 000 000 SEK",
            "DE-BLIND-1.0.0",
            "Kaikki kolme numeerista testiä läpäistiin",
            "toimittaja-arvot ja tarkat poikkeamat pidetään salassa",
            "Saksa ei ole donor",
            "Euromonitor on 1/6",
            "EI PISTEYTETTY",
            "laajempi paketti HOLD",
            "retail-arvo not_computed",
        ),
        "en": (
            "280,684,512.81",
            "274,180,410.21",
            "189,402,451.96",
            "84,709,409.85",
            "68,548.40",
            "2,137,085.24",
            "4,367,017.37",
            "189,640,890",
            "203,340,531",
            "6,270,209",
            "183,370,681",
            "197,070,322",
            "1.495224911",
            "1.391282094",
            "2,763,284,338",
            "4.99 billion",
            "1,219,160,000",
            "5.031748%",
            "D1–D10",
            "174 observations from 54 sources",
            "152 official observations",
            "116 market measures",
            "36 Swedish FHM register-structure counts",
            "578 observed World Bank",
            "194/195",
            "190/195",
            "28 reviewed_method_plan",
            "0 reviewed_source_lead",
            "15 regional_tpd_pattern_only",
            "152 proxy_only_unscoped",
            "Austria has a reviewed method plan",
            "Belgium’s official partial-period tax figure",
            "83,333,333 ml",
            "83,333 litres",
            "9 month tax-volume indicator",
            "Switzerland has a two-rate method route",
            "Luxembourg has an excise and fiscal-mark method",
            "consumption-release volumes",
            "price list is not multiplied by unverified pack counts",
            "Norway has no current Article 20(7)",
            "EUR 55,910,871.89",
            "EUR 84,309,841.41",
            "this is not retail market value",
            "1,451,529 litres in 2020",
            "4,382,500",
            "62,500",
            "broad vaporisation-device-group",
            "neither is e-cigarette-specific",
            "are not sales",
            "26,000 litres",
            "SEK 80,000,000",
            "DE-BLIND-1.0.0",
            "All three numerical tests passed",
            "vendor values and exact deviations remain withheld",
            "Germany is not a donor",
            "Euromonitor is 1/6",
            "NOT SCORED",
            "wider package is HOLD",
            "retail value remains not_computed",
        ),
    }
    for language, rows in csv_rows_by_language.items():
        joined = "\n".join("\t".join(row) for row in rows)
        for marker in register_markers[language]:
            if marker not in joined:
                errors.append(f"{language} Evidence Register lacks v43 marker {marker!r}")
    errors.extend(
        validate_register_parity(
            csv_rows_by_language["fi"],
            csv_rows_by_language["en"],
        )
    )
    headers_by_language = {"fi": REGISTER_HEADERS, "en": EN_REGISTER_HEADERS}
    for artifact_id, expected in EXPECTED_ARTIFACTS.items():
        item = artifact_by_id.get(artifact_id)
        if not isinstance(item, dict):
            continue
        required = {"id", "kind", "language", "titleFi", "titleEn", "fileName", "path", "sha256", "bytes"}
        if expected["kind"] == "pptx":
            required.add("slideCount")
        else:
            required.add("rowCount")
        if set(item) != required:
            errors.append(f"manifest artifact {artifact_id} has an unexpected schema")
        if (
            item.get("kind") != expected["kind"]
            or item.get("language") != expected["language"]
            or item.get("path") != expected["path"]
        ):
            errors.append(f"manifest artifact {artifact_id} kind/language/path differs from allowlist")
        if item.get("fileName") != Path(expected["path"]).name:
            errors.append(f"manifest artifact {artifact_id} filename differs from path")
        if not str(item.get("titleFi", "")).strip() or not str(item.get("titleEn", "")).strip():
            errors.append(f"manifest artifact {artifact_id} requires bilingual titles")
        relative = str(item.get("path", ""))
        path = ROOT / "site" / relative
        if path.parent != ROOT / "site" / "downloads" or not path.is_file():
            errors.append(f"download missing or outside allowlist: {relative}")
            continue
        if path.stat().st_size > 12 * 1024 * 1024:
            errors.append(f"{relative}: file exceeds 12 MiB")
        if item.get("bytes") != path.stat().st_size:
            errors.append(f"{relative}: manifest byte count differs")
        if not SHA256_RE.fullmatch(str(item.get("sha256", ""))) or item.get("sha256") != sha256(path):
            errors.append(f"{relative}: manifest SHA-256 differs")
        is_english = expected["language"] == "en"
        validate_ooxml(
            path,
            errors,
            require_deterministic_zip=False,
            allow_notes=expected["kind"] == "pptx",
        )
        if expected["kind"] == "pptx":
            validate_slide_source_notes(path, fx, errors)
            texts = slide_texts(path, errors)
            if len(texts) != expected["slideCount"] or item.get("slideCount") != expected["slideCount"]:
                errors.append(f"{relative}: expected exactly {expected['slideCount']} slides")
            for index, text in enumerate(texts, start=1):
                if not text.strip():
                    errors.append(f"{relative}: slide {index} has no readable text")
                validate_text(f"{relative} slide {index}", text, errors)
            combined = "\n".join(texts).casefold()
            expected_boundary = "independent public evidence" if is_english else "julkinen riippumaton"
            if expected_boundary not in combined:
                errors.append(f"{relative}: public-boundary disclosure is missing")
            release_deck_markers = (
                (
                    "274,180 milj. nzd",
                    "116 markkinamittaria",
                    "36 ruotsin fhm",
                    "578 wb-havaintoa",
                    "28 / 0 / 15 / 152",
                    "0/3",
                    "7/10",
                    "1,219160 mrd cad",
                    "uusi-seelanti",
                    "euromonitor",
                    "1/6",
                    "hold",
                    "de-blind",
                    "ei pisteytetty",
                    "arvot salassa",
                    RELEASE_DATE,
                    RELEASE_VERSION,
                )
                if not is_english
                else (
                    "nzd 274.180m",
                    "116 market measures",
                    "36 swedish fhm register",
                    "578 wb records",
                    "28 / 0 / 15 / 152",
                    "0/3",
                    "7/10",
                    "cad 1.219160bn",
                    "new zealand",
                    "euromonitor",
                    "1/6",
                    "hold",
                    "de-blind",
                    "not scored",
                    "values remain withheld",
                    RELEASE_DATE,
                    RELEASE_VERSION,
                )
            )
            (
                nz_observed_fx_marker,
                nz_model_fx_marker,
                ftc_fx_marker,
                canada_retail_fx_marker,
                canada_shipments_fx_marker,
            ) = deck_fx_markers(
                expected_eur_rows,
                "en" if is_english else "fi",
            )
            release_deck_markers = (
                *release_deck_markers,
                nz_observed_fx_marker,
                canada_retail_fx_marker,
                canada_shipments_fx_marker,
            )
            if expected["slideCount"] == 30:
                release_deck_markers = (
                    *release_deck_markers,
                    "2,763 mrd usd" if not is_english else "usd 2.763bn",
                    "4,99 mrd eur" if not is_english else "eur 4.99bn",
                    "4 382 500" if not is_english else "4,382,500",
                    "62 500" if not is_english else "62,500",
                    (
                        "laajan höyrystyslaiteryhmän"
                        if not is_english
                        else "broad vaporisation-device-group"
                    ),
                    (
                        "ei sähkötupakkakohtainen"
                        if not is_english
                        else "not e-cigarette-only"
                    ),
                    "puola" if not is_english else "poland",
                    "euromonitor",
                    ftc_fx_marker,
                    nz_model_fx_marker,
                    "nzd 533.7–731.2m" if is_english else "533,7–731,2 milj. nzd",
                    (
                        "päiväpaketti muodostetaan enintään kerran asia/nicosia-kalenteripäivässä"
                        if not is_english
                        else "the daily package is generated at most once per asia/nicosia calendar day"
                    ),
                )
            for marker in release_deck_markers:
                if marker not in combined:
                    errors.append(f"{relative}: v43 market marker is missing: {marker!r}")
            if expected["slideCount"] == 6:
                short_only_markers = (
                    (
                        "nz:n rajaproxyt eivät ole retail-arvoja",
                        "yksityinen de-blind-audit",
                        "vuosien 2023 ja 2024 vuosirajat sekä yhteisrajan",
                        "arvot salassa",
                        "1/6",
                        "ei pisteytetty",
                        "hold",
                    )
                    if not is_english
                    else (
                        "nz border proxies are not retail values",
                        "private de-blind audit",
                        "2023, 2024 and combined caps",
                        "values remain withheld",
                        "1/6",
                        "not scored",
                        "hold",
                    )
                )
                for marker in short_only_markers:
                    if marker not in combined:
                        errors.append(
                            f"{relative}: v43 concise-deck marker is missing: {marker!r}"
                        )
            if expected["slideCount"] == 30:
                large_only_markers = (
                    (
                        "d8 on suljettu virallisella veroperustalla",
                        "183,371/197,070 milj. nzd:n nettorajaproxyt ovat vain d10-diagnostiikkaa",
                        "lisensoitu vuosien 2022–2025 saksa-ote vastaanotettiin",
                        "vuoden 2023 ja 2024 vuosirajat sekä kahden vuoden yhteisrajan",
                        "toimittaja-arvot ja tarkat poikkeamat pidetään salassa",
                        "saksa ei ole donor",
                        "laajempi 25/50/78 maan paketti hold",
                    )
                    if not is_english
                    else (
                        "official tax evidence closes d8",
                        "nzd 183.371m/197.070m net border proxies are d10 diagnostics only",
                        "licensed germany 2022–2025 extract was received",
                        "2023 and 2024 annual caps and the two-year combined cap",
                        "vendor values and exact deviations remain withheld",
                        "germany is not a donor",
                        "wider 25/50/78-country package hold",
                    )
                )
                for marker in large_only_markers:
                    if marker not in combined:
                        errors.append(
                            f"{relative}: v43 extended-deck marker is missing: {marker!r}"
                        )
        else:
            csv_rows = csv_rows_by_language[expected["language"]]
            row_count = validate_workbook(
                path,
                csv_rows,
                headers_by_language[expected["language"]],
                expected_eur_rows,
                not allow_reviewed_input_drift,
                errors,
            )
            if item.get("rowCount") != row_count or row_count != len(csv_rows):
                errors.append(f"{relative}: manifest/workbook/CSV row counts differ")


def main() -> None:
    errors: list[str] = []
    validate_v22_market_bindings(errors)
    validate_v43_vendor_and_global_boundary(errors)
    validate_manifest(errors)
    if errors:
        for error in errors:
            print(f"ERROR: {error}", file=sys.stderr)
        print(f"Bank-package validation failed with {len(errors)} error(s).", file=sys.stderr)
        raise SystemExit(1)
    print(
        "Validated bilingual public bank package: English and Finnish 6- and 30-slide decks, "
        "Evidence Register parity, release-lock and SHA-256 integrity, safe OOXML and public-data-only boundary."
    )


if __name__ == "__main__":
    main()
