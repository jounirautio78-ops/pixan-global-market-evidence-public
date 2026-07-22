#!/usr/bin/env python3
"""Build the public, bilingual Pixan bank diligence package.

The builder deliberately reads only the repository's sanitised public data in
``site/data``.  It must never be pointed at a private data room. Finnish OOXML
outputs are rebuilt deterministically; reviewed English derivatives are bound
to the same release and Finnish source hashes through a fail-closed lock.
"""

from __future__ import annotations

import argparse
import copy
import csv
import hashlib
import json
import math
import os
import re
import tempfile
import zipfile
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree as ET

from openpyxl import Workbook, load_workbook
from openpyxl.formatting.rule import FormulaRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from bank_register_parity import validate_register_parity
except ModuleNotFoundError:  # Support importing this file as scripts.build_bank_package.
    from scripts.bank_register_parity import validate_register_parity


ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = ROOT / "site" / "data"
DOWNLOAD_DIR = ROOT / "site" / "downloads"

INPUT_FILES = (
    DATA_DIR / "atlas.json",
    DATA_DIR / "market-values.json",
    DATA_DIR / "patent-history.json",
    DATA_DIR / "changelog.json",
)

OUTPUTS = {
    "short-deck": DOWNLOAD_DIR / "pixan-bank-deck-short-fi.pptx",
    "medium-deck": DOWNLOAD_DIR / "pixan-bank-deck-medium-fi.pptx",
    "large-deck": DOWNLOAD_DIR / "pixan-bank-deck-large-fi.pptx",
    "evidence-register": DOWNLOAD_DIR / "pixan-bank-evidence-register-fi.xlsx",
}
CSV_OUTPUT = DATA_DIR / "bank-evidence-register.csv"
EN_OUTPUTS = {
    "short-deck-en": DOWNLOAD_DIR / "pixan-bank-deck-short-en.pptx",
    "medium-deck-en": DOWNLOAD_DIR / "pixan-bank-deck-medium-en.pptx",
    "large-deck-en": DOWNLOAD_DIR / "pixan-bank-deck-large-en.pptx",
    "evidence-register-en": DOWNLOAD_DIR / "pixan-bank-evidence-register-en.xlsx",
}
EN_REGISTER_SOURCE = ROOT / "source" / "bank-evidence-register-en.json"
EN_DECK_TRANSLATIONS_SOURCE = ROOT / "source" / "bank-deck-en-translations.json"
EN_LOCK_SOURCE = ROOT / "source" / "bank-package-en-lock.json"
EN_CSV_OUTPUT = DATA_DIR / "bank-evidence-register-en.csv"
MANIFEST_OUTPUT = DATA_DIR / "bank-package-manifest.json"

REGISTER_HEADERS = [
    "Väite",
    "Dia/osio",
    "Todiste",
    "Lähde",
    "Päivämäärä",
    "Laskentatapa",
    "Oletukset",
    "Luottamustaso",
    "Puutteet / tarvittava lisänäyttö",
]
ALLOWED_STATUSES = {"Vahvistettu", "Tuettu", "Oletus", "Puuttuu"}
EN_REGISTER_HEADERS = [
    "Claim",
    "Slide/section",
    "Evidence",
    "Source",
    "Date",
    "Calculation method",
    "Assumptions",
    "Confidence",
    "Gaps / additional evidence needed",
]
EN_ALLOWED_STATUSES = {"Confirmed", "Supported", "Assumption", "Missing"}

# Restrictive boundary scan.  Generic finance terms are intentionally allowed;
# names and local path fragments from private work are not.
FORBIDDEN_PUBLIC_TERMS = (
    "/Users/",
    "\\Users\\",
    "file://",
)
PRIVATE_IDENTIFIER_FINGERPRINTS = frozenset(
    {
        (7, "46d7415f6182ece9e933e8e9f780957e449361e0dbe10e34f46c186cad3382a1"),
        (7, "f910f0bbe95037851d18ca33b91ee7fc9f334c6cfcd02deaf66af4501c8a884c"),
        (9, "7e6578c2e34b53136741c6efe7799a2dce739651c22404a7894b48d42aa88b41"),
        (13, "933536a17b00f1b39ba9d3585427bd7232d44960ab35754318c1da8e4cf6c5be"),
        (25, "40f45830e7e3e21d88245728fe87f76b2e8919543a502aad248a465487cacee3"),
    }
)

NAVY = "071A2B"
BLUE = "0D5F86"
TEAL = "00A4A6"
PALE = "EAF3F6"
PALE_TEAL = "E3F6F3"
GOLD = "F4B942"
RED = "C84B4B"
GREEN = "138A72"
INK = "182935"
MUTED = "5B6B75"
WHITE = "FFFFFF"
LINE = "CBD8DE"
LIGHT = "F6F9FA"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def parse_iso_date(value: str) -> datetime:
    return datetime.strptime(value[:10], "%Y-%m-%d").replace(tzinfo=timezone.utc)


def euro_m(value: float, currency: str = "€") -> str:
    return f"{value / 1_000_000:,.1f} M{currency}".replace(",", " ")


def billion(value: float, currency: str) -> str:
    return f"{value / 1_000_000_000:,.2f} mrd {currency}".replace(",", " ")


def million_litres(value: float) -> str:
    return f"{format_local_number(value / 1_000_000, 3, decimal_separator='.')} milj. l"


def format_local_number(
    value: int | float,
    decimals: int = 0,
    *,
    decimal_separator: str = ",",
    grouped: bool = True,
) -> str:
    """Format a public-data value deterministically for Finnish register prose."""

    rendered = f"{value:,.{decimals}f}"
    if not grouped:
        rendered = rendered.replace(",", "")
    else:
        rendered = rendered.replace(",", "\N{NO-BREAK SPACE}")
    rendered = rendered.replace(".", decimal_separator)
    return rendered.replace("\N{NO-BREAK SPACE}", " ")


def format_compact_local_number(
    value: int | float,
    *,
    minimum_decimals: int = 0,
    maximum_decimals: int = 2,
    decimal_separator: str = ",",
    grouped: bool = True,
) -> str:
    rendered = format_local_number(
        value,
        maximum_decimals,
        decimal_separator=decimal_separator,
        grouped=grouped,
    )
    if maximum_decimals <= minimum_decimals or decimal_separator not in rendered:
        return rendered
    integer, fraction = rendered.rsplit(decimal_separator, 1)
    fraction = fraction.rstrip("0")
    if len(fraction) < minimum_decimals:
        fraction += "0" * (minimum_decimals - len(fraction))
    return integer if not fraction else integer + decimal_separator + fraction


def observation_year(observation: dict[str, Any]) -> int:
    if isinstance(observation.get("year"), int):
        return int(observation["year"])
    match = re.match(r"^[A-Z]+-(\d{4})-", str(observation.get("observationId", "")))
    if not match:
        raise ValueError(f"Observation year is missing from {observation.get('observationId')!r}")
    return int(match.group(1))


def fi_date(iso_date: str) -> str:
    parsed = datetime.strptime(iso_date, "%Y-%m-%d")
    return f"{parsed.day}.{parsed.month}.{parsed.year}"


def fi_cardinal(value: int) -> str:
    words = {
        0: "nolla",
        1: "yksi",
        2: "kaksi",
        3: "kolme",
        4: "neljä",
        5: "viisi",
        6: "kuusi",
        7: "seitsemän",
        8: "kahdeksan",
        9: "yhdeksän",
        10: "kymmenen",
    }
    return words.get(value, str(value))


def fi_elative_cardinal(value: int) -> str:
    words = {0: "nollasta", 1: "yhdestä", 2: "kahdesta", 3: "kolmesta", 4: "neljästä", 5: "viidestä"}
    return words.get(value, str(value))


def fi_join(values: list[str]) -> str:
    if not values:
        return ""
    if len(values) == 1:
        return values[0]
    return ", ".join(values[:-1]) + " ja " + values[-1]


def required_match(pattern: str, value: str, label: str) -> re.Match[str]:
    match = re.search(pattern, value, flags=re.IGNORECASE)
    if not match:
        raise ValueError(f"Could not derive {label} from public data")
    return match


def proceeding_process_state(item: dict[str, Any]) -> str:
    """Classify a proceeding's appeal/finality state without inventing one."""

    explicit_values = [item.get("status"), item.get("appealStatus")]
    for value in explicit_values:
        normalised = str(value or "").strip().casefold().replace("-", "_").replace(" ", "_")
        if normalised in {"pending", "appeal_pending", "not_final", "open", "on_appeal"}:
            return "pending"
        if normalised in {"final", "closed", "concluded", "appeal_concluded", "no_appeal"}:
            return "final"
        if normalised in {"unknown", "unverified", "not_verified"}:
            return "unverified"

    text = " ".join(
        str(item.get(field) or "")
        for field in ("titleEn", "titleFi", "finalityEn", "finalityFi", "reference")
    ).casefold()
    if "no appeal pending" in text or "ei valitusta vireillä" in text:
        return "final"
    if any(marker in text for marker in ("not final", "appeal pending", "valitus vireillä", "vireillä oleva valitus")):
        return "pending"
    if any(
        marker in text
        for marker in (
            "not independently verified",
            "not verified",
            "ei ole tässä julkaisussa riippumattomasti vahvistettu",
            "vahvistamatta",
        )
    ):
        return "unverified"
    if any(marker in text for marker in ("final judgment", "is final", "lainvoimainen", "closed", "concluded")):
        return "final"
    return "unverified"


def proceeding_appeal_reference(item: dict[str, Any]) -> str | None:
    explicit = item.get("appealReference")
    if isinstance(explicit, str) and explicit.strip():
        return explicit.strip()
    match = re.search(r"(?:appeal|valitus)\s+(.+)$", str(item.get("reference", "")), flags=re.IGNORECASE)
    return match.group(1).strip() if match else None


def required_single(items: Iterable[dict[str, Any]], label: str) -> dict[str, Any]:
    matches = list(items)
    if len(matches) != 1:
        raise ValueError(f"Expected exactly one {label}, got {len(matches)}")
    return matches[0]


def unique_index(items: Iterable[dict[str, Any]], key: str, label: str) -> dict[str, dict[str, Any]]:
    indexed: dict[str, dict[str, Any]] = {}
    for position, item in enumerate(items, start=1):
        identifier = item.get(key)
        if not isinstance(identifier, str) or not identifier:
            raise ValueError(f"{label} item {position} lacks a non-empty {key}")
        if identifier in indexed:
            raise ValueError(f"Duplicate {label} {key}: {identifier}")
        indexed[identifier] = item
    return indexed


def fi_unit(unit: str | None) -> str:
    if not unit:
        raise ValueError("Public observation is missing its unit")
    aliases = {"litre": "l", "unit": "yks."}
    if unit in aliases:
        return aliases[unit]
    per_match = re.fullmatch(r"([A-Z]{3})_per_([a-z]+)", unit)
    if per_match:
        return f"{per_match.group(1)}/{per_match.group(2)}"
    return unit


def fi_finality(finality: str | None) -> str:
    labels = {
        "final": "Lopullinen",
        "provisional": "Alustava",
        "published": "Julkaistu",
        "official_response": "Virallinen vastaus",
        "official_rounded": "Virallinen pyöristys",
        "external_estimate": "Ulkoinen arvio",
        "current_listing": "Nykyinen listahinta",
    }
    if finality not in labels:
        raise ValueError(f"Unsupported public-data finality: {finality!r}")
    return labels[finality]


def concise_publisher(publisher: str) -> str:
    aliases = {
        "Statistisches Bundesamt (Destatis)": "Destatis",
        "Finnish Tax Administration": "Vero",
        "Sejm of the Republic of Poland": "Sejm",
        "Government Offices of Sweden": "Ruotsin hallitus",
        "IMARC Group": "IMARC",
        "Grand View Research": "GVR",
        "Fortune Business Insights": "Fortune",
        "European Patent Office": "EPO",
        "European Patent Office / Agreement on a Unified Patent Court": "EPO / UPC",
        "Finnish Patent and Registration Office": "PRH",
        "United States Patent and Trademark Office": "USPTO",
    }
    return aliases.get(publisher, publisher)


def unique_strings(values: Iterable[str]) -> list[str]:
    return list(dict.fromkeys(value for value in values if value))


def year_label(observations: Iterable[dict[str, Any]]) -> str:
    years = sorted({observation_year(item) for item in observations})
    if not years:
        raise ValueError("Cannot format an empty observation-year range")
    if len(years) == 1:
        return str(years[0])
    if years == list(range(years[0], years[-1] + 1)):
        return f"{years[0]}–{years[-1]}"
    return " / ".join(str(year) for year in years)


def format_billions(value: int | float, currency: str) -> str:
    return (
        f"{format_compact_local_number(value / 1_000_000_000, minimum_decimals=1, grouped=False)} "
        f"mrd {currency}"
    )


def format_millions(value: int | float, unit: str) -> str:
    return (
        f"{format_compact_local_number(value / 1_000_000, maximum_decimals=3)} "
        f"milj. {unit}"
    )


def public_text_scan(values: Iterable[str]) -> None:
    texts = list(values)
    joined = "\n".join(texts)
    for term in FORBIDDEN_PUBLIC_TERMS:
        if term.casefold() in joined.casefold():
            raise ValueError("Forbidden private/public-boundary path marker found")
    for text in texts:
        normalised = re.sub(r"[^a-z0-9]+", "", text.casefold())
        for length, expected in PRIVATE_IDENTIFIER_FINGERPRINTS:
            if any(
                hashlib.sha256(normalised[index:index + length].encode("utf-8")).hexdigest() == expected
                for index in range(max(0, len(normalised) - length + 1))
            ):
                raise ValueError("Forbidden private identifier fingerprint found")


def source_url(source_id: str, market_sources: dict[str, Any], patent_sources: dict[str, Any]) -> str:
    if source_id in market_sources and source_id in patent_sources:
        raise ValueError(f"Ambiguous public source id appears in market and patent data: {source_id}")
    source = market_sources.get(source_id) or patent_sources.get(source_id)
    if not source:
        raise ValueError(f"Unknown public source id: {source_id}")
    url = source.get("pageUrl") or source.get("url")
    if not isinstance(url, str) or not url.strip():
        raise ValueError(f"Public source lacks a usable URL: {source_id}")
    return url


def build_context() -> dict[str, Any]:
    atlas = read_json(DATA_DIR / "atlas.json")
    market = read_json(DATA_DIR / "market-values.json")
    patent = read_json(DATA_DIR / "patent-history.json")
    changelog = read_json(DATA_DIR / "changelog.json")

    if not changelog.get("releases"):
        raise ValueError("changelog.json must contain at least one release")
    release = max(changelog["releases"], key=lambda item: item["publishedAt"])
    as_of = changelog["asOf"]
    if any(data.get("meta", {}).get("asOf", data.get("asOf")) != as_of for data in (atlas, market, patent)):
        raise ValueError("Public inputs do not share the changelog as-of date")

    observations = unique_index(market["observations"], "observationId", "market observation")
    models = unique_index(market["models"], "modelId", "market model")
    proceedings = unique_index(patent["proceedings"], "proceedingId", "patent proceeding")
    alerts = unique_index(patent["diligenceAlerts"], "alertId", "diligence alert")
    market_sources = unique_index(market["sources"], "sourceId", "market source")
    patent_sources = unique_index(patent["sources"], "sourceId", "patent source")
    source_id_collisions = sorted(set(market_sources) & set(patent_sources))
    if source_id_collisions:
        raise ValueError(
            "Market and patent source ids must be globally unique: " + ", ".join(source_id_collisions)
        )

    official_country_codes = sorted(
        {
            item["countryIso2"]
            for item in market["observations"]
            if item["geography"] != "Global" and item["evidenceStatus"].startswith("official")
        }
    )
    retail_donors = int(market["meta"]["modelReadiness"]["comparableFullYearMarketValueDonors"])
    grade_counts = atlas["summary"]["gradeCounts"]

    return {
        "atlas": atlas,
        "market": market,
        "patent_history": patent,
        "changelog": changelog,
        "release": release,
        "as_of": as_of,
        "observations": observations,
        "models": models,
        "proceedings": proceedings,
        "alerts": alerts,
        "market_sources": market_sources,
        "patent_sources": patent_sources,
        "official_country_codes": official_country_codes,
        "retail_donors": retail_donors,
        "grade_counts": grade_counts,
    }


def canonical_facts(ctx: dict[str, Any]) -> dict[str, Any]:
    """Derive every mutable bank-deck fact from the public canonical inputs."""

    patent_history = ctx["patent_history"]
    patent = patent_history["patent"]
    observations = list(ctx["market"]["observations"])
    observations_by_id = unique_index(observations, "observationId", "market observation")
    unique_index(ctx["market"]["models"], "modelId", "market model")
    unique_index(patent_history["proceedings"], "proceedingId", "patent proceeding")
    unique_index(patent_history["diligenceAlerts"], "alertId", "diligence alert")
    unique_index(patent_history["familyMembers"], "publicationNumber", "patent family member")
    unique_index(patent_history["monetisation"]["sequence"], "phaseId", "monetisation phase")
    market_sources = unique_index(ctx["market"]["sources"], "sourceId", "market source")
    patent_sources = unique_index(patent_history["sources"], "sourceId", "patent source")
    source_id_collisions = sorted(set(market_sources) & set(patent_sources))
    if source_id_collisions:
        raise ValueError(
            "Market and patent source ids must be globally unique: " + ", ".join(source_id_collisions)
        )
    atlas_by_iso = unique_index(ctx["atlas"]["countries"], "iso2", "atlas country")

    def require_source_ids(item: dict[str, Any], source_map: dict[str, Any], label: str) -> None:
        source_ids = item.get("sourceIds")
        if not isinstance(source_ids, list) or not source_ids or len(source_ids) != len(set(source_ids)):
            raise ValueError(f"{label} must have a non-empty, duplicate-free sourceIds list")
        missing = [source_id for source_id in source_ids if source_id not in source_map]
        if missing:
            raise ValueError(f"{label} refers to unknown source ids: {', '.join(missing)}")

    def require_observation(
        item: dict[str, Any],
        label: str,
        *,
        country: str | None,
        geography: str,
        metric: str,
        period: str,
        unit: str,
        currency: str | None,
        evidence_status: str | set[str],
        finality: str | set[str],
    ) -> None:
        expected_statuses = {evidence_status} if isinstance(evidence_status, str) else evidence_status
        expected_finalities = {finality} if isinstance(finality, str) else finality
        expected = {
            "countryIso2": country,
            "geography": geography,
            "metric": metric,
            "period": period,
            "unit": unit,
            "currency": currency,
        }
        for field, value in expected.items():
            if item.get(field) != value:
                raise ValueError(f"{label} {field} must be {value!r}, got {item.get(field)!r}")
        if item.get("evidenceStatus") not in expected_statuses:
            raise ValueError(
                f"{label} evidenceStatus must be one of {sorted(expected_statuses)!r}, "
                f"got {item.get('evidenceStatus')!r}"
            )
        if item.get("finality") not in expected_finalities:
            raise ValueError(
                f"{label} finality must be one of {sorted(expected_finalities)!r}, "
                f"got {item.get('finality')!r}"
            )
        if (
            not isinstance(item.get("value"), (int, float))
            or isinstance(item.get("value"), bool)
            or not math.isfinite(float(item["value"]))
            or float(item["value"]) <= 0
        ):
            raise ValueError(f"{label} must have a positive finite numeric value")
        require_source_ids(item, market_sources, label)

    for item in patent_history["familyMembers"]:
        require_source_ids(
            item,
            patent_sources,
            f"patent family member {item.get('publicationNumber')!r}",
        )

    def market_observations(country: str, metric: str) -> list[dict[str, Any]]:
        return sorted(
            (
                item
                for item in observations
                if item.get("countryIso2") == country and item.get("metric") == metric
            ),
            key=lambda item: (observation_year(item), item["observationId"]),
        )

    def market_observation(country: str, metric: str) -> dict[str, Any]:
        matches = market_observations(country, metric)
        if not matches:
            raise ValueError(f"Expected at least one {country} {metric} observation")
        latest_year = max(observation_year(item) for item in matches)
        return required_single(
            (item for item in matches if observation_year(item) == latest_year),
            f"latest-year {country} {metric} observation",
        )

    def market_observation_for_year(country: str, metric: str, year: int) -> dict[str, Any]:
        return required_single(
            (item for item in market_observations(country, metric) if observation_year(item) == year),
            f"{year} {country} {metric} observation",
        )

    def proceeding(
        jurisdiction: str,
        proceeding_type: str,
        *,
        required_evidence_tier: str | None = None,
        required_source_kind: str | None = None,
    ) -> dict[str, Any]:
        matches = [
            item
            for item in patent_history["proceedings"]
            if item.get("jurisdictionCode") == jurisdiction
            and item.get("proceedingType") == proceeding_type
        ]
        if required_evidence_tier is not None or required_source_kind is not None:
            matches = [
                item
                for item in matches
                if (required_evidence_tier is None or item.get("evidenceTier") == required_evidence_tier)
                and any(
                    (required_evidence_tier is None or patent_sources[source_id].get("evidenceTier") == required_evidence_tier)
                    and (required_source_kind is None or patent_sources[source_id].get("sourceKind") == required_source_kind)
                    for source_id in item.get("sourceIds", [])
                    if source_id in patent_sources
                )
            ]
        if not matches:
            source_requirement = ""
            if required_evidence_tier or required_source_kind:
                source_requirement = (
                    f" with source tier {required_evidence_tier!r}"
                    f" and kind {required_source_kind!r}"
                )
            raise ValueError(
                f"Expected at least one {jurisdiction} {proceeding_type} proceeding{source_requirement}"
            )
        latest_date = max(item["eventDate"] for item in matches)
        selected = required_single(
            (item for item in matches if item["eventDate"] == latest_date),
            f"latest {jurisdiction} {proceeding_type} proceeding",
        )
        require_source_ids(selected, patent_sources, f"{jurisdiction} {proceeding_type} proceeding")
        return selected

    def country_name(code: str) -> str:
        if code == "EP":
            return "Eurooppa"
        country = atlas_by_iso.get(code)
        if not country:
            raise ValueError(f"Atlas lacks Finnish country name for {code}")
        return country["nameFi"]

    def publishers(source_ids: Iterable[str], source_map: dict[str, Any]) -> list[str]:
        return unique_strings(concise_publisher(source_map[source_id]["publisher"]) for source_id in source_ids)

    def observation_publishers(items: Iterable[dict[str, Any]]) -> list[str]:
        return publishers(
            (source_id for item in items for source_id in item["sourceIds"]),
            market_sources,
        )

    b2_claim_count = int(
        required_match(
            r"B2 specification contains (\d+) claims",
            patent["claimScopeSummaryEn"],
            "amended B2 claim count",
        ).group(1)
    )
    epo_proceeding = proceeding("EP", "opposition_and_appeal", required_evidence_tier="official")
    de_nullity = proceeding(
        "DE",
        "patent_nullity",
        required_evidence_tier="official",
        required_source_kind="official_judgment",
    )
    de_infringement = proceeding(
        "DE",
        "patent_infringement",
        required_evidence_tier="official",
        required_source_kind="official_judgment",
    )
    cn_review = proceeding("CN", "review_request")
    infringement_match = required_match(
        r"claims (\d+) and (\d+)",
        de_infringement["detailEn"],
        "German judgment claim numbers",
    )
    infringement_claims = [int(infringement_match.group(1)), int(infringement_match.group(2))]
    nullity_appeal_reference = proceeding_appeal_reference(de_nullity)
    german_official_proceedings = sorted(
        (
            item
            for item in patent_history["proceedings"]
            if item.get("jurisdictionCode") == "DE"
            and item.get("evidenceTier") == "official"
            and any(
                patent_sources[source_id].get("evidenceTier") == "official"
                and patent_sources[source_id].get("sourceKind") == "official_judgment"
                for source_id in item.get("sourceIds", [])
                if source_id in patent_sources
            )
        ),
        key=lambda item: (item["eventDate"], item["proceedingId"]),
    )
    if de_nullity not in german_official_proceedings or de_infringement not in german_official_proceedings:
        raise ValueError("Latest German nullity and infringement narratives must belong to official judgment records")

    cn_grant = required_single(
        (
            item
            for item in patent_history["familyMembers"]
            if item.get("jurisdictionCode") == "CN"
            and item.get("recordType") == "national_grant_publication"
        ),
        "Chinese national grant publication",
    )
    national_register_members = [
        item
        for item in patent_history["familyMembers"]
        if item.get("verificationLevel") == "official_national_record"
        and any(
            patent_sources[source_id].get("evidenceTier") == "official"
            and patent_sources[source_id].get("sourceKind")
            in {"official_national_register", "official_national_api"}
            for source_id in item.get("sourceIds", [])
            if source_id in patent_sources
        )
    ]
    national_register_jurisdictions = sorted(
        {item["jurisdictionCode"] for item in national_register_members}
    )

    universe_match = required_match(
        r"^UN(\d+)\+VA\+PS$",
        ctx["atlas"]["summary"]["universe"],
        "UN research-universe components",
    )
    official_country_codes = sorted(
        {
            item["countryIso2"]
            for item in observations
            if item.get("geography") != "Global"
            and isinstance(item.get("countryIso2"), str)
            and str(item.get("evidenceStatus", "")).startswith("official")
        }
    )
    official_country_names_fi = [country_name(code) for code in official_country_codes]

    ca_value = market_observation("CA", "manufacturer_importer_shipments_value")
    ca_year = observation_year(ca_value)
    ca_units = market_observation_for_year("CA", "manufacturer_importer_shipments_units", ca_year)
    ca_litres = market_observation_for_year("CA", "manufacturer_importer_shipments_liquid_volume", ca_year)
    require_observation(
        ca_value,
        "Canada shipment-value observation",
        country="CA",
        geography="Canada",
        metric="manufacturer_importer_shipments_value",
        period="calendar_year",
        unit="CAD",
        currency="CAD",
        evidence_status="official_observed",
        finality="published",
    )
    require_observation(
        ca_units,
        "Canada shipment-unit observation",
        country="CA",
        geography="Canada",
        metric="manufacturer_importer_shipments_units",
        period="calendar_year",
        unit="unit",
        currency=None,
        evidence_status="official_observed",
        finality="published",
    )
    require_observation(
        ca_litres,
        "Canada shipment-liquid observation",
        country="CA",
        geography="Canada",
        metric="manufacturer_importer_shipments_liquid_volume",
        period="calendar_year",
        unit="litre",
        currency=None,
        evidence_status="official_observed",
        finality="published",
    )
    de_volumes = market_observations("DE", "taxed_substitutes_volume")
    de_excise = market_observations("DE", "substitutes_excise_receipts")
    if not de_volumes or {observation_year(item) for item in de_volumes} != {
        observation_year(item) for item in de_excise
    }:
        raise ValueError("Germany volume and excise series must cover the same years")
    de_volume_by_year = {
        year: required_single(
            (item for item in de_volumes if observation_year(item) == year),
            f"Germany {year} taxed-liquid volume",
        )
        for year in {observation_year(item) for item in de_volumes}
    }
    de_excise_by_year = {
        year: required_single(
            (item for item in de_excise if observation_year(item) == year),
            f"Germany {year} substitutes-excise receipt",
        )
        for year in {observation_year(item) for item in de_excise}
    }
    if any(
        de_volume_by_year[year].get("finality") != de_excise_by_year[year].get("finality")
        for year in de_volume_by_year
    ):
        raise ValueError("Germany volume and excise finality must match within each year")
    german_status_by_finality = {"final": "official_observed", "provisional": "official_provisional"}
    for year, volume in de_volume_by_year.items():
        finality = volume.get("finality")
        if finality not in german_status_by_finality:
            raise ValueError(f"Germany {year} volume has unsupported finality {finality!r}")
        require_observation(
            volume,
            f"Germany {year} taxed-liquid volume",
            country="DE",
            geography="Germany",
            metric="taxed_substitutes_volume",
            period="calendar_year",
            unit="litre",
            currency=None,
            evidence_status=german_status_by_finality[finality],
            finality=finality,
        )
        require_observation(
            de_excise_by_year[year],
            f"Germany {year} substitutes-excise receipt",
            country="DE",
            geography="Germany",
            metric="substitutes_excise_receipts",
            period="calendar_year",
            unit="EUR",
            currency="EUR",
            evidence_status=german_status_by_finality[finality],
            finality=finality,
        )

    fi_volume = market_observation("FI", "nicotine_e_liquid_taxed_volume")
    fi_excise = market_observation("FI", "nicotine_e_liquid_excise_receipts")
    pl_volume = market_observation("PL", "reported_e_liquid_volume")
    pl_excise = market_observation("PL", "e_liquid_excise_amount")
    se_volume = market_observation("SE", "nicotine_e_liquid_taxed_volume")
    se_excise = market_observation("SE", "nicotine_e_liquid_excise_receipts")
    national_observation_schemas = (
        (fi_volume, "Finland taxed-liquid volume", "FI", "Finland", "nicotine_e_liquid_taxed_volume", "litre", None, "published"),
        (fi_excise, "Finland liquid-excise receipt", "FI", "Finland", "nicotine_e_liquid_excise_receipts", "EUR", "EUR", "published"),
        (pl_volume, "Poland reported-liquid volume", "PL", "Poland", "reported_e_liquid_volume", "litre", None, "official_response"),
        (pl_excise, "Poland liquid-excise amount", "PL", "Poland", "e_liquid_excise_amount", "PLN", "PLN", "official_response"),
        (se_volume, "Sweden taxed-liquid volume", "SE", "Sweden", "nicotine_e_liquid_taxed_volume", "litre", None, "official_rounded"),
        (se_excise, "Sweden liquid-excise receipt", "SE", "Sweden", "nicotine_e_liquid_excise_receipts", "SEK", "SEK", "official_rounded"),
    )
    for item, label, country, geography, metric, unit, currency, finality in national_observation_schemas:
        require_observation(
            item,
            label,
            country=country,
            geography=geography,
            metric=metric,
            period="calendar_year",
            unit=unit,
            currency=currency,
            evidence_status="official_observed",
            finality=finality,
        )

    model_candidates = [
        item
        for item in ctx["market"]["models"]
        if item.get("countryIso2") == "DE"
        and item.get("atlasEstimate") is True
        and item.get("marketValueBasis") == "retail_equivalent_plausibility_range_not_observed_sales"
    ]
    if not model_candidates:
        raise ValueError("Expected at least one German taxed-liquid retail-equivalent model")
    latest_model_year = max(int(item["year"]) for item in model_candidates)
    model = required_single(
        (item for item in model_candidates if int(item["year"]) == latest_model_year),
        "latest German taxed-liquid retail-equivalent model",
    )
    supported_model_fields = {
        "countryIso2": "DE",
        "evidenceStatus": "modelled",
        "confidence": "low",
        "productScope": "taxed_substitutes_for_tobacco_liquid_only",
        "marketValueBasis": "retail_equivalent_plausibility_range_not_observed_sales",
        "comparableMarketValue": False,
        "atlasEstimate": True,
        "formula": "volume_litres * 1000 * retail_price_eur_per_ml",
        "currency": "EUR",
    }
    for field, expected in supported_model_fields.items():
        if model.get(field) != expected:
            raise ValueError(f"Germany model {field} must be {expected!r}, got {model.get(field)!r}")
    if not isinstance(model.get("year"), int) or isinstance(model.get("year"), bool):
        raise ValueError("Germany model year must be an integer")
    range_input_map = model.get("rangeInputMap")
    if not isinstance(range_input_map, dict) or list(range_input_map) != ["low", "central", "high"]:
        raise ValueError("Germany model rangeInputMap must be ordered low, central, high")
    price_input_ids = [range_input_map[key] for key in ("low", "central", "high")]
    if any(not isinstance(input_id, str) or not input_id for input_id in price_input_ids):
        raise ValueError("Germany model price-input ids must be non-empty strings")
    if len(set(price_input_ids)) != 3:
        raise ValueError("Germany model low, central and high must use three distinct price inputs")
    input_ids = model.get("inputIds")
    if (
        not isinstance(input_ids, list)
        or any(not isinstance(input_id, str) or not input_id for input_id in input_ids)
        or len(input_ids) != len(set(input_ids))
    ):
        raise ValueError("Germany model inputIds must be a duplicate-free list")
    missing_input_ids = [input_id for input_id in input_ids if input_id not in observations_by_id]
    if missing_input_ids:
        raise ValueError(f"Germany model refers to unknown input ids: {', '.join(missing_input_ids)}")
    if len(input_ids) != 4 or input_ids[1:] != price_input_ids:
        raise ValueError(
            "Germany model inputIds must contain one volume input followed by low, central and high price inputs"
        )
    if input_ids[0] in price_input_ids:
        raise ValueError("Germany model volume input cannot also be a range price input")
    model_prices = {
        key: observations_by_id[range_input_map[key]]
        for key in ("low", "central", "high")
    }
    model_volume = observations_by_id[input_ids[0]]
    volume_finality = model_volume.get("finality")
    if volume_finality not in german_status_by_finality:
        raise ValueError(f"Germany model volume has unsupported finality {volume_finality!r}")
    require_observation(
        model_volume,
        "Germany model volume input",
        country="DE",
        geography="Germany",
        metric="taxed_substitutes_volume",
        period="calendar_year",
        unit="litre",
        currency=None,
        evidence_status=german_status_by_finality[volume_finality],
        finality=volume_finality,
    )
    model_year = int(model["year"])
    volume_year = observation_year(model_volume)
    if volume_year != model_year:
        raise ValueError("Germany model year must equal its taxed-liquid volume-input year")
    for key, price in model_prices.items():
        require_observation(
            price,
            f"Germany model {key} price input",
            country="DE",
            geography="Germany",
            metric="retail_price_input",
            period="point_in_time",
            unit="EUR_per_ml",
            currency="EUR",
            evidence_status="published_price_input",
            finality="current_listing",
        )
    model_price_years = {observation_year(item) for item in model_prices.values()}
    if len(model_price_years) != 1:
        raise ValueError("Germany model must use one deterministic retail-price year")
    has_price_year_mismatch = any(
        observation_year(price) != model_year for price in model_prices.values()
    )
    if has_price_year_mismatch != (model.get("yearMismatch") is True):
        raise ValueError("Germany model yearMismatch must exactly describe retail-price year mismatches")
    price_values = [float(model_prices[key]["value"]) for key in ("low", "central", "high")]
    if not price_values[0] < price_values[1] < price_values[2]:
        raise ValueError("Germany model retail-price inputs must increase low < central < high")
    for key in ("low", "central", "high"):
        if (
            not isinstance(model.get(key), (int, float))
            or isinstance(model.get(key), bool)
            or not math.isfinite(float(model[key]))
            or float(model[key]) <= 0
        ):
            raise ValueError(f"Germany model {key} total must be a positive finite number")
    for key, price in model_prices.items():
        calculated = model_volume["value"] * 1_000 * price["value"]
        if not math.isclose(float(model[key]), float(calculated), rel_tol=0, abs_tol=0.5):
            raise ValueError(f"Germany model {key} total differs from volume × 1,000 × price")
    if not float(model["low"]) < float(model["central"]) < float(model["high"]):
        raise ValueError("Germany model totals must increase low < central < high")
    model_source_ids = unique_strings(
        source_id
        for input_id in input_ids
        for source_id in observations_by_id[input_id]["sourceIds"]
    )

    global_observations = [item for item in observations if item.get("geography") == "Global"]
    if not global_observations:
        raise ValueError("At least one public commercial global estimate is required")
    for item in global_observations:
        require_observation(
            item,
            "Global commercial market estimate",
            country=None,
            geography="Global",
            metric="commercial_market_estimate",
            period="calendar_year_estimate",
            unit="USD",
            currency="USD",
            evidence_status="commercial_estimate",
            finality="external_estimate",
        )
    global_years = {observation_year(item) for item in global_observations}
    global_currencies = {item.get("currency") for item in global_observations}
    if len(global_years) != 1 or global_currencies != {"USD"}:
        raise ValueError("Commercial global estimates must share one explicit year and USD currency")

    alerts = list(patent_history["diligenceAlerts"])
    for item in alerts:
        require_source_ids(item, patent_sources, f"diligence alert {item.get('alertId')!r}")
    au_alert = required_single(
        (item for item in alerts if item.get("jurisdictionCode") == "AU"),
        "Australian diligence alert",
    )
    fi_alert = required_single(
        (item for item in alerts if item.get("jurisdictionCode") == "FI"),
        "Finnish diligence alert",
    )
    us_alert = required_single(
        (item for item in alerts if item.get("jurisdictionCode") == "US"),
        "US diligence alert",
    )
    fi_fee_year = int(
        required_match(r"year-(\d+) fee", fi_alert["detailEn"], "Finnish renewal-fee year").group(1)
    )
    us_maintenance_year = float(
        required_match(r"(\d+(?:\.\d+)?)-year", us_alert["titleEn"], "US maintenance-fee year").group(1)
    )
    finance_phase = next(
        (
            index
            for index, phase in enumerate(patent_history["monetisation"]["sequence"], start=1)
            if phase["phaseId"] == "PHASE-5-FINANCE"
        ),
        None,
    )
    if finance_phase is None:
        raise ValueError("Finance phase is missing from public patent-history data")

    summary = patent_history["summary"]
    if summary["familyRecordCount"] != len(patent_history["familyMembers"]):
        raise ValueError("Patent family summary count differs from familyMembers")
    if summary["proceedingCount"] != len(patent_history["proceedings"]):
        raise ValueError("Patent proceeding summary count differs from proceedings")
    if summary["diligenceAlertCount"] != len(alerts):
        raise ValueError("Patent alert summary count differs from diligenceAlerts")

    market_country_rows: list[list[str]] = []
    official_observations = [
        item
        for item in observations
        if item.get("countryIso2") in official_country_codes
        and str(item.get("evidenceStatus", "")).startswith("official")
    ]
    for code in official_country_codes:
        country_items = [item for item in official_observations if item.get("countryIso2") == code]
        shipment_values = [item for item in country_items if item.get("metric") == "manufacturer_importer_shipments_value"]
        liquid_volumes = [item for item in country_items if item.get("unit") == "litre"]
        liquid_excise = [
            item
            for item in country_items
            if "excise" in str(item.get("metric"))
            and "device" not in str(item.get("metric"))
            and "component" not in str(item.get("metric"))
        ]
        if shipment_values:
            shipment_year = max(observation_year(item) for item in shipment_values)
            shipment = required_single(
                (item for item in shipment_values if observation_year(item) == shipment_year),
                f"latest-year {code} shipment-value observation",
            )
            liquid = required_single(
                (item for item in liquid_volumes if observation_year(item) == shipment_year),
                f"{shipment_year} {code} liquid-volume observation",
            )
            summary_text = (
                f"{format_billions(shipment['value'], shipment['currency'])} toimitusarvo; "
                f"{format_local_number(liquid['value'])} {fi_unit(liquid['unit'])} nestettä"
            )
        elif len(liquid_volumes) > 1:
            liquid_volumes.sort(key=observation_year)
            summary_text = (
                f"{format_millions(liquid_volumes[0]['value'], fi_unit(liquid_volumes[0]['unit']))} → "
                f"{format_millions(liquid_volumes[-1]['value'], fi_unit(liquid_volumes[-1]['unit']))} "
                "verotettua nestettä"
            )
        else:
            liquid = required_single(liquid_volumes, f"{code} liquid-volume observation")
            excise = required_single(liquid_excise, f"{code} liquid-excise observation")
            summary_text = (
                f"{format_compact_local_number(liquid['value'], maximum_decimals=3)} {fi_unit(liquid['unit'])}; "
                f"{format_millions(excise['value'], excise['currency'])} valmisteveroa"
            )
        market_country_rows.append([country_name(code), year_label(country_items), summary_text])

    return {
        "patent": patent,
        "market_source_map": market_sources,
        "patent_source_map": patent_sources,
        "b2_claim_count": b2_claim_count,
        "priority_date": patent["earliestPriorityDate"],
        "priority_year": int(patent["earliestPriorityDate"][:4]),
        "b2_date": patent["epCentralStatusDate"],
        "ep_publication": patent["epPublication"],
        "epo_proceeding": epo_proceeding,
        "de_nullity": de_nullity,
        "de_infringement": de_infringement,
        "cn_review": cn_review,
        "cn_grant": cn_grant,
        "infringement_claims": infringement_claims,
        "nullity_appeal_reference": nullity_appeal_reference,
        "de_nullity_process_state": proceeding_process_state(de_nullity),
        "de_infringement_process_state": proceeding_process_state(de_infringement),
        "german_official_proceedings": german_official_proceedings,
        "german_official_proceeding_count": len(german_official_proceedings),
        "family_count": summary["familyRecordCount"],
        "official_source_count": summary["officialSourceCount"],
        "proceeding_count": summary["proceedingCount"],
        "alert_count": summary["diligenceAlertCount"],
        "national_register_members": national_register_members,
        "national_register_jurisdictions": national_register_jurisdictions,
        "national_register_check_count": len(national_register_jurisdictions),
        "atlas_country_count": ctx["atlas"]["summary"]["countryCount"],
        "atlas_universe": ctx["atlas"]["summary"]["universe"],
        "un_member_count": int(universe_match.group(1)),
        "atlas_evidence_count": ctx["atlas"]["summary"]["evidenceCount"],
        "grade_counts": dict(ctx["atlas"]["summary"]["gradeCounts"]),
        "grade_labels": sorted(ctx["atlas"]["summary"]["gradeCounts"]),
        "official_country_codes": official_country_codes,
        "official_country_names_fi": official_country_names_fi,
        "official_country_count": len(official_country_codes),
        "retail_donors": int(ctx["market"]["meta"]["modelReadiness"]["comparableFullYearMarketValueDonors"]),
        "minimum_required_donors": int(ctx["market"]["meta"]["modelReadiness"]["minimumRequiredDonors"]),
        "market_readiness": ctx["market"]["meta"]["modelReadiness"],
        "country_name": country_name,
        "ca_value_observation": ca_value,
        "ca_units_observation": ca_units,
        "ca_litres_observation": ca_litres,
        "de_volume_observations": de_volumes,
        "de_excise_observations": de_excise,
        "de_excise_by_year": de_excise_by_year,
        "fi_volume_observation": fi_volume,
        "fi_excise_observation": fi_excise,
        "pl_volume_observation": pl_volume,
        "pl_excise_observation": pl_excise,
        "se_volume_observation": se_volume,
        "se_excise_observation": se_excise,
        "model": model,
        "model_prices": model_prices,
        "model_price_year": next(iter(model_price_years)),
        "model_input_years": sorted(
            {observation_year(observations_by_id[input_id]) for input_id in input_ids}
        ),
        "model_has_year_mismatch": has_price_year_mismatch,
        "model_volume_observation": model_volume,
        "model_source_ids": model_source_ids,
        "global_observations": global_observations,
        "global_values": [item["value"] for item in global_observations],
        "global_year": next(iter(global_years)),
        "global_currency": next(iter(global_currencies)),
        "global_publishers": observation_publishers(global_observations),
        "global_source_ids": [source_id for item in global_observations for source_id in item["sourceIds"]],
        "market_country_rows": market_country_rows,
        "market_publishers": observation_publishers(official_observations),
        "alerts": alerts,
        "au_alert": au_alert,
        "fi_alert": fi_alert,
        "us_alert": us_alert,
        "alert_rows": [
            [country_name(item["jurisdictionCode"]), fi_date(item["targetDate"]), item["titleFi"]]
            for item in alerts
        ],
        "alert_publishers": publishers(
            (source_id for item in alerts for source_id in item["sourceIds"]),
            patent_sources,
        ),
        "fi_fee_year": fi_fee_year,
        "us_maintenance_year": us_maintenance_year,
        "finance_phase": finance_phase,
    }


def evidence_rows(ctx: dict[str, Any]) -> list[dict[str, str]]:
    facts = canonical_facts(ctx)
    p = facts["patent"]
    model = facts["model"]
    market_sources = facts["market_source_map"]
    patent_sources = facts["patent_source_map"]
    as_of = ctx["as_of"]

    def sources(*ids: str) -> str:
        return " ; ".join(source_url(item, market_sources, patent_sources) for item in ids)

    def row(
        claim: str,
        section: str,
        proof: str,
        source: str,
        date: str,
        method: str,
        assumptions: str,
        status: str,
        gap: str,
    ) -> dict[str, str]:
        if status not in ALLOWED_STATUSES:
            raise ValueError(f"Invalid evidence status: {status}")
        return dict(zip(REGISTER_HEADERS, (claim, section, proof, source, date, method, assumptions, status, gap)))

    de_volume_observations = facts["de_volume_observations"]
    de23 = de_volume_observations[0]
    de25 = de_volume_observations[-1]
    de_latest_excise = facts["de_excise_by_year"][observation_year(de25)]
    ca_value = facts["ca_value_observation"]
    ca_litres = facts["ca_litres_observation"]
    fi_volume = facts["fi_volume_observation"]
    fi_excise = facts["fi_excise_observation"]
    pl_volume = facts["pl_volume_observation"]
    pl_excise = facts["pl_excise_observation"]
    se_volume = facts["se_volume_observation"]
    epo_proceeding = facts["epo_proceeding"]
    de_nullity = facts["de_nullity"]
    de_infringement = facts["de_infringement"]
    cn_review = facts["cn_review"]
    de_final_years = [str(observation_year(item)) for item in de_volume_observations if item.get("finality") == "final"]
    de_provisional_years = [
        str(observation_year(item)) for item in de_volume_observations if item.get("finality") == "provisional"
    ]
    if not de_final_years or not de_provisional_years:
        raise ValueError("Germany volume series must expose final and provisional years")
    de_volume_status_proof = (
        f"{fi_join(de_final_years)} {'lopullisia' if len(de_final_years) > 1 else 'lopullinen'}; "
        f"{fi_join(de_provisional_years)} {'alustavia' if len(de_provisional_years) > 1 else 'alustava'}."
    )
    global_observations = facts["global_observations"]
    global_values = facts["global_values"]
    global_year = facts["global_year"]
    global_currency = facts["global_currency"]
    global_source_ids = facts["global_source_ids"]

    def global_publisher(item: dict[str, Any]) -> str:
        publisher = market_sources[item["sourceIds"][0]]["publisher"]
        return "IMARC" if publisher == "IMARC Group" else publisher

    b2_claim_count = facts["b2_claim_count"]
    infringement_claims = facts["infringement_claims"]
    un_member_count = facts["un_member_count"]
    official_country_names_fi = facts["official_country_names_fi"]
    model_price_year = facts["model_price_year"]
    finance_phase = facts["finance_phase"]
    fi_fee_year = facts["fi_fee_year"]
    us_maintenance_year = facts["us_maintenance_year"]
    model_input_year_label = fi_join([str(year) for year in facts["model_input_years"]])
    if facts["model_has_year_mismatch"]:
        model_date_label = f"{model['year']} määrä / {model_price_year} hinnat"
        model_assumptions = (
            "Kolme yksittäistä verkkokauppahintaa; "
            f"syötevuodet {model_input_year_label}; vain verotettu neste."
        )
        model_gap = (
            "Tarvitaan edustava hintakori, tuotemix ja kanavamarginaalit sekä "
            f"mallivuoden {model['year']} hinnat."
        )
    else:
        model_date_label = f"{model['year']} määrä ja hinnat"
        model_assumptions = (
            "Kolme yksittäistä verkkokauppahintaa; "
            f"kaikki syötteet vuodelta {model['year']}; vain verotettu neste."
        )
        model_gap = "Tarvitaan edustava hintakori, tuotemix ja kanavamarginaalit."

    def legal_assumption(item: dict[str, Any]) -> str:
        finality = item.get("finalityFi")
        if not isinstance(finality, str) or not finality.strip():
            raise ValueError(f"Proceeding {item.get('proceedingId')!r} lacks finalityFi")
        return f"Maantieteellinen ja prosessuaalinen rajaus: Saksa. {finality.strip()}"

    def legal_event_title(item: dict[str, Any], process_state: str) -> str:
        title = str(item.get("titleFi") or "").strip()
        if not title:
            raise ValueError(f"Proceeding {item.get('proceedingId')!r} lacks titleFi")
        if process_state == "final" and any(
            marker in title.casefold()
            for marker in ("vireillä", "ei lainvoimainen", "pending", "not final")
        ):
            detail = str(item.get("detailFi") or "").strip()
            if not detail:
                raise ValueError(f"Final proceeding {item.get('proceedingId')!r} lacks a usable detailFi")
            return detail.split(".", 1)[0]
        return title

    def legal_gap(item: dict[str, Any], process_state: str) -> str:
        if process_state == "pending":
            appeal_reference = proceeding_appeal_reference(item)
            reference_text = f" {appeal_reference}" if appeal_reference else ""
            return (
                f"Seuraa muutoksenhakua{reference_text}, sen lopputulosta, täytäntöönpanoa "
                "ja mahdollisia maksettuja korvauksia."
            )
        if process_state == "final":
            return (
                "Lopputulos on merkitty lopulliseksi; varmista täytäntöönpano, "
                "korvausten toteutuminen ja myöhemmät rekisterimuutokset."
            )
        return (
            "Vahvista valitustila, lainvoimaisuus, täytäntöönpano ja maksetut korvaukset."
        )

    rows = [
        row(
            "Julkisen rekisteriaineiston mukaan patentin kirjattu haltija on Pixan Oy.",
            "IP-status",
            f"EPO:n päätietueessa kirjattu haltija: {p['recordedProprietor']}.",
            sources("EPO-REGISTER-MAIN"),
            ctx["patent_history"]["meta"].get("asOf", as_of),
            "Suora rekisterihavainto.",
            "Kirjaus vastaa tarkastelupäivän julkista tietuetta.",
            "Vahvistettu",
            "Täydellinen omistusketju, siirtokirjat ja rasitteet on tarkastettava erikseen.",
        ),
        row(
            f"Patenttiperheen varhaisin prioriteetti on {p['earliestPriorityNumber']} päivältä {p['earliestPriorityDate']}.",
            "IP-status",
            f"EPO-tietue: {p['familyLabel']}.",
            sources("EPO-REGISTER-MAIN", "EPO-REGISTER-FAMILY"),
            p["earliestPriorityDate"],
            "Suora rekisterihavainto.",
            "Ei oletuksia.",
            "Vahvistettu",
            "Ei olennaista puutetta prioriteettiväitteen osalta.",
        ),
        row(
            "EPO:n keskitetty väite- ja valitusmenettely päättyi patentin pysyttämiseen muutettuna.",
            "IP-status",
            epo_proceeding["detailFi"],
            sources(*epo_proceeding["sourceIds"]),
            epo_proceeding["eventDate"],
            "Virallisten EPO-tapahtumien ja B2-julkaisun ristiintarkastus.",
            "Kansallinen voimassaolo käsitellään erikseen.",
            "Vahvistettu",
            "Kunkin kansallisen oikeuden nykytila on vahvistettava kansallisesta rekisteristä.",
        ),
        row(
            f"Muutetussa {p['epPublication']}-julkaisussa on {fi_cardinal(b2_claim_count)} vaatimusta.",
            "Patentoitu ratkaisu",
            p["claimScopeSummaryFi"],
            sources("EPO-B2-SPECIFICATION"),
            p["epCentralStatusDate"],
            "B2-julkaisun vaatimusten lukumäärä.",
            "Yleiskielinen tiivistelmä ei ole claim construction.",
            "Vahvistettu",
            "Maakohtainen asiantuntijan vaatimusvertailu tarvitaan kaupallisiin tai oikeudellisiin johtopäätöksiin.",
        ),
        row(
            f"Patenttiperheessä on {ctx['patent_history']['summary']['familyRecordCount']} julkaisupohjaista tietuetta.",
            "Maantieteellinen kattavuus",
            f"Julkisen patenttihistoriadatan familyRecordCount = {ctx['patent_history']['summary']['familyRecordCount']}.",
            sources("EPO-REGISTER-FAMILY"),
            as_of,
            "EPO-perhejulkaisujen lukumäärä; ei voimassa olevien oikeuksien lukumäärä.",
            "Julkaisutietue rinnastetaan vain perhereitiksi.",
            "Tuettu",
            "Vahvista omistus, maksut, käytettävät vaatimukset ja voimassaolo jokaisessa maassa.",
        ),
        row(
            f"Saksan virallinen tapahtuma: "
            f"{legal_event_title(de_nullity, facts['de_nullity_process_state'])} "
            f"({fi_date(de_nullity['eventDate'])}).",
            "Oikeudellinen näyttö",
            de_nullity["detailFi"],
            sources(*de_nullity["sourceIds"]),
            de_nullity["eventDate"],
            "Virallisen ratkaisun ja valitusviitteen lukeminen.",
            legal_assumption(de_nullity),
            "Vahvistettu",
            legal_gap(de_nullity, facts["de_nullity_process_state"]),
        ),
        row(
            f"Saksan virallinen tapahtuma koski vaatimuksia "
            f"{fi_join([str(value) for value in infringement_claims])}: "
            f"{legal_event_title(de_infringement, facts['de_infringement_process_state'])}.",
            "Oikeudellinen näyttö",
            de_infringement["detailFi"],
            sources(*de_infringement["sourceIds"]),
            de_infringement["eventDate"],
            "Virallisen tuomion rajattu kuvaus.",
            legal_assumption(de_infringement),
            "Vahvistettu",
            legal_gap(de_infringement, facts["de_infringement_process_state"]),
        ),
        row(
            f"Kiinan tunnistettu menettely: {cn_review['titleFi']}.",
            "Oikeudellinen näyttö",
            cn_review["detailFi"],
            sources(*cn_review["sourceIds"]),
            cn_review["eventDate"],
            "Sekundäärisen docket-tiedon luokittelu virallisen prosessiohjeen avulla.",
            "Nimivastaavuus on vahva mutta virallista päätöstä ei saatu.",
            "Tuettu",
            "Hanki CNIPA:n virallinen päätös ja sen perustelut.",
        ),
        row(
            f"Julkinen atlas sisältää muuttumattoman {ctx['atlas']['summary']['countryCount']} maan tutkimusuniversumin.",
            "Markkinan rajaus",
            f"countryCount = {ctx['atlas']['summary']['countryCount']}; universe = {ctx['atlas']['summary']['universe']}.",
            "https://www.un.org/en/about-us/member-states ; https://www.un.org/en/about-us/non-member-states",
            as_of,
            f"{un_member_count} YK:n jäsenvaltiota + Pyhä istuin + Palestiinan valtio.",
            "Universumi on tutkimusrunko, ei todistettu markkinapeitto.",
            "Vahvistettu",
            "Ei puutetta universumin määrittelyssä; evidenssipeitto on erillinen asia.",
        ),
        row(
            f"Hyväksyttyjä vuosittaisia virallisia määrähavaintoja on {fi_elative_cardinal(facts['official_country_count'])} maasta.",
            "Markkinan koko",
            f"Maat: {fi_join(official_country_names_fi)}.",
            "site/data/market-values.json (julkisen sivuston koneellisesti luettava lähdetiedosto)",
            as_of,
            "Uniikit maat virallisiksi luokitelluista vuosihavainnoista.",
            "Mittarit eivät ole keskenään samanlaisia.",
            "Vahvistettu",
            "Tarvitaan vertailukelpoiset laite- ja nestemyyntisarjat muista maista.",
        ),
        row(
            f"Hyväksyttyjä virallisia koko vuoden kansallisia kuluttajavähittäisarvon luovuttajamarkkinoita on {fi_cardinal(facts['retail_donors'])}.",
            "Markkinan koko",
            f"comparableFullYearMarketValueDonors = {facts['retail_donors']}.",
            "site/data/market-values.json (modelReadiness)",
            as_of,
            "Yhteensopivien koko vuoden kuluttajavähittäisarvojen hyväksyntäkriteeri.",
            "Toimitus-, vero- ja määräluvut eivät ole vähittäisarvoja.",
            "Vahvistettu",
            "Hanki vähintään kolme yhteensopivaa luovuttajamarkkinaa sekä alue- ja sääntelytyyppien peitto.",
        ),
        row(
            "Kanadan vuoden "
            f"{observation_year(ca_value)} "
            "valmistaja- ja maahantuojatoimitusten arvo oli "
            f"{format_local_number(ca_value['value'] / 1_000_000_000, 11, grouped=False)} mrd "
            f"{ca_value['currency']}.",
            "Markkinan koko",
            ca_value["limitationFi"],
            sources(*ca_value["sourceIds"]),
            str(observation_year(ca_value)),
            "Health Canadan neljän raportoidun tuoteryhmän toimitusarvon summa.",
            "Ei kuluttajavähittäismyyntiä.",
            "Vahvistettu",
            "Tarvitaan vähittäismarginaali-, kanava- ja varastomuutostäsmäytys.",
        ),
        row(
            "Kanadan vuoden "
            f"{observation_year(ca_litres)} raportoidut toimitukset sisälsivät "
            f"{format_local_number(ca_litres['value'])} litraa nestettä.",
            "Markkinan koko",
            ca_litres["limitationFi"],
            sources(*ca_litres["sourceIds"]),
            str(observation_year(ca_litres)),
            "Raportoitujen nestettä sisältävien tuoteryhmien summa.",
            "Fyysinen määrä, ei markkina-arvo.",
            "Vahvistettu",
            "Tarvitaan tuotemixin, keskihinnan ja vähittäiskanavan tiedot.",
        ),
        row(
            f"Saksan verotettu nestemäärä kasvoi {million_litres(de23['value'])}:sta {million_litres(de25['value'])}:aan "
            f"vuosina {observation_year(de23)}–{observation_year(de25)}.",
            "Markkinan koko",
            de_volume_status_proof,
            sources(*unique_strings(source_id for item in de_volume_observations for source_id in item["sourceIds"])),
            f"{observation_year(de23)}–{observation_year(de25)}",
            "Destatis-taulukon nettomäärät, ei ekstrapolointia.",
            "Mittaa verotettua nestettä, ei laitteita tai laitonta myyntiä.",
            "Vahvistettu",
            f"Vuoden {observation_year(de25)} lopullinen luku ja vähittäismyyntitiedot puuttuvat.",
        ),
        row(
            "Saksan tupakankorvikkeiden valmisteverotulot olivat "
            f"{format_local_number(de_latest_excise['value'] / 1_000_000)} milj. euroa "
            f"vuonna {observation_year(de_latest_excise)}"
            f"{' (alustava)' if de_latest_excise.get('finality') == 'provisional' else ''}.",
            "Markkinan koko",
            de_latest_excise["limitationFi"],
            sources(*de_latest_excise["sourceIds"]),
            str(observation_year(de_latest_excise)),
            "Virallinen kassaperusteinen verohavainto.",
            "Verotulo ei ole myyntitulo tai vähittäismarkkina-arvo.",
            "Vahvistettu",
            f"Tarvitaan lopullinen {observation_year(de_latest_excise)} tilasto ja tuoteryhmäkohtainen veropohja.",
        ),
        row(
            "Suomessa verotettiin "
            f"{format_local_number(fi_volume['value'], 3)} litraa nikotiininestettä "
            f"vuonna {observation_year(fi_volume)}.",
            "Markkinan koko",
            fi_volume["limitationFi"],
            sources(*fi_volume["sourceIds"]),
            str(observation_year(fi_volume)),
            "Verohallinnon PXWeb-havainto.",
            "Verotettu nikotiinineste ei kata koko sähkötupakkamarkkinaa.",
            "Vahvistettu",
            "Laitteet, nikotiinittomat tuotteet, veroton ja laiton kauppa puuttuvat.",
        ),
        row(
            "Suomen nikotiininesteiden valmisteverotulot olivat "
            f"{format_local_number(fi_excise['value'])} euroa "
            f"vuonna {observation_year(fi_excise)}.",
            "Markkinan koko",
            fi_excise["limitationFi"],
            sources(*fi_excise["sourceIds"]),
            str(observation_year(fi_excise)),
            "Verohallinnon PXWeb-havainto.",
            "Verotulo ei ole kuluttajamyynti.",
            "Vahvistettu",
            "Tarvitaan vähittäismyynti- ja hintadata.",
        ),
        row(
            "Puolassa raportoitu sähkötupakkanesteiden määrä oli "
            f"{format_local_number(pl_volume['value'])} litraa "
            f"vuonna {observation_year(pl_volume)}.",
            "Markkinan koko",
            pl_volume["limitationFi"],
            sources(*pl_volume["sourceIds"]),
            str(observation_year(pl_volume)),
            "Parlamentaariseen vastaukseen raportoitu määrä.",
            "Mittarin kattavuus on luettava alkuperäislähteen rajauksin.",
            "Vahvistettu",
            "Tarvitaan uudempi vuosisarja ja vähittäisarvo.",
        ),
        row(
            f"Puolan vuoden {observation_year(pl_excise)} ilmoitettu e-nestevalmisteveron määrä oli "
            f"{format_local_number(pl_excise['value'] / 1_000_000, 1)} milj. "
            f"{pl_excise['currency']}.",
            "Markkinan koko",
            pl_excise["limitationFi"],
            sources(*pl_excise["sourceIds"]),
            str(observation_year(pl_excise)),
            "Parlamentaariseen vastaukseen raportoitu veromäärä.",
            "Veromäärä ei ole vähittäismyynti.",
            "Vahvistettu",
            "Tarvitaan toteuman lopullisuus ja veropohjan täsmäytys.",
        ),
        row(
            "Ruotsissa verotettiin "
            f"{format_local_number(se_volume['value'])} litraa nikotiininestettä "
            f"vuonna {observation_year(se_volume)}.",
            "Markkinan koko",
            se_volume["limitationFi"],
            sources(*se_volume["sourceIds"]),
            str(observation_year(se_volume)),
            "Hallituksen laskentaperusteessa ilmoitettu määrä.",
            "Verotettu nikotiinineste ei kata koko sähkötupakkamarkkinaa.",
            "Vahvistettu",
            "Tarvitaan toteutunut vähittäismyynti ja laitedata.",
        ),
        row(
            f"{fi_cardinal(len(global_observations)).capitalize()} kaupallista vuoden {global_year} "
            "globaaliarviota muodostavat "
            f"{format_compact_local_number(min(global_values) / 1_000_000_000, minimum_decimals=1, grouped=False)}–"
            f"{format_compact_local_number(max(global_values) / 1_000_000_000, minimum_decimals=1, grouped=False)} mrd "
            f"{global_currency} haarukan.",
            "Markkinan koko",
            "; ".join(
                f"{global_publisher(item)} "
                f"{format_compact_local_number(item['value'] / 1_000_000_000, minimum_decimals=1, grouped=False)}"
                for item in global_observations
            )
            + f" mrd {global_currency}.",
            sources(*global_source_ids),
            str(global_year),
            f"Minimi {min(global_values):.0f}, maksimi {max(global_values):.0f}; arvioita verrataan, ei summata.",
            "Tuoterajaukset ja menetelmät eivät välttämättä ole yhteismitallisia.",
            "Tuettu",
            "Hanki raporttien täydet metodit ja harmonisoi tuoterajaus.",
        ),
        row(
            f"Saksan vuoden {model['year']} verotetun nesteen vähittäismyyntivastaavuuden mallihaitari on "
            f"{format_local_number(model['low'] / 1_000_000, 2)}–{format_local_number(model['high'] / 1_000_000, 2)} milj. euroa.",
            "Taloudellinen malli",
            model["limitationFi"],
            sources(*facts["model_source_ids"]),
            model_date_label,
            model["formula"],
            model_assumptions,
            "Oletus",
            model_gap,
        ),
        row(
            "Atlas ei ole vielä lainanantajavalmis markkina-arvio.",
            "Rahoitusteesi",
            ctx["atlas"]["readiness"]["status"] + ": " + " ".join(ctx["atlas"]["readiness"]["blockers"]),
            "site/data/atlas.json (readiness)",
            as_of,
            "Julkisen datasetin oma readiness-luokitus.",
            "Markkinadata ja IP-arvo pidetään erillään.",
            "Vahvistettu",
            "Täytä blocker-lista ennen vakuusarvo- tai yritysarvopäätelmää.",
        ),
        row(
            "Julkinen aineisto ei sisällä riippumatonta IVS-arvonmääritystä.",
            "Taloudellinen malli",
            "Readiness-blocker ilmoittaa riippumattoman arvonmäärityksen puuttuvan.",
            "site/data/atlas.json (readiness.blockers)",
            as_of,
            "Aineiston puute.",
            "Ei oletuksia.",
            "Puuttuu",
            "Tilaa riippumaton IP- ja tarvittaessa yritysarvonmääritys skenaarioineen.",
        ),
        row(
            "Julkinen aineisto ei osoita toteutunutta lisenssi- tai vahingonkorvauskassavirtaa.",
            "Kaupallistamismalli",
            "Readiness-blocker ilmoittaa realisoituneen kassavirran puuttuvan.",
            "site/data/atlas.json (readiness.blockers)",
            as_of,
            "Aineiston puute.",
            "Ei oletuksia.",
            "Puuttuu",
            "Lisää allekirjoitetut sopimukset, laskut, maksutositteet ja saamisten täsmäytys.",
        ),
        row(
            "Julkinen aineisto ei vahvista testituloksia tai riippumatonta teknistä validointia.",
            "Tekninen erottautuminen",
            "Patenttiselitys kuvaa ratkaisun; erillistä testirekisteriä ei ole julkisessa datassa.",
            "site/data/patent-history.json",
            as_of,
            "Aineiston puute suhteessa tekniseen suorituskykyväitteeseen.",
            "Patenttijulkaisu ei ole tuotetestiraportti.",
            "Puuttuu",
            "Hanki testiprotokolla, riippumaton laboratorio, tulokset ja raakadata.",
        ),
        row(
            "Julkinen aineisto ei vahvista piloteja tai nimettyä asiakasnäyttöä.",
            "Validointi",
            "Julkisessa evidenssirekisterissä ei ole pilotti- tai asiakasvalidointia.",
            "site/data/evidence.csv",
            as_of,
            "Aineiston puute.",
            "Ei päätellä asiakkaiden olemassaoloa tai puuttumista yhtiötasolla.",
            "Puuttuu",
            "Lisää allekirjoitetut pilotit, toimitusnäyttö, asiakasreferenssit ja käyttödata.",
        ),
        row(
            "Julkinen aineisto ei vahvista nykyistä maakohtaista omistusta, vuosimaksuja ja rasitteettomuutta koko perheessä.",
            "IP-status",
            "Perhejulkaisut eivät yksin osoita kansallista post-grant-tilaa.",
            sources("EPO-REGISTER-FAMILY", "EPO-REGISTER-LEGAL", "EPO-NATIONAL-VALIDATION"),
            as_of,
            "Diligence-rajaus.",
            f"{fi_cardinal(facts['national_register_check_count']).capitalize()} kansallista rekisterijurisdiktiota ei kata koko perhettä.",
            "Puuttuu",
            "Laadi asiamiehen allekirjoittama maamatriisi ja liitä tuoreet rekisteriotteet sekä maksukuitit.",
        ),
        row(
            "Julkinen aineisto ei sisällä maakohtaisia tuote–vaatimusvertailuja.",
            "Tekninen erottautuminen",
            "Patenttihistorian guardrail edellyttää counsel-reviewed claim chartia.",
            sources("EPO-B2-SPECIFICATION", *de_infringement["sourceIds"]),
            as_of,
            "Vaatimuspiirteiden tuotekohtainen kartoitus puuttuu.",
            "Samankaltaisuus ei ole loukkaus.",
            "Puuttuu",
            "Hanki tuotteet, säilytä hallussapitoketju, tee teardown ja asiamiehen tarkastama claim chart.",
        ),
        row(
            "Markkinan kokonaisarvo ei ole automaattisesti rojaltipohja.",
            "Taloudellinen malli",
            "Patenttihistorian kaupallistamisrajaus erottaa markkinan, kohdistettavan myynnin ja kassavirran.",
            sources("WIPO-IP-VALUATION"),
            as_of,
            "Arvonmäärityksen perusrajaus.",
            "Rojaltipohja vaatii patentin maantieteellisen, ajallisen ja tuotekohtaisen osuvuuden.",
            "Vahvistettu",
            "Tarvitaan addressable/in-scope/claim-mapped sales -silta.",
        ),
        row(
            "Patenttiarvo, yritysarvo, osakearvo ja vakuusarvo ovat eri mittareita.",
            "Rahoitusteesi",
            "WIPO:n IP-arvonmääritys- ja rahoituskehys sekä aineiston oma guardrail.",
            sources("WIPO-IP-VALUATION", "WIPO-IP-FINANCE"),
            as_of,
            "Arvokäsitteiden erottelu.",
            "Ei oletuksia.",
            "Vahvistettu",
            "Rahoitusrakenne ja vakuusarvo vaativat erillisen oikeudellisen ja taloudellisen analyysin.",
        ),
        row(
            "Asiantuntijavetoinen rajattu lisensointi- tai sovintopilotti on mahdollinen kaupallistamishypoteesi.",
            "Kaupallistamismalli",
            "Patenttihistorian vaiheistus suosittelee kovia portteja, claim chartia ja auditoitavaa yhteydenottoa.",
            sources("WIPO-ASSIGNMENT-LICENSING", "WIPO-DISPUTE-RESOLUTION", "EPO-IPSCORE"),
            as_of,
            "Hypoteesi, ei toteutunut kaupallinen näyttö.",
            "Kohdemaat ja vastapuolet valitaan vasta oikeus- ja evidenssiporttien jälkeen.",
            "Oletus",
            "Hyväksytä pilottimalli johdolla ja patenttiasiantuntijalla; dokumentoi kustannus, vasteet ja termit.",
        ),
        row(
            "Patenttivakuudellinen rahoitus on arvioitava vain vahvistettujen oikeuksien ja kassavirtojen pohjalta.",
            "Rahoitusteesi",
            f"WIPO IP Finance ja patenttihistorian vaihe {finance_phase}.",
            sources("WIPO-IP-FINANCE", "WIPO-IP-VALUATION"),
            as_of,
            "Rahoitusperiaate.",
            "Ei näyttöä olemassa olevasta vakuusarvosta.",
            "Tuettu",
            "Tarvitaan riippumaton arvonmääritys, oikeuksien due diligence, toteutuneet tai sopimuspohjaiset kassavirrat ja downside-analyysi.",
        ),
        row(
            "Australian oikeuden seuraava uusiminen on vahvistettava välittömästi.",
            "Riskit",
            facts["au_alert"]["detailFi"],
            sources(*facts["au_alert"]["sourceIds"]),
            facts["au_alert"]["targetDate"],
            "Virallisen API-tiedon diligence-hälytys.",
            "API:n in-force-through-päivä ei ole tässä lakisääteinen eräpäiväväite.",
            "Vahvistettu",
            facts["au_alert"]["actionFi"],
        ),
        row(
            f"Suomen {fi_fee_year}. vuosimaksun eräpäiväksi on rekisteröity "
            f"{fi_date(facts['fi_alert']['targetDate'])}.",
            "Riskit",
            facts["fi_alert"]["detailFi"],
            sources(*facts["fi_alert"]["sourceIds"]),
            facts["fi_alert"]["targetDate"],
            "PRH:n rekisterihavainto.",
            "Tavoite on maksaa normaalina eräpäivänä.",
            "Vahvistettu",
            facts["fi_alert"]["actionFi"],
        ),
        row(
            "Yhdysvaltain "
            f"{format_compact_local_number(us_maintenance_year, minimum_decimals=1, maximum_decimals=1, grouped=False)} "
            "vuoden ylläpitomaksun tila on vahvistettava.",
            "Riskit",
            facts["us_alert"]["detailFi"],
            sources(*facts["us_alert"]["sourceIds"]),
            facts["us_alert"]["targetDate"],
            "Virallisen maksuikkunan tarkistus.",
            "Julkaisu ei väitä, että maksu on tai ei ole suoritettu.",
            "Vahvistettu",
            facts["us_alert"]["actionFi"],
        ),
        row(
            "Julkinen aineisto ei sisällä auditoituja historiallisia tilinpäätöksiä tai kassavirtaennustetta.",
            "Taloudellinen malli",
            "Markkina- ja patenttiaineisto ei ole yhtiön talousaineisto.",
            "Julkisen paketin rajaus",
            as_of,
            "Aineiston puute.",
            "Yhtiön taloudellisesta tilanteesta ei tehdä päätelmää.",
            "Puuttuu",
            "Lisää 3–5 vuoden tilinpäätökset, tuore pääkirja, budjetti, kassavirta, velat ja verovelkatodistus.",
        ),
        row(
            "Julkinen aineisto ei yksilöi rahoitustarvetta, käyttötarkoitusta tai takaisinmaksulähdettä.",
            "Rahoitusteesi",
            "Rahoitusparametreja ei ole markkina- tai patenttidatassa.",
            "Julkisen paketin rajaus",
            as_of,
            "Aineiston puute.",
            "Ei oleteta lainamäärää, maturiteettia tai vakuusrakennetta.",
            "Puuttuu",
            "Määritä lainamäärä, käyttötarkoitus, maturiteetti, lyhennysprofiili, takaisinmaksu ja kovenantit.",
        ),
        row(
            "Julkinen aineisto ei osoita osakekantaa, cap tablea tai osakkeiden rasitteita.",
            "Rahoitusteesi",
            "Julkinen paketti on rajattu markkina- ja patenttievidenssiin.",
            "Julkisen paketin rajaus",
            as_of,
            "Aineiston puute.",
            "Ei oleteta omistusosuuksia tai vakuuskelpoisuutta.",
            "Puuttuu",
            "Toimita varmennettu osakasluettelo, yhtiöjärjestys, päätökset, panttaus- ja rasiteselvitys suljettuun datahuoneeseen.",
        ),
        row(
            "Julkinen aineisto ei sisällä vertailukelpoista kilpailija- ja vaihtoehtoanalyysiä.",
            "Kilpailu",
            "Patenttiselitys ja oikeustapaukset eivät yksin muodosta markkinakilpailukarttaa.",
            "site/data/patent-history.json",
            as_of,
            "Aineiston puute.",
            "Kilpailu on erotettava patenttien päällekkäisyydestä ja ei-patentoiduista vaihtoehdoista.",
            "Puuttuu",
            "Laadi tuotteet, valmistajat, vaihtoehtoiset teknologiat, patenttiperheet, hinnat ja claim-map-status kattava vertailu.",
        ),
        row(
            "Asiakassegmentit voidaan alustavasti jäsentää valmistajiin, teknologiatoimittajiin ja IP-rahoittajiin.",
            "Asiakkaat",
            "Patenttihistorian kaupallistamisreitit sisältävät lisensoinnin, luovutuksen ja IP-rahoituksen.",
            sources("WIPO-ASSIGNMENT-LICENSING", "WIPO-IP-FINANCE"),
            as_of,
            "Segmentointihypoteesi kaupallistamisreiteistä.",
            "Ei todennettua ostoaikomusta tai asiakaskantaa.",
            "Oletus",
            "Validoi 10–15 strukturoitua ostaja- ja rahoittajahaastattelua sekä dokumentoi päätöskriteerit.",
        ),
        row(
            "Saksan ratkaisut voivat toimia evidenssiankkurina, mutta eivät maailmanlaajuisena loukkausratkaisuna.",
            "Rahoitusteesi",
            "Ratkaisut ovat kansallisia, tuote-, osapuoli- ja ajanjaksokohtaisia.",
            sources(*unique_strings(de_nullity["sourceIds"] + de_infringement["sourceIds"])),
            as_of,
            "Oikeudellisen vaikutuksen alueellinen rajaus.",
            "Muiden maiden tulokset riippuvat paikallisista oikeuksista ja tuotteista.",
            "Vahvistettu",
            "Laadi maakohtaiset oikeus- ja claim chart -paketit ennen yhteydenottoja.",
        ),
        row(
            f"Aineiston A-luokan maita on {fi_cardinal(facts['grade_counts']['A'])}, mutta "
            f"{facts['grade_counts']['D']} maata on D-luokassa.",
            "Markkinan koko",
            f"Luokat: A {facts['grade_counts']['A']}, B {facts['grade_counts']['B']}, C {facts['grade_counts']['C']}, D {facts['grade_counts']['D']}.",
            "site/data/atlas.json (summary.gradeCounts)",
            as_of,
            "Maakohtaisten evidenssiluokkien laskenta.",
            "Luokka mittaa evidenssikypsyyttä, ei markkinan houkuttelevuutta.",
            "Vahvistettu",
            "Priorisoi D-maiden tietopyynnöt talouskoon, patenttistatuksen ja sääntelyn mukaan.",
        ),
        row(
            "Maailmanlaajuista atlasestimaattia ei ole hyväksytty julkaistavaksi.",
            "Markkinan koko",
            f"{fi_cardinal(facts['retail_donors']).capitalize()} vertailukelpoista vähittäisarvoluovuttajaa alittaa "
            f"{fi_cardinal(facts['minimum_required_donors'])}n luovuttajan minimikynnyksen.",
            "site/data/market-values.json (modelReadiness)",
            as_of,
            "Hard gate: vähintään "
            f"{fi_cardinal(facts['minimum_required_donors'])} yhteensopivaa luovuttajaa + alueellinen validointi.",
            "Kaupallisia globaaliarvioita käytetään vain sanity check -haarukkana.",
            "Vahvistettu",
            "Älä esitä yhtä maailmanlukua ennen metodin porttien täyttymistä.",
        ),
        row(
            "Korkean luottamuksen rahoituspaketti vaatii suljetun datahuoneen julkisen paketin rinnalle.",
            "Seuraavat vaiheet",
            "Julkinen paketti osoittaa lähteet ja puutteet mutta ei sisällä luottamuksellista talous-, sopimus- tai omistusaineistoa.",
            "Julkisen paketin rajaus",
            as_of,
            "Diligence-arkkitehtuurin suositus.",
            "Pääsy rajataan ja lokitetaan.",
            "Tuettu",
            "Perusta indeksoitu datahuone: Corporate, IP, Legal, Commercial, Finance, Valuation, Security.",
        ),
    ]
    public_text_scan(value for item in rows for value in item.values())
    return rows


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, radius: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_footer(slide, ctx: dict[str, Any], page: int, source_ids: str = "") -> None:
    add_rect(slide, 0.38, 7.13, 12.57, 0.015, LINE)
    footer = f"Julkinen riippumaton evidenssikooste · {ctx['release']['version']} · {ctx['as_of']}"
    if source_ids:
        footer += f" · Lähteet: {source_ids}"
    add_text(slide, footer, 0.42, 7.17, 11.85, 0.19, size=8.5, color=MUTED)
    add_text(slide, str(page), 12.45, 7.15, 0.45, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title: str, section: str, ctx: dict[str, Any], page: int, sources: str = "") -> None:
    add_text(slide, section.upper(), 0.55, 0.32, 3.8, 0.28, size=10, color=TEAL, bold=True)
    add_text(slide, title, 0.55, 0.68, 12.15, 0.6, size=34, color=NAVY, bold=True)
    add_rect(slide, 0.55, 1.34, 1.05, 0.055, TEAL)
    add_footer(slide, ctx, page, sources)


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, size: float = 18, color: str = INK) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for idx, text in enumerate(bullets):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = text
        para.level = 0
        para.font.name = "Aptos"
        para.font.size = Pt(size)
        para.font.color.rgb = rgb(color)
        para.space_after = Pt(12)
        para.line_spacing = 1.05
        para.text = "•  " + para.text


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, color: str = TEAL) -> None:
    add_rect(slide, x, y, w, 1.5, LIGHT, LINE, True)
    add_text(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.58, size=27, color=color, bold=True)
    add_text(slide, label, x + 0.18, y + 0.82, w - 0.36, 0.48, size=14, color=MUTED)


def add_callout(slide, title: str, body: str, x: float, y: float, w: float, h: float, fill: str = PALE_TEAL) -> None:
    add_rect(slide, x, y, w, h, fill, fill, True)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.3, size=17, color=NAVY, bold=True)
    add_text(slide, body, x + 0.18, y + 0.55, w - 0.36, h - 0.68, size=15, color=INK)


def add_table(slide, headers: list[str], rows: list[list[str]], x: float, y: float, w: float, h: float, widths: list[float] | None = None, font_size: float = 13.5) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        total = sum(widths)
        for idx, value in enumerate(widths):
            table.columns[idx].width = Inches(w * value / total)
    header_height = min(0.48, h / (len(rows) + 1) * 1.2)
    table.rows[0].height = Inches(header_height)
    body_height = (h - header_height) / max(1, len(rows))
    for row_index in range(1, len(table.rows)):
        table.rows[row_index].height = Inches(body_height)
    for col, value in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(NAVY)
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        for para in cell.text_frame.paragraphs:
            para.font.name = "Aptos"
            para.font.size = Pt(font_size)
            para.font.bold = True
            para.font.color.rgb = rgb(WHITE)
    for ridx, row_values in enumerate(rows, start=1):
        for cidx, value in enumerate(row_values):
            cell = table.cell(ridx, cidx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(WHITE if ridx % 2 else LIGHT)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            for para in cell.text_frame.paragraphs:
                para.font.name = "Aptos"
                para.font.size = Pt(font_size)
                para.font.color.rgb = rgb(INK)
                para.alignment = PP_ALIGN.LEFT


def new_presentation(ctx: dict[str, Any], title: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = title
    prs.core_properties.subject = "Julkinen pankki- ja rahoitusarvioinnin evidenssipaketti"
    prs.core_properties.author = "Pixan Global Market Evidence Atlas"
    prs.core_properties.keywords = "Pixan, patentti, markkinaevidenssi, due diligence"
    stamp = parse_iso_date(ctx["as_of"])
    prs.core_properties.created = stamp
    prs.core_properties.modified = stamp
    return prs


def cover_slide(prs: Presentation, ctx: dict[str, Any], subtitle: str, slide_count: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY)
    add_rect(slide, 0.58, 0.62, 1.2, 0.07, TEAL)
    add_text(slide, "PIXAN · JULKINEN EVIDENSSIPAKETTI", 0.58, 1.0, 7.4, 0.35, size=13, color=TEAL, bold=True)
    # Keep the cover title and subtitle in distinct vertical bands.  The title
    # can wrap to three lines in LibreOffice/PowerPoint depending on font
    # metrics, so reserve enough height instead of relying on a two-line fit.
    add_text(slide, "Patentista pankkikelpoiseksi\ntodistelupaketiksi", 0.58, 1.52, 8.05, 2.25, size=40, color=WHITE, bold=True)
    add_text(slide, subtitle, 0.62, 4.05, 7.7, 0.48, size=18, color="C6D8E2")
    add_rect(slide, 9.05, 0.0, 4.28, 7.5, BLUE)
    add_text(slide, "Aineiston tila", 9.55, 1.25, 3.2, 0.4, size=19, color=WHITE, bold=True)
    add_text(slide, ctx["release"]["version"], 9.55, 1.86, 3.15, 0.55, size=29, color=GOLD, bold=True)
    add_text(slide, f"Päivitetty {ctx['as_of']}\n{slide_count} diaa · suomeksi", 9.55, 2.58, 3.05, 0.8, size=17, color=WHITE)
    add_text(slide, "Lähtökohta", 9.55, 4.18, 3.0, 0.35, size=13, color="B8DFE2", bold=True)
    add_text(slide, "Vahvistettu näyttö erotetaan tuetusta tiedosta, oletuksista ja puuttuvasta näytöstä.", 9.55, 4.65, 3.05, 1.2, size=18, color=WHITE)
    add_text(slide, "Ei Pixan Oy:n virallinen kanta · ei arvo-, sijoitus-, laina- tai oikeudellinen lausunto", 0.62, 6.84, 7.7, 0.34, size=10, color="9DB3C0")


def slide_claim(prs: Presentation, ctx: dict[str, Any], page: int, title: str, section: str, claim: str, bullets: list[str], side_title: str, side_body: str, sources: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title, section, ctx, page, sources)
    # Finnish compound words and evidence-qualified claims vary materially in
    # length.  Reserve claim height before placing bullets so the two text
    # regions can never collide after Office/LibreOffice line wrapping.
    if len(claim) <= 115:
        claim_size, claim_height = 23, 1.12
    elif len(claim) <= 175:
        claim_size, claim_height = 20, 1.42
    else:
        claim_size, claim_height = 17.5, 1.78
    bullet_y = 1.78 + claim_height + 0.18
    add_text(slide, claim, 0.6, 1.7, 7.4, claim_height, size=claim_size, color=BLUE, bold=True)
    add_bullets(slide, bullets, 0.64, bullet_y, 7.25, 6.45 - bullet_y, size=16.5)
    add_callout(slide, side_title, side_body, 8.45, 1.7, 4.25, 4.75, PALE_TEAL)


def slide_metrics(prs: Presentation, ctx: dict[str, Any], page: int, title: str, section: str, metrics: list[tuple[str, str]], takeaway: str, note: str, sources: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title, section, ctx, page, sources)
    count = len(metrics)
    card_w = (12.15 - 0.28 * (count - 1)) / count
    for idx, (value, label) in enumerate(metrics):
        add_metric(slide, value, label, 0.58 + idx * (card_w + 0.28), 1.88, card_w, [TEAL, BLUE, GOLD, GREEN][idx % 4])
    add_text(slide, takeaway, 0.62, 3.86, 11.95, 0.86, size=25, color=NAVY, bold=True)
    add_callout(slide, "Tulkinta", note, 0.62, 4.96, 11.95, 1.25, PALE)


def slide_table(prs: Presentation, ctx: dict[str, Any], page: int, title: str, section: str, headers: list[str], rows: list[list[str]], takeaway: str, sources: str, widths: list[float] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title, section, ctx, page, sources)
    add_table(slide, headers, rows, 0.58, 1.68, 12.15, 4.35, widths, font_size=13.2)
    add_text(slide, takeaway, 0.66, 6.19, 11.9, 0.55, size=17, color=BLUE, bold=True)


def slide_process(prs: Presentation, ctx: dict[str, Any], page: int, title: str, section: str, steps: list[tuple[str, str]], takeaway: str, sources: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title, section, ctx, page, sources)
    top = 1.78
    row_h = 1.05
    for idx, (head, body) in enumerate(steps):
        y = top + idx * 1.18
        add_rect(slide, 0.72, y, 0.62, 0.62, TEAL if idx < 3 else BLUE, None, True)
        add_text(slide, str(idx + 1), 0.72, y + 0.08, 0.62, 0.36, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, head, 1.62, y - 0.02, 3.0, 0.34, size=18, color=NAVY, bold=True)
        add_text(slide, body, 4.5, y - 0.02, 7.75, row_h, size=16, color=INK)
    add_text(slide, takeaway, 0.72, 6.44, 11.65, 0.44, size=16, color=BLUE, bold=True)


def closing_slide(prs: Presentation, ctx: dict[str, Any], page: int, title: str, bullets: list[str], sources: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY)
    add_text(slide, "PÄÄTÖSKEHYS", 0.62, 0.55, 3.0, 0.3, size=12, color=TEAL, bold=True)
    add_text(slide, title, 0.62, 1.1, 11.9, 1.0, size=39, color=WHITE, bold=True)
    add_bullets(slide, bullets, 0.68, 2.5, 10.95, 3.55, size=19, color=WHITE)
    add_text(slide, f"Versio {ctx['release']['version']} · {ctx['as_of']} · {page} / {len(prs.slides)}", 0.65, 6.86, 11.2, 0.25, size=10, color="9DB3C0")


def common_values(ctx: dict[str, Any]) -> dict[str, Any]:
    facts = canonical_facts(ctx)
    patent = facts["patent"]
    ca_value = facts["ca_value_observation"]
    ca_units = facts["ca_units_observation"]
    ca_litres = facts["ca_litres_observation"]
    de_volumes = facts["de_volume_observations"]
    de_excise_by_year = facts["de_excise_by_year"]
    model = facts["model"]
    model_prices = facts["model_prices"]
    global_observations = facts["global_observations"]
    global_values = facts["global_values"]
    global_currency = facts["global_currency"]
    market_sources = facts["market_source_map"]

    if {item.get("currency") for item in model_prices.values()} != {model["currency"]}:
        raise ValueError("Germany model price currencies differ from the model currency")
    if {item.get("unit") for item in model_prices.values()} != {f"{model['currency']}_per_ml"}:
        raise ValueError("Germany model price inputs must share the model currency per ml")

    de_rows = []
    for volume in de_volumes:
        year = observation_year(volume)
        excise = de_excise_by_year[year]
        de_rows.append(
            [
                str(year),
                format_millions(volume["value"], fi_unit(volume["unit"])),
                format_millions(excise["value"], excise["currency"]),
                fi_finality(volume.get("finality")),
            ]
        )

    global_metrics = []
    for item in global_observations:
        publisher = concise_publisher(market_sources[item["sourceIds"][0]]["publisher"])
        global_metrics.append(
            (format_billions(item["value"], item["currency"]), f"{publisher} {observation_year(item)}")
        )

    nullity_reference = facts["de_nullity"]["reference"].split(" · ", 1)[0]
    german_references = "; ".join(
        item["reference"].split(" · ", 1)[0]
        for item in facts["german_official_proceedings"]
    )
    confidence_fi = {"low": "matala", "medium": "keskitaso", "high": "korkea"}.get(
        model["confidence"], model["confidence"]
    )
    grade_labels = facts["grade_labels"]

    return {
        **facts,
        "ca_country": facts["country_name"](ca_value["countryIso2"]),
        "ca_year": observation_year(ca_value),
        "ca_value": format_billions(ca_value["value"], ca_value["currency"]),
        "ca_units": (
            f"{format_local_number(ca_units['value'] / 1_000_000, 1, grouped=False)} milj. "
            f"{fi_unit(ca_units['unit'])}"
        ),
        "ca_litres": format_millions(ca_litres["value"], fi_unit(ca_litres["unit"])),
        "ca_scope": ca_value["limitationFi"].split(".", 1)[0] + ".",
        "ca_publishers": unique_strings(
            concise_publisher(market_sources[source_id]["publisher"])
            for item in (ca_value, ca_units, ca_litres)
            for source_id in item["sourceIds"]
        ),
        "de_country": facts["country_name"](model["countryIso2"]),
        "de_rows": de_rows,
        "de_years": year_label(de_volumes),
        "de_first": format_millions(de_volumes[0]["value"], fi_unit(de_volumes[0]["unit"])),
        "de_latest": format_millions(de_volumes[-1]["value"], fi_unit(de_volumes[-1]["unit"])),
        "de_latest_year": observation_year(de_volumes[-1]),
        "de_latest_finality": fi_finality(de_volumes[-1].get("finality")).lower(),
        "de_model": (
            f"{format_millions(model['low'], model['currency'])}–"
            f"{format_millions(model['high'], model['currency'])}"
        ),
        "de_model_central": format_millions(model["central"], model["currency"]),
        "de_model_metrics": [
            (
                format_millions(model[key], model["currency"]),
                f"{format_compact_local_number(model_prices[key]['value'], maximum_decimals=2, grouped=False)} "
                f"{fi_unit(model_prices[key]['unit'])}",
            )
            for key in ("low", "central", "high")
        ],
        "de_model_price_count": len(model_prices),
        "de_model_confidence_fi": confidence_fi,
        "global_range": (
            f"{format_compact_local_number(min(global_values) / 1_000_000_000, minimum_decimals=1, grouped=False)}–"
            f"{format_compact_local_number(max(global_values) / 1_000_000_000, minimum_decimals=1, grouped=False)} "
            f"mrd {global_currency}"
        ),
        "global_metrics": global_metrics,
        "global_sources": "; ".join(facts["global_publishers"]),
        "market_sources": "; ".join(facts["market_publishers"]),
        "alert_sources": "; ".join(facts["alert_publishers"]),
        "priority_date_fi": fi_date(facts["priority_date"]),
        "b2_date_fi": fi_date(facts["b2_date"]),
        "infringement_claims_fi": fi_join([str(value) for value in facts["infringement_claims"]]),
        "nullity_reference": nullity_reference,
        "german_references": german_references,
        "cn_grant_date_fi": fi_date(facts["cn_grant"]["publicationDate"]),
        "grade_range": grade_labels[0] if len(grade_labels) == 1 else f"{grade_labels[0]}–{grade_labels[-1]}",
        "ep_publication": patent["epPublication"],
    }


def presentation_slide_strings(prs: Presentation) -> list[list[str]]:
    slides: list[list[str]] = []
    for slide in prs.slides:
        values: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    values.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            values.append(text)
        slides.append(values)
    return slides


def validate_finnish_deck_fact_tokens(prs: Presentation, ctx: dict[str, Any], deck_kind: str) -> None:
    """Fail closed if a mutable canonical fact is absent from its Finnish slide."""

    v = common_values(ctx)
    slide_strings = presentation_slide_strings(prs)
    exact: list[tuple[int, str]] = []
    contains: list[tuple[int, str]] = []

    if deck_kind == "short":
        exact.extend(
            (3, value)
            for value in (
                str(v["b2_claim_count"]),
                str(v["family_count"]),
                str(v["german_official_proceeding_count"]),
            )
        )
        exact.extend(
            (4, value)
            for value in (
                str(v["atlas_country_count"]),
                str(v["official_country_count"]),
                str(v["retail_donors"]),
            )
        )
        contains.extend(
            (slide, value)
            for slide, value in (
                (2, v["ep_publication"]),
                (2, str(v["b2_claim_count"])),
                (3, v["ep_publication"]),
                (4, v["global_range"]),
                (4, v["ca_value"]),
                (4, v["de_latest"]),
                (4, str(v["de_latest_year"])),
                (4, v["de_latest_finality"]),
            )
        )
    elif deck_kind == "medium":
        exact.extend(
            (4, value)
            for value in (str(v["priority_year"]), str(v["b2_claim_count"]), str(v["family_count"]))
        )
        exact.extend(
            (6, value)
            for value in (str(v["atlas_country_count"]), str(v["official_country_count"]), v["global_range"])
        )
        exact.append((11, f"{v['retail_donors']} retail-luovuttajaa"))
        contains.extend(
            (slide, value)
            for slide, value in (
                (3, v["ep_publication"]),
                (4, v["b2_date_fi"]),
                (4, v["patent"]["recordedProprietor"]),
                (4, v["ep_publication"]),
                (5, v["infringement_claims_fi"]),
                (5, v["de_infringement"]["reference"]),
                (6, str(v["global_year"])),
                (6, v["de_model"]),
                (6, str(v["retail_donors"])),
                (9, str(v["official_country_count"])),
            )
        )
    elif deck_kind == "large":
        exact.extend(
            (7, value)
            for value in (v["priority_date_fi"], v["b2_date_fi"], str(v["b2_claim_count"]))
        )
        exact.extend(
            (8, value)
            for value in (
                str(v["family_count"]),
                str(v["proceeding_count"]),
                str(v["national_register_check_count"]),
            )
        )
        exact.extend(
            (9, value)
            for value in (
                v["nullity_reference"],
                v["de_infringement"]["reference"],
                f"{v['de_nullity']['titleFi']} ({fi_date(v['de_nullity']['eventDate'])})",
                f"{v['de_nullity']['detailFi']} {v['de_nullity']['finalityFi']}",
                f"{v['de_infringement']['titleFi']} ({fi_date(v['de_infringement']['eventDate'])})",
                f"{v['de_infringement']['detailFi']} {v['de_infringement']['finalityFi']}",
            )
        )
        for alert in v["alerts"]:
            exact.extend(
                (
                    (11, v["country_name"](alert["jurisdictionCode"])),
                    (11, fi_date(alert["targetDate"])),
                    (11, alert["titleFi"]),
                )
            )
        for grade in v["grade_labels"]:
            exact.extend(((12, str(v["grade_counts"][grade])), (12, f"{grade}-luokan maata")))
        for row in v["market_country_rows"]:
            exact.extend((13, cell) for cell in row)
        exact.extend((14, value) for value in (v["ca_value"], v["ca_units"], v["ca_litres"]))
        for row in v["de_rows"]:
            exact.extend((15, cell) for cell in row)
        for value, label in v["de_model_metrics"]:
            exact.extend(((16, value), (16, label)))
        for value, label in v["global_metrics"]:
            exact.extend(((17, value), (17, label)))
        contains.extend(
            (slide, value)
            for slide, value in (
                (2, v["ep_publication"]),
                (2, str(v["german_official_proceeding_count"])),
                (2, str(v["atlas_country_count"])),
                (2, str(v["retail_donors"])),
                (6, fi_cardinal(v["b2_claim_count"])),
                (6, v["ep_publication"]),
                (10, v["cn_grant"]["publicationNumber"]),
                (10, v["cn_grant_date_fi"]),
                (11, fi_cardinal(v["alert_count"]).capitalize()),
                (11, str(v["fi_fee_year"])),
                (
                    11,
                    format_compact_local_number(
                        v["us_maintenance_year"],
                        minimum_decimals=1,
                        maximum_decimals=1,
                        grouped=False,
                    ),
                ),
                (12, str(v["atlas_country_count"])),
                (12, v["atlas_universe"]),
                (12, str(v["atlas_evidence_count"])),
                (12, str(v["un_member_count"])),
                (13, str(v["official_country_count"])),
                (14, str(v["ca_year"])),
                (16, str(v["model"]["year"])),
                (16, str(v["model_price_year"])),
                (16, v["model"]["limitationFi"]),
                (18, str(v["minimum_required_donors"])),
                (18, str(v["retail_donors"])),
            )
        )
    else:
        raise ValueError(f"Unknown Finnish deck kind: {deck_kind}")

    errors = []
    for slide_number, token in exact:
        if token not in slide_strings[slide_number - 1]:
            errors.append(f"slide {slide_number}: exact token {token!r}")
    for slide_number, token in contains:
        if not any(token in value for value in slide_strings[slide_number - 1]):
            errors.append(f"slide {slide_number}: substring {token!r}")
    if errors:
        raise AssertionError(
            f"{deck_kind} Finnish deck is not bound to canonical public facts:\n- " + "\n- ".join(errors)
        )


def regression_check_finnish_deck_fact_binding(ctx: dict[str, Any]) -> None:
    """Exercise positive fact propagation and fail-closed schema mutations."""

    def observation(context: dict[str, Any], observation_id: str) -> dict[str, Any]:
        return required_single(
            (
                item
                for item in context["market"]["observations"]
                if item.get("observationId") == observation_id
            ),
            f"regression observation {observation_id}",
        )

    def selected_model(context: dict[str, Any]) -> dict[str, Any]:
        return required_single(context["market"]["models"], "regression market model")

    def register_row(rows: list[dict[str, str]], prefix: str) -> dict[str, str]:
        return required_single(
            (item for item in rows if item["Väite"].startswith(prefix)),
            f"Evidence Register row beginning {prefix!r}",
        )

    def add_market_source(
        context: dict[str, Any],
        template_id: str,
        source_id: str,
        url: str,
    ) -> None:
        template = required_single(
            (item for item in context["market"]["sources"] if item.get("sourceId") == template_id),
            f"regression market source {template_id}",
        )
        source = copy.deepcopy(template)
        source["sourceId"] = source_id
        source["pageUrl"] = url
        source.pop("url", None)
        context["market"]["sources"].append(source)

    def expect_rejection(label: str, mutate) -> None:
        changed = copy.deepcopy(ctx)
        mutate(changed)
        try:
            canonical_facts(changed)
        except ValueError:
            return
        raise AssertionError(f"Invalid public-data mutation was accepted: {label}")

    with tempfile.TemporaryDirectory(prefix="pixan-fi-deck-fact-binding-") as directory:
        root = Path(directory)

        # Short deck: a changed commercial source value must alter the visible range.
        short_ctx = copy.deepcopy(ctx)
        commercial = [item for item in short_ctx["market"]["observations"] if item.get("geography") == "Global"]
        target = min(commercial, key=lambda item: item["value"])
        original_range = common_values(ctx)["global_range"]
        target["value"] += 1_000_000_000
        changed_range = common_values(short_ctx)["global_range"]
        if changed_range == original_range:
            raise AssertionError("Canonical market mutation did not change the Finnish global-range token")
        short_path = root / "short-fi.pptx"
        build_short_deck(short_ctx, short_path)
        short_strings = presentation_slide_strings(Presentation(short_path))
        if not any(changed_range in value for value in short_strings[3]):
            raise AssertionError("Canonical market mutation did not reach visible Finnish short-deck text")

        # Medium deck: amended-claim count is read from the current patent payload.
        medium_ctx = copy.deepcopy(ctx)
        medium_ctx["patent_history"]["patent"]["claimScopeSummaryEn"] = re.sub(
            r"B2 specification contains \d+ claims",
            "B2 specification contains 8 claims",
            medium_ctx["patent_history"]["patent"]["claimScopeSummaryEn"],
        )
        medium_path = root / "medium-fi.pptx"
        build_medium_deck(medium_ctx, medium_path)
        medium_strings = presentation_slide_strings(Presentation(medium_path))
        if "8" not in medium_strings[3]:
            raise AssertionError("Canonical patent mutation did not reach visible Finnish medium-deck text")

        # Large deck + register: latest CA observations, latest official judgment,
        # alert date and price-input citation must all propagate from one context.
        large_ctx = copy.deepcopy(ctx)
        add_market_source(
            large_ctx,
            "CA-HC-VAPING-SALES-2024",
            "TEST-CA-HC-VAPING-SALES-2025",
            "https://example.invalid/ca-2025",
        )
        for metric, value in (
            ("manufacturer_importer_shipments_value", 1_250_000_000),
            ("manufacturer_importer_shipments_units", 125_000_000),
            ("manufacturer_importer_shipments_liquid_volume", 1_300_000),
        ):
            base = required_single(
                (
                    item
                    for item in large_ctx["market"]["observations"]
                    if item.get("countryIso2") == "CA" and item.get("metric") == metric
                ),
                f"regression Canada {metric}",
            )
            latest = copy.deepcopy(base)
            latest["observationId"] = base["observationId"].replace("CA-2024-", "CA-2025-")
            latest["year"] = 2025
            latest["value"] = value
            latest["sourceIds"] = ["TEST-CA-HC-VAPING-SALES-2025"]
            large_ctx["market"]["observations"].append(latest)

        add_market_source(
            large_ctx,
            "INTASTE-SAMURAI-2026",
            "TEST-DE-PRICE-CENTRAL-2026",
            "https://example.invalid/de-price-central",
        )
        central_price = observation(large_ctx, "DE-2026-RETAIL-PRICE-BASE-EUR-PER-ML")
        central_price["value"] = 0.80
        central_price["sourceIds"] = ["TEST-DE-PRICE-CENTRAL-2026"]
        model = selected_model(large_ctx)
        volume = observation(large_ctx, model["inputIds"][0])
        model["central"] = volume["value"] * 1_000 * central_price["value"]

        legal_source = copy.deepcopy(
            required_single(
                (
                    item
                    for item in large_ctx["patent_history"]["sources"]
                    if item.get("sourceId") == "DE-LGMUC-7O3341-24"
                ),
                "regression German official-judgment source",
            )
        )
        legal_source["sourceId"] = "TEST-DE-LGMUC-7O9999-26"
        legal_source["pageUrl"] = "https://example.invalid/de-judgment-2026"
        legal_source.pop("url", None)
        large_ctx["patent_history"]["sources"].append(legal_source)
        legal_event = copy.deepcopy(
            required_single(
                (
                    item
                    for item in large_ctx["patent_history"]["proceedings"]
                    if item.get("proceedingId") == "DE-LGMUC-7O3341-24"
                ),
                "regression German infringement proceeding",
            )
        )
        legal_event.update(
            {
                "proceedingId": "TEST-DE-LGMUC-7O9999-26",
                "eventDate": "2026-05-01",
                "reference": "7 O 9999/26",
                "titleFi": "Uudempi virallinen loukkausratkaisu",
                "detailFi": "Virallinen ratkaisu koski vaatimuksia 2 ja 7 testituotteissa.",
                "detailEn": "The official judgment concerned claims 2 and 7 in the test products.",
                "finalityFi": "Testissä ilmoitettu muutoksenhakutila.",
                "sourceIds": ["TEST-DE-LGMUC-7O9999-26"],
            }
        )
        large_ctx["patent_history"]["proceedings"].append(legal_event)
        large_ctx["patent_history"]["summary"]["proceedingCount"] += 1
        fi_alert = required_single(
            (
                item
                for item in large_ctx["patent_history"]["diligenceAlerts"]
                if item.get("jurisdictionCode") == "FI"
            ),
            "regression Finnish alert",
        )
        fi_alert["targetDate"] = "2026-09-15"

        changed_values = common_values(large_ctx)
        large_path = root / "large-fi.pptx"
        build_large_deck(large_ctx, large_path)
        large_strings = presentation_slide_strings(Presentation(large_path))
        required_large_tokens = {
            9: [
                legal_event["reference"],
                f"{legal_event['titleFi']} ({fi_date(legal_event['eventDate'])})",
                f"{legal_event['detailFi']} {legal_event['finalityFi']}",
            ],
            11: [fi_date(fi_alert["targetDate"])],
            14: ["2025", changed_values["ca_value"]],
            16: [changed_values["de_model_central"]],
        }
        for slide_number, tokens in required_large_tokens.items():
            for token in tokens:
                if not any(token in value for value in large_strings[slide_number - 1]):
                    raise AssertionError(
                        f"Canonical mutation did not reach Finnish large deck slide {slide_number}: {token!r}"
                    )

        changed_rows = evidence_rows(large_ctx)
        ca_row = register_row(changed_rows, "Kanadan vuoden 2025 valmistaja-")
        if "https://example.invalid/ca-2025" not in ca_row["Lähde"]:
            raise AssertionError("Latest Canada source did not reach the Evidence Register")
        legal_row = register_row(changed_rows, "Saksan virallinen tapahtuma koski vaatimuksia 2 ja 7")
        if (
            legal_row["Päivämäärä"] != legal_event["eventDate"]
            or legal_event["detailFi"] != legal_row["Todiste"]
            or "https://example.invalid/de-judgment-2026" not in legal_row["Lähde"]
        ):
            raise AssertionError("Latest German judgment did not bind proof, date and source in the register")
        model_row = register_row(changed_rows, "Saksan vuoden 2025 verotetun nesteen")
        if "https://example.invalid/de-price-central" not in model_row["Lähde"]:
            raise AssertionError("Current model-input citation did not reach the Evidence Register")

    # A final legal outcome must drive both assumptions and gaps without stale
    # appeal-pending or finality-verification language.
    final_ctx = copy.deepcopy(ctx)
    for proceeding_type in ("patent_nullity", "patent_infringement"):
        final_event = required_single(
            (
                item
                for item in final_ctx["patent_history"]["proceedings"]
                if item.get("jurisdictionCode") == "DE"
                and item.get("proceedingType") == proceeding_type
            ),
            f"regression final German {proceeding_type}",
        )
        final_event["status"] = "final"
        final_event["appealStatus"] = "appeal_concluded"
        final_event["finalityEn"] = "The outcome is final."
        final_event["finalityFi"] = "Lopputulos on lainvoimainen."
    final_facts = canonical_facts(final_ctx)
    final_rows = evidence_rows(final_ctx)
    for selected, state_key in (
        (final_facts["de_nullity"], "de_nullity_process_state"),
        (final_facts["de_infringement"], "de_infringement_process_state"),
    ):
        if final_facts[state_key] != "final":
            raise AssertionError(f"Final proceeding was not classified final: {selected['proceedingId']}")
        final_row = required_single(
            (item for item in final_rows if item["Todiste"] == selected["detailFi"]),
            f"final Evidence Register row {selected['proceedingId']}",
        )
        if selected["finalityFi"] not in final_row["Oletukset"]:
            raise AssertionError("Final proceeding finality did not reach Evidence Register assumptions")
        if "vireillä" in " ".join(
            final_row[field] for field in ("Väite", "Oletukset", "Puutteet / tarvittava lisänäyttö")
        ).casefold():
            raise AssertionError("Final proceeding retained appeal-pending language")
        forbidden_final_gap = ("valitustila", "lainvoimaisuus")
        if any(token in final_row["Puutteet / tarvittava lisänäyttö"].casefold() for token in forbidden_final_gap):
            raise AssertionError("Final proceeding retained pending/unverified-finality gap language")

    # Same-year model inputs must remove the year-gap assumption and request.
    same_year_ctx = copy.deepcopy(ctx)
    same_year_model = selected_model(same_year_ctx)
    for input_id in same_year_model["rangeInputMap"].values():
        observation(same_year_ctx, input_id)["year"] = int(same_year_model["year"])
    same_year_model["yearMismatch"] = False
    same_year_model["limitationFi"] = "Ei havaittua myyntiä; vain verotettu neste ja kolme verkkokauppahintaa."
    same_year_row = register_row(evidence_rows(same_year_ctx), "Saksan vuoden 2025 verotetun nesteen")
    if same_year_row["Päivämäärä"] != "2025 määrä ja hinnat":
        raise AssertionError("Same-year model did not produce a same-year date label")
    if "kaikki syötteet vuodelta 2025" not in same_year_row["Oletukset"]:
        raise AssertionError("Same-year model assumptions did not describe the actual input years")
    if any(
        token in same_year_row["Puutteet / tarvittava lisänäyttö"].casefold()
        for token in ("saman vuoden", "mallivuoden", "vuosiero")
    ):
        raise AssertionError("Same-year model retained a year-gap request")

    current_model_row = register_row(evidence_rows(ctx), "Saksan vuoden 2025 verotetun nesteen")
    if "syötevuodet 2025 ja 2026" not in current_model_row["Oletukset"]:
        raise AssertionError("Year-mismatched model assumptions do not expose actual input years")
    if "mallivuoden 2025 hinnat" not in current_model_row["Puutteet / tarvittava lisänäyttö"]:
        raise AssertionError("Year-mismatched model gap does not request the missing model-year prices")

    # A newer non-official item must not replace the current official judgment.
    secondary_ctx = copy.deepcopy(ctx)
    secondary_source = copy.deepcopy(
        required_single(
            (
                item
                for item in secondary_ctx["patent_history"]["sources"]
                if item.get("sourceId") == "DE-LGMUC-7O3341-24"
            ),
            "regression German judgment source",
        )
    )
    secondary_source.update(
        {
            "sourceId": "TEST-DE-SECONDARY",
            "evidenceTier": "secondary",
            "pageUrl": "https://example.invalid/de-secondary",
        }
    )
    secondary_ctx["patent_history"]["sources"].append(secondary_source)
    secondary_event = copy.deepcopy(canonical_facts(ctx)["de_infringement"])
    secondary_event.update(
        {
            "proceedingId": "TEST-DE-SECONDARY",
            "eventDate": "2026-12-31",
            "reference": "Secondary only",
            "sourceIds": ["TEST-DE-SECONDARY"],
        }
    )
    secondary_ctx["patent_history"]["proceedings"].append(secondary_event)
    secondary_ctx["patent_history"]["summary"]["proceedingCount"] += 1
    if canonical_facts(secondary_ctx)["de_infringement"]["proceedingId"] != "DE-LGMUC-7O3341-24":
        raise AssertionError("A secondary source displaced the latest official German judgment")

    secondary_proceeding_ctx = copy.deepcopy(ctx)
    secondary_proceeding = copy.deepcopy(canonical_facts(ctx)["de_infringement"])
    secondary_proceeding.update(
        {
            "proceedingId": "TEST-DE-SECONDARY-PROCEEDING",
            "eventDate": "2027-01-01",
            "reference": "Secondary proceeding with official source",
            "evidenceTier": "secondary",
        }
    )
    secondary_proceeding_ctx["patent_history"]["proceedings"].append(secondary_proceeding)
    secondary_proceeding_ctx["patent_history"]["summary"]["proceedingCount"] += 1
    secondary_proceeding_facts = canonical_facts(secondary_proceeding_ctx)
    if secondary_proceeding_facts["de_infringement"]["proceedingId"] != "DE-LGMUC-7O3341-24":
        raise AssertionError("A secondary proceeding displaced the official German judgment via an official source")
    if secondary_proceeding_facts["german_official_proceeding_count"] != canonical_facts(ctx)["german_official_proceeding_count"]:
        raise AssertionError("A secondary proceeding inflated the official German judgment count")

    # All official German judgment records contribute to count/reference,
    # while latest-per-type records alone drive the current narratives.
    extra_judgment_ctx = copy.deepcopy(ctx)
    extra_judgment = copy.deepcopy(canonical_facts(ctx)["de_infringement"])
    extra_judgment.update(
        {
            "proceedingId": "TEST-DE-OFFICIAL-COST-ORDER",
            "proceedingType": "cost_order",
            "eventDate": "2026-06-01",
            "reference": "Test official cost order",
            "titleFi": "Virallinen kuluratkaisu",
            "titleEn": "Official cost order",
            "detailFi": "Virallinen kuluratkaisu testaa kokonaislukumäärää.",
            "detailEn": "Official cost order used to test the total count.",
        }
    )
    extra_judgment_ctx["patent_history"]["proceedings"].append(extra_judgment)
    extra_judgment_ctx["patent_history"]["summary"]["proceedingCount"] += 1
    base_facts = canonical_facts(ctx)
    extra_judgment_facts = canonical_facts(extra_judgment_ctx)
    if extra_judgment_facts["german_official_proceeding_count"] != base_facts["german_official_proceeding_count"] + 1:
        raise AssertionError("An additional official German judgment did not increase the all-record count")
    if extra_judgment_facts["de_nullity"]["proceedingId"] != base_facts["de_nullity"]["proceedingId"]:
        raise AssertionError("A different German judgment type displaced the nullity narrative")
    if extra_judgment_facts["de_infringement"]["proceedingId"] != base_facts["de_infringement"]["proceedingId"]:
        raise AssertionError("A different German judgment type displaced the infringement narrative")
    extra_references = common_values(extra_judgment_ctx)["german_references"]
    if "Test official cost order" not in extra_references:
        raise AssertionError("An additional official German judgment was omitted from the reference footer")

    # National-register metric is unique official jurisdictions and is independent of alerts.
    same_jurisdiction_ctx = copy.deepcopy(ctx)
    same_member = copy.deepcopy(same_jurisdiction_ctx["patent_history"]["familyMembers"][0])
    same_member["publicationNumber"] = "TEST-AU-SECOND-RECORD"
    same_member["jurisdictionCode"] = "AU"
    same_member["verificationLevel"] = "official_national_record"
    same_member["sourceIds"] = ["IPAU-PATENT-API"]
    same_jurisdiction_ctx["patent_history"]["familyMembers"].append(same_member)
    same_jurisdiction_ctx["patent_history"]["summary"]["familyRecordCount"] += 1
    if canonical_facts(same_jurisdiction_ctx)["national_register_check_count"] != canonical_facts(ctx)["national_register_check_count"]:
        raise AssertionError("A second record in one jurisdiction inflated the national-register count")
    extra_alert_ctx = copy.deepcopy(ctx)
    extra_alert = copy.deepcopy(extra_alert_ctx["patent_history"]["diligenceAlerts"][0])
    extra_alert.update({"alertId": "TEST-CA-ALERT", "jurisdictionCode": "CA"})
    extra_alert_ctx["patent_history"]["diligenceAlerts"].append(extra_alert)
    extra_alert_ctx["patent_history"]["summary"]["diligenceAlertCount"] += 1
    if canonical_facts(extra_alert_ctx)["national_register_check_count"] != canonical_facts(ctx)["national_register_check_count"]:
        raise AssertionError("An extra diligence alert changed the national-register count")

    # Negative regressions: malformed ids, observations and model schemas must fail closed.
    expect_rejection(
        "duplicate observationId",
        lambda changed: changed["market"]["observations"].append(copy.deepcopy(changed["market"]["observations"][0])),
    )
    expect_rejection(
        "duplicate modelId",
        lambda changed: changed["market"]["models"].append(copy.deepcopy(changed["market"]["models"][0])),
    )
    expect_rejection(
        "duplicate proceedingId",
        lambda changed: changed["patent_history"]["proceedings"].append(
            copy.deepcopy(changed["patent_history"]["proceedings"][0])
        ),
    )
    expect_rejection(
        "duplicate alertId",
        lambda changed: changed["patent_history"]["diligenceAlerts"].append(
            copy.deepcopy(changed["patent_history"]["diligenceAlerts"][0])
        ),
    )
    expect_rejection(
        "unsupported model formula",
        lambda changed: selected_model(changed).__setitem__("formula", "volume_litres * retail_price_eur_per_ml"),
    )
    expect_rejection(
        "duplicate model range input",
        lambda changed: selected_model(changed)["rangeInputMap"].__setitem__(
            "central", selected_model(changed)["rangeInputMap"]["low"]
        ),
    )
    expect_rejection(
        "range price missing from inputIds",
        lambda changed: selected_model(changed)["inputIds"].remove(
            selected_model(changed)["rangeInputMap"]["central"]
        ),
    )
    expect_rejection(
        "model volume-year mismatch",
        lambda changed: selected_model(changed).__setitem__("year", 2024),
    )
    expect_rejection(
        "model price wrong country",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("countryIso2", "CA"),
    )
    expect_rejection(
        "model price wrong metric",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("metric", "retail_price"),
    )
    expect_rejection(
        "model price wrong status",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("evidenceStatus", "commercial_estimate"),
    )
    expect_rejection(
        "model price wrong unit",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("unit", "EUR_per_litre"),
    )
    expect_rejection(
        "model price wrong currency",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("currency", "USD"),
    )
    expect_rejection(
        "model price wrong period",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("period", "calendar_year"),
    )
    expect_rejection(
        "non-increasing model range",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("value", 2.0),
    )
    expect_rejection(
        "model total mismatch",
        lambda changed: selected_model(changed).__setitem__("central", selected_model(changed)["central"] + 1_000_000),
    )
    expect_rejection(
        "Canada value wrong currency",
        lambda changed: observation(changed, "CA-2024-MANUFACTURER-IMPORTER-SHIPMENTS-VALUE").__setitem__("currency", "USD"),
    )
    expect_rejection(
        "Canada unit-count wrong unit",
        lambda changed: observation(changed, "CA-2024-MANUFACTURER-IMPORTER-SHIPMENTS-UNITS").__setitem__("unit", "box"),
    )
    expect_rejection(
        "Canada annual observation wrong period",
        lambda changed: observation(changed, "CA-2024-MANUFACTURER-IMPORTER-SHIPMENTS-VALUE").__setitem__("period", "point_in_time"),
    )
    expect_rejection(
        "Germany volume wrong unit",
        lambda changed: observation(changed, "DE-2025-TAXED-LIQUID-VOLUME-L").__setitem__("unit", "kg"),
    )
    expect_rejection(
        "Germany model volume wrong country",
        lambda changed: observation(changed, "DE-2025-TAXED-LIQUID-VOLUME-L").__setitem__("countryIso2", "CA"),
    )
    expect_rejection(
        "Germany model volume wrong metric",
        lambda changed: observation(changed, "DE-2025-TAXED-LIQUID-VOLUME-L").__setitem__("metric", "reported_e_liquid_volume"),
    )
    expect_rejection(
        "Germany model volume wrong status",
        lambda changed: observation(changed, "DE-2025-TAXED-LIQUID-VOLUME-L").__setitem__("evidenceStatus", "official_observed"),
    )
    expect_rejection(
        "Germany model volume wrong currency",
        lambda changed: observation(changed, "DE-2025-TAXED-LIQUID-VOLUME-L").__setitem__("currency", "EUR"),
    )
    expect_rejection(
        "Germany excise wrong currency",
        lambda changed: observation(changed, "DE-2025-SUBSTITUTES-EXCISE-RECEIPTS").__setitem__("currency", "PLN"),
    )
    expect_rejection(
        "global estimate wrong unit",
        lambda changed: observation(changed, "GLOBAL-2025-IMARC-COMMERCIAL-ESTIMATE").__setitem__("unit", "EUR"),
    )
    expect_rejection(
        "global estimate wrong currency",
        lambda changed: observation(changed, "GLOBAL-2025-IMARC-COMMERCIAL-ESTIMATE").__setitem__("currency", "EUR"),
    )
    expect_rejection(
        "Finland excise wrong currency",
        lambda changed: observation(changed, "FI-2025-NICOTINE-E-LIQUID-EXCISE-RECEIPTS").__setitem__("currency", "PLN"),
    )
    expect_rejection(
        "negative market observation",
        lambda changed: observation(changed, "DE-2026-RETAIL-PRICE-LOW-EUR-PER-ML").__setitem__("value", -0.44),
    )


def build_short_deck(ctx: dict[str, Any], path: Path) -> None:
    v = common_values(ctx)
    prs = new_presentation(ctx, "Pixan · suppea julkinen pankkidekki")
    cover_slide(prs, ctx, "Suppea päätöksenteon tiivistelmä", 6)
    slide_claim(
        prs, ctx, 2, "Rahoitusteesi perustuu näyttöön — ei markkinahypeen", "Rahoitusteesi",
        "Patentilla on dokumentoitu tekninen ydin ja merkittävää virallista oikeusnäyttöä, mutta vakuusarvo ei ole vielä todistettu.",
        [f"{v['patent']['epCentralStatusFi']} {v['ep_publication']} sisältää {v['b2_claim_count']} vaatimusta.", f"{v['de_country']}: {fi_cardinal(v['german_official_proceeding_count'])} virallista ratkaisua; lopullisuus- ja aluerajat säilyvät.", "Rahoituskelpoisuus syntyy, kun oikeudet, claim chartit, kassavirta ja riippumaton arvonmääritys sidotaan yhteen."],
        "Pankin lukutapa", "Nykyinen aineisto tukee jatkodiligenceä. Se ei vielä yksin tue tiettyä lainamäärää, osakearvoa tai patentin vakuusarvoa.",
        f"EPO; {v['de_country']} — viralliset tuomiot; WIPO",
    )
    slide_metrics(
        prs, ctx, 3, "Patenttihistoria antaa vahvan mutta rajatun ankkurin", "Patentti ja IP",
        [(str(v["b2_claim_count"]), f"vaatimusta {v['ep_publication']}:ssa"), (str(v["family_count"]), "perhejulkaisutietuetta"), (str(v["german_official_proceeding_count"]), f"virallista ratkaisua: {v['de_country']}")],
        "Todistettu oikeusnäyttö on arvokasta vain maakohtaisen voimassaolo- ja tuotekytkennän kanssa.",
        "Kirjattu haltija, prioriteetti ja EPO:n keskitetty lopputulos ovat vahvistettavissa. Koko perheen nykyinen kansallinen tila, rasitteet ja täytäntöönpano vaativat lisänäyttöä.",
        f"EPO Register; {v['ep_publication']}; {v['german_references']}",
    )
    slide_metrics(
        prs, ctx, 4, "Markkina-aineisto on läpinäkyvä mutta ei vielä valmis globaaliksi arvoksi", "Markkina",
        [(str(v["atlas_country_count"]), "maan tutkimusuniversumi"), (str(v["official_country_count"]), "maata virallisilla vuosihavainnoilla"), (str(v["retail_donors"]), "hyväksyttyä retail-arvon luovuttajaa")],
        f"Kaupallinen {v['global_range']} haarukka on sanity check — ei atlaksen oma maailmanestimaatti.",
        f"{v['ca_country']}: {v['ca_value']} toimitusarvo. {v['de_country']}: {v['de_latest']} verotettua nestettä vuonna {v['de_latest_year']} ({v['de_latest_finality']}). Mittareita ei summata keskenään.",
        v["market_sources"],
    )
    slide_table(
        prs, ctx, 5, "Puutteet ovat rajattavissa konkreettiseksi työohjelmaksi", "Riskit ja näyttö",
        ["Puutuva näyttö", "Miksi ratkaiseva", "Seuraava todiste"],
        [["Koko perheen kansallinen tila", "Määrittää toteuttamiskelpoiset oikeudet", "Asiamiehen allekirjoittama maamatriisi"], ["Tuote–vaatimusvertailut", "Rajaa relevantin myynnin", "Näytteet, teardown, claim chart"], ["Kassavirta ja sopimukset", "Määrittää takaisinmaksun", "Sopimukset, laskut, maksut"], ["Riippumaton arvonmääritys", "Määrittää skenaariot ja downside-riskin", "IVS-yhteensopiva arvio"]],
        "Pankkikelpoisuus paranee eniten, kun nämä neljä aukkoa suljetaan ennen arvokeskustelua.", "Evidence Register; WIPO IP valuation", [2.7, 3.7, 4.8],
    )
    closing_slide(
        prs, ctx, 6, "Ehdollinen eteneminen: 90 päivän diligence ennen rahoituspäätöstä",
        ["0–30 päivää: omistus-, rasite-, maksu- ja oikeusmatriisi.", "31–60 päivää: priorisoidut claim chartit, markkinamyynnin rajaus ja vastapuolidata.", "61–90 päivää: riippumaton arvonmääritys, kassavirtaskenaariot ja rahoitusehdot.", "Päätösportti: vain vahvistettuun oikeuteen ja takaisinmaksulähteeseen perustuva rakenne."],
        "Evidence Register; WIPO IP Finance",
    )
    if len(prs.slides) != 6:
        raise AssertionError("Short deck must have 6 slides")
    validate_finnish_deck_fact_tokens(prs, ctx, "short")
    save_presentation(prs, path)


def build_medium_deck(ctx: dict[str, Any], path: Path) -> None:
    v = common_values(ctx)
    p = ctx["patent_history"]["patent"]
    prs = new_presentation(ctx, "Pixan · keskikokoinen julkinen pankkidekki")
    cover_slide(prs, ctx, "12 dian rahoitus- ja teknologia-arvio", 12)
    add_text(prs.slides[0], "1 · Rahoitusteesi", 0.62, 4.72, 6.0, 0.35, size=15, color=TEAL, bold=True)
    # Slide 1 is the cover and carries the financing thesis.
    slide_claim(prs, ctx, 2, "Ongelma on yhtä paljon evidenssissä kuin teknologiassa", "2 · Ongelma",
                "Laaja markkinapuhe ei muutu pankkikelpoiseksi vakuudeksi ilman maakohtaista oikeus-, tuote-, myynti- ja kassavirtaketjua.",
                ["Markkinamittarit ovat hajanaisia: toimitukset, verot, litrat ja kaupalliset arviot eivät ole sama asia.", "Perhejulkaisu ei todista nykyistä voimassaoloa tai rasitteettomuutta.", "Oikeustuomio ei yksin todista maksettua korvausta tai globaalia rojaltipohjaa."],
                "Pankin ydinkysymys", "Mikä todennettu omaisuus ja kassavirta kattaa lainan myös downside-tapauksessa?", "Atlas readiness; WIPO" )
    slide_claim(prs, ctx, 3, "Ratkaisu ohjaa höyrystimen tehoa mitatun resistanssin perusteella", "3 · Patentoitu ratkaisu",
                p["inventionSummaryFi"],
                ["Ohjaus perustuu tallennettuihin resistanssi–tehoarvoihin.", "Tehosuhde kuvataan ei-suoraan-verrannolliseksi.", "Käyttäjäsäätö tapahtuu nollasta poikkeavan minimin ja tallennetun maksimin välillä."],
                "Rajaus", "Tämä on julkaistun patenttiselityksen yleiskielinen tiivistelmä. Vain kyseisen maan operative claims ja tuotekohtainen claim chart ratkaisevat.", v["ep_publication"] )
    slide_metrics(prs, ctx, 4, "IP-ydin on dokumentoitu; maakohtainen käyttökelpoisuus on vielä täsmäytettävä", "4 · Patentti ja IP-status",
                  [(str(v["priority_year"]), "varhaisin prioriteettivuosi"), (str(v["b2_claim_count"]), "B2-vaatimusta"), (str(v["family_count"]), "perhejulkaisutietuetta")],
                  f"{v['patent']['epCentralStatusFi']} B2 julkaistiin {v['b2_date_fi']}.",
                  f"Julkisen rekisterin kirjattu haltija: {v['patent']['recordedProprietor']}. Koko perheen nykyistä kansallista voimassaoloa, maksuja ja rasitteita ei ole vielä varmennettu.", f"EPO Register; {v['ep_publication']}" )
    slide_claim(prs, ctx, 5, "Tekninen erottautuminen on vaatimuksissa — ei tuotteen ulkonäössä", "5 · Tekninen erottautuminen",
                "Rahoitusarvioinnin tulee jäljittää jokainen olennainen vaatimuspiirre tuotteeseen ja maahan.",
                ["Patenttiselitys antaa teknisen hypoteesin ja dokumentoidun suojatekstin.", f"{v['de_country']}: ratkaisu tukee tiettyjen tuotteiden ja vaatimusten {v['infringement_claims_fi']} kytkentää.", "Riippumattomat testit, teardown-dossierit ja muiden maiden claim chartit puuttuvat julkisesta aineistosta."],
                "Todistusketju", "Näyte → hallussapitoketju → teardown → mittausdata → claim chart → asiamiehen tarkastus → relevantti myynti.", f"{v['ep_publication']}; {v['de_infringement']['reference']}" )
    slide_metrics(prs, ctx, 6, "Markkinakoko on tällä hetkellä haarukka ja havaintokokoelma — ei yksi luku", "6 · Markkinan koko ja rajaus",
                  [(str(v["atlas_country_count"]), "maan tutkimusrivit"), (str(v["official_country_count"]), "maata vuosihavainnoilla"), (v["global_range"], f"kaupallinen globaali {v['global_year']} haarukka")],
                  f"{v['de_country']}: nestemallin luottamustaso on {v['de_model_confidence_fi']} ja vaihteluväli {v['de_model']}; sitä ei saa esittää havaittuna myyntinä.",
                  f"Hyväksyttyjä virallisia kansallisia consumer-retail-arvoja on {v['retail_donors']}. Atlas ei julkaise omaa maailmanestimaattia ennen metodisten porttien täyttymistä.", f"Market-values; {v['global_sources']}" )
    slide_claim(prs, ctx, 7, "Asiakkuus on vielä validoitava kolmessa eri päätöksentekologiikassa", "7 · Asiakkaat ja ostoperuste",
                "Mahdolliset segmentit ovat valmistajat, teknologiatoimittajat sekä IP-rahoittajat tai ostajat.",
                ["Valmistaja ostaa toimintarauhaa, oikeusvarmuutta tai teknologiaa — jos relevantti tuote ja alue osoitetaan.", "Teknologiatoimittaja arvioi integroitavuutta, vapautta toimia ja yksikkötaloutta.", "Rahoittaja arvioi omistusta, realisoitavaa kassavirtaa, kontrollia ja downside-arvoa."],
                "Nykyinen näyttö", "Julkinen aineisto ei vahvista nimettyjä asiakkaita, piloteja, ostoaikomuksia tai allekirjoitettuja lisenssejä.", "WIPO licensing; WIPO IP Finance" )
    slide_table(prs, ctx, 8, "Kilpailu on erotettava vaihtoehtoisista teknologioista ja patenttiriskistä", "8 · Kilpailu ja vaihtoehdot",
                ["Vertailutaso", "Kysymys", "Nykytila"],
                [["Tuoteratkaisut", "Ratkaiseeko tuote saman ongelman eri tavalla?", "Puuttuu"], ["Patenttiperheet", "Onko suojan päällekkäisyys tai FTO-riski?", "Puuttuu"], ["Kaupalliset vaihtoehdot", "Lisenssi, luovutus, sovinto vai oma käyttö?", "Oletus"], ["Ei-patentoidut vaihtoehdot", "Voiko suorituskyvyn saavuttaa kiertämällä vaatimukset?", "Puuttuu"]],
                "Kilpailija-analyysi on pankin teknologia-arvioinnin olennainen avoin työpaketti.", "Evidence Register; EPO; WIPO", [2.6, 5.6, 2.1] )
    slide_table(prs, ctx, 9, "Validoinnin vahvuus on epätasainen ja siksi näkyvästi luokiteltu", "9 · Validointi ja nykyinen näyttö",
                ["Väitealue", "Luokitus", "Mitä on", "Mitä puuttuu"],
                [["Patenttiydin", "Vahvistettu", "EPO-rekisteri ja B2", "Maakohtainen statusmatriisi"], [f"{v['de_country']}n prosessit", "Vahvistettu", f"{fi_cardinal(v['german_official_proceeding_count']).capitalize()} virallista ratkaisua", "Lopullisuus, täytäntöönpano, kassa"], ["Markkinahavainnot", "Vahvistettu", f"{v['official_country_count']} maan viralliset luvut", "Yhteismitallinen retail-arvo"], ["Tekninen testaus", "Puuttuu", "Patenttiselitys", "Riippumaton testi ja raakadata"], ["Kaupallinen näyttö", "Puuttuu", "Kaupallistamiskehys", "Asiakkaat, sopimukset, maksut"]],
                "Todistettu oikeusnäyttö ei korvaa teknistä, kaupallista tai taloudellista validointia.", "Evidence Register", [2.5, 1.8, 3.8, 4.4] )
    slide_process(prs, ctx, 10, "Kaupallistaminen etenee porttien kautta, ei massaväitteillä", "10 · Kaupallistamismalli",
                  [("Oikeusmatriisi", "Vahvista maa, haltija, maksu, rasite, operative claims ja määräpäivät."), ("Tuotedossier", "Hanki näyte, dokumentoi ketju ja tee asiamiehen tarkastama claim chart."), ("Kohdepisteytys", "Arvioi oikeus, näyttö, myynti, vastapuoli, kustannus ja täytäntöönpano."), ("Rajattu pilotti", "Testaa lisenssi-, sovinto-, luovutus- tai rahoitusreitti auditoitavin ehdoin.")],
                  "Vasta toteutunut sopimus tai saaminen luo kassavirtanäyttöä.", "WIPO licensing; dispute resolution; IPscore" )
    slide_table(prs, ctx, 11, "Taloudellinen malli on rakennettava lähteistä, ei markkinaosuusoletuksesta", "11 · Taloudellinen malli ja herkkyydet",
                ["Silta", "Todennettava syöte", "Nykytila"],
                [["Kokonaismarkkina", "Yhteismitallinen maakohtainen myynti", f"{v['retail_donors']} retail-luovuttajaa"], ["Kohdistettava myynti", "Voimassaolo × tuote × aika × maa", "Puuttuu"], ["Rojaltipohja", "Claim-mapped net sales", "Puuttuu"], ["Kassavirta", "Sopimusehdot, kulut, verot, viive", "Puuttuu"], ["Vakuusarvo", "Downside-realisointi ja kontrolli", "Puuttuu"]],
                "Näytä downside/base/upside vasta, kun jokainen sillan syöte on dokumentoitu.", "WIPO IP valuation; Market-values", [2.5, 5.7, 1.8] )
    closing_slide(prs, ctx, 12, "Rahoituspäätös vasta neljän kriittisen aukon sulkeuduttua",
                  ["1. Asiamiehen allekirjoittama omistus-, rasite-, maksu- ja oikeusmatriisi.", "2. Priorisoitujen tuotteiden claim chartit ja dokumentoitu relevantti myynti.", "3. Toteutuneet tai sopimuspohjaiset kassavirrat sekä auditoidut taloustiedot.", "4. Riippumaton arvonmääritys ja downside-vakuusanalyysi.", "Seuraava päätös: hyväksytäänkö 90 päivän kontrolloitu diligence-vaihe?"], "Evidence Register; WIPO" )
    add_text(prs.slides[-1], "12 · Riskit, hallintatoimet ja seuraavat vaiheet", 0.64, 0.82, 8.7, 0.28, size=12, color=TEAL, bold=True)
    if len(prs.slides) != 12:
        raise AssertionError("Medium deck must have 12 slides")
    validate_finnish_deck_fact_tokens(prs, ctx, "medium")
    save_presentation(prs, path)


def build_large_deck(ctx: dict[str, Any], path: Path) -> None:
    v = common_values(ctx)
    p = ctx["patent_history"]["patent"]
    obs = ctx["observations"]
    model = v["model"]
    prs = new_presentation(ctx, "Pixan · laaja julkinen pankkidekki")
    cover_slide(prs, ctx, "Laaja 30 dian tutkija- ja rahoituspäätöspaketti", 30)
    slide_claim(prs, ctx, 2, "Rahoitettavuus on mahdollisuus, ei nykyinen johtopäätös", "Rahoitusteesi",
                "Viralliset patentti- ja oikeuslähteet oikeuttavat jatkodiligencen; ne eivät vielä osoita vakuusarvoa tai takaisinmaksua.",
                [f"EPO:n muutettu {v['ep_publication']} ja {v['de_country']}n {fi_cardinal(v['german_official_proceeding_count'])} virallista ratkaisua muodostavat oikeusnäytön ankkurin.", f"Markkina-aineisto kattaa {v['atlas_country_count']} maan tutkimusrungon, mutta yhteismitallisia retail-luovuttajia on {v['retail_donors']}.", "Rahoitusrakenne tarvitsee kansalliset oikeudet, claim-mapped sales -sillan, kassavirran ja riippumattoman arvonmäärityksen."],
                "Suositus", "Avaa 90 päivän ehdollinen diligence, ei lopullista luottopäätöstä.", f"EPO; {v['de_country']} — viralliset tuomiot; WIPO" )
    slide_table(prs, ctx, 3, "Kolme perustetta jatkaa — ja kolme rajaa olla kiirehtimättä", "Rahoitusteesi",
                ["Vahva signaali", "Mitä se tukee", "Mitä se ei todista"],
                [[v["epo_proceeding"]["titleFi"], "Dokumentoitu patenttiydin", "Kansallinen voimassaolo"], [v["de_nullity"]["titleFi"], "Pätevyyden kansallinen signaali", "Lainvoima tai globaali pätevyys"], [v["de_infringement"]["titleFi"], "Tuote- ja claim-kohtainen signaali", "Muu tuote, maa tai maksettu korvaus"]],
                "Pankin tulee hinnoitella vain se osa näytöstä, jonka omistus, toteutettavuus ja kassavirta on vahvistettu.", f"EPO; {v['german_references']}", [3.1,4.6,4.6] )
    slide_claim(prs, ctx, 4, "Todennusketju katkeaa tällä hetkellä ennen kassavirtaa", "Ongelma",
                "Patentti → voimassa oleva maaoikeus → relevantti tuote → relevantti myynti → sopimus tai tuomio → maksu → velanhoito.",
                ["Ketjun alku on osin vahvistettu virallisilla lähteillä.", f"Tuote- ja myyntikytkentä on vahvistettu vain rajatusti: {v['de_country']}, {v['de_infringement']['reference']}.", "Sopimus-, maksu-, yhtiötalous- ja vakuusarvonäyttö puuttuu julkisesta paketista."],
                "Luottoriski", "Yksi vahvistamaton lenkki voi muuttaa nimellisen markkina-arvon nollaksi realisoitavassa downside-skenaariossa.", "Evidence Register; WIPO" )
    slide_claim(prs, ctx, 5, "Patentoitu ratkaisu ohjaa lämmitystehoa resistanssitiedon avulla", "Ratkaisu",
                p["inventionSummaryFi"],
                ["Mittaus: lämmityselementin resistanssi.", "Tieto: tallennetut resistanssi–tehoarvot.", "Ohjaus: lämmittimelle syötetty teho ja käyttäjäsäädön rajat."],
                "Tulkintaraja", "Yleiskielinen kuvaus ei korvaa maakohtaista claim constructionia tai tuotekohtaista vaatimusanalyysiä.", v["ep_publication"] )
    slide_table(prs, ctx, 6, f"B2-tekstin {fi_cardinal(v['b2_claim_count'])} vaatimusta ovat arvioinnin lähtöpiste", "Patentti",
                ["Elementti", "Julkinen kuvaus", "Diligence-testi"],
                [["Ohjain / menetelmä", "Sähköisen höyrystimen hallinta", "Mikä tuoteversio ja maa?"], ["Resistanssin mittaus", "Lämmityselementin resistanssi", "Miten ja milloin mitataan?"], ["Tallennetut arvot", "Resistanssi–teho-kytkentä", "Missä data sijaitsee?"], ["Ei-proportionaalisuus", "Tehosuhteen määritelty luonne", "Täyttyykö vaatimuspiirre?"], ["Säätörajat", "Minimi ja tallennettu maksimi", "Miten käyttöliittymä toteuttaa rajat?"]],
                "Jokainen vaatimuspiirre on osoitettava todisteella; ominaisuusluettelo ei yksin riitä.", v["ep_publication"], [2.2,5.9,3.1] )
    slide_metrics(prs, ctx, 7, "IP-historian ydintapahtumat ovat virallisesti jäljitettävissä", "Patentti",
                  [(v["priority_date_fi"], "varhaisin prioriteetti"), (v["b2_date_fi"], "B2-julkaisu"), (str(v["b2_claim_count"]), "muutettua vaatimusta")],
                  v["epo_proceeding"]["finalityFi"],
                  "Kansalliset post-grant-oikeudet elävät erillään EPO:n keskusmenettelystä.", "EPO Register; B2" )
    slide_metrics(prs, ctx, 8, f"Perhejulkaisuja on {v['family_count']} — voimassa olevien maiden määrää ei vielä väitetä", "Maantieteellinen kattavuus",
                  [(str(v["family_count"]), "perhejulkaisutietuetta"), (str(v["proceeding_count"]), "tunnistettua menettelyä"), (str(v["national_register_check_count"]), "virallista kansallista rekisteritarkistusta"), ("?", "täytäntöönpanokelpoista maata")],
                  "Julkaisureitti on kartta, ei omistus- tai voimassaolotodistus.",
                  "Asiamiehen maamatriisin tulee sisältää haltija, operative claims, vuosimaksu, kuitti, rasitteet, UPC-asema ja seuraava määräpäivä.", "EPO family; kansalliset rekisterit" )
    slide_table(prs, ctx, 9, f"{v['de_country']}n ratkaisut ovat vahvaa näyttöä tarkasti rajatussa kehyksessä", "Oikeudellinen näyttö",
                ["Viite", "Virallinen tapahtuma", "Rajaus ja lopullisuus"],
                [
                    [
                        v["nullity_reference"],
                        f"{v['de_nullity']['titleFi']} ({fi_date(v['de_nullity']['eventDate'])})",
                        f"{v['de_nullity']['detailFi']} {v['de_nullity']['finalityFi']}",
                    ],
                    [
                        v["de_infringement"]["reference"],
                        f"{v['de_infringement']['titleFi']} ({fi_date(v['de_infringement']['eventDate'])})",
                        f"{v['de_infringement']['detailFi']} {v['de_infringement']['finalityFi']}",
                    ],
                ],
                f"{v['de_country']}n näyttö ei automaattisesti siirry toiseen tuotteeseen, vastapuoleen tai valtioon.", f"Viralliset tuomiot: {v['german_references']}", [3.6,4.2,5.2] )
    slide_claim(prs, ctx, 10, f"{v['country_name']('CN')}n asia ei ole loukkausvoitto", "Oikeudellinen näyttö",
                "Julkinen sekundäärinen docket-tieto viittaa hylätyn hakemuksen hakijapuolen uudelleentarkastukseen.",
                [f"Menettely luokitellaan {v['cn_review']['proceedingType'].replace('_', ' ')} -asiaksi.", "Virallista päätöstä ja tarkkoja perusteluja ei saatu julkiseen pakettiin.", f"{v['cn_grant']['publicationNumber']} julkaistiin myöhemmin myönnettynä {v['cn_grant_date_fi']}."],
                "Käyttöraja", f"Asiaa ei saa kuvata {v['country_name']('CN')}n loukkausoikeudenkäynniksi, vastapuolen mitätöintiasiaksi tai kassavirtanäytöksi.", "RPX; CNIPA guidance; EPO family" )
    slide_table(prs, ctx, 11, f"{fi_cardinal(v['alert_count']).capitalize()} ajankohtaista IP-hälytystä vaatii dokumentoidun omistajan", "Riskit",
                ["Maa / alue", "Päivä", "Toimi"],
                v["alert_rows"],
                "Määräpäivävalvonta on vakuusarvon operatiivinen kontrolli, ei hallinnollinen sivuseikka.", v["alert_sources"], [2.4,2.0,7.1] )
    slide_metrics(prs, ctx, 12, f"Atlas: {v['atlas_country_count']} maata ({v['atlas_universe']}); tutkimusrunko ei ole markkina-arvo", "Markkina",
                  [(str(v["grade_counts"][grade]), f"{grade}-luokan maata") for grade in v["grade_labels"]],
                  f"{v['atlas_evidence_count']} evidenssimerkintää; universumin YK-osuus on {v['un_member_count']} jäsenvaltiota.",
                  f"Luokat {v['grade_range']} kuvaavat evidenssikypsyyttä, eivät markkinan kokoa tai kaupallista houkuttelevuutta.", "UN; Atlas" )
    slide_table(prs, ctx, 13, f"Viralliset havainnot {v['official_country_count']} maasta mittaavat eri asioita", "Markkina",
                ["Maa", "Vuosi", "Virallinen havainto"],
                v["market_country_rows"],
                "Litroja, toimitusarvoja ja valmisteveroja ei summata yhdeksi markkinaksi.", v["market_sources"], [1.6,2.1,8.3] )
    slide_metrics(prs, ctx, 14, f"{v['ca_country']} on vahva toimitusmyynnin ankkuri, ei retail-arvon luovuttaja", v["ca_country"],
                  [(v["ca_value"], "valmistaja-/maahantuojatoimitukset"), (v["ca_units"], "raportoitua yksikkömäärää"), (v["ca_litres"], "raportoitua nestettä")],
                  v["ca_scope"],
                  v["ca_value_observation"]["limitationFi"], f"{'; '.join(v['ca_publishers'])} {v['ca_year']}" )
    slide_table(prs, ctx, 15, f"{v['de_country']}n virallinen nestemäärä kasvoi, mutta arvo vaatii hintaoletuksen", v["de_country"],
                ["Vuosi", "Verotettu neste", "Valmistevero", "Lopullisuus"],
                v["de_rows"],
                "Verotettu nestemäärä ei sisällä laitteita, verotonta tai laitonta kauppaa eikä ole retail-arvo.", "; ".join(unique_strings(concise_publisher(v["market_source_map"][source_id]["publisher"]) for item in v["de_volume_observations"] + v["de_excise_observations"] for source_id in item["sourceIds"])), [1.4,2.3,2.1,1.8] )
    slide_metrics(prs, ctx, 16, f"{v['de_country']}n mallihaitari näyttää herkkyyden — ei havaittua myyntiä", "Mallinnus",
                  v["de_model_metrics"],
                  f"Kaava: {model['formula']}.",
                  model["limitationFi"], "; ".join(unique_strings(concise_publisher(v["market_source_map"][source_id]["publisher"]) for source_id in v["model_source_ids"])) )
    slide_metrics(prs, ctx, 17, "Kaupalliset globaaliarviot ovat sanity check, eivät oma estimaatti", "Globaali markkina",
                  v["global_metrics"],
                  "Haarukka on leveä, koska tuoterajaukset ja metodit voivat poiketa.",
                  "Arvioita verrataan keskenään. Niitä ei summata, eikä haarukkaa käytetä automaattisesti rojaltipohjana.", v["global_sources"] )
    slide_claim(prs, ctx, 18, f"Hyväksyttävä maailmanestimaatti tarvitsee vähintään {v['minimum_required_donors']} yhteensopivaa luovuttajaa", "Metodi",
                f"Nykyinen comparable consumer-retail donor count on {v['retail_donors']}; hard gate ei täyty.",
                ["Sama tuoterajaus ja kalenterivuosi.", "Kuluttajavähittäisarvo, ei toimitus-, vero- tai volyymiproxy.", "Alue- ja sääntelytyyppien peitto sekä suora validointi suurissa talouksissa."],
                "Vasta sitten", "Trianguloi kysyntä-, vero-, tulli-, yritys- ja hintamenetelmät. Vertaa tuloksia; älä lisää vaihtoehtoisia arvioita yhteen.", "Market-values modelReadiness" )
    slide_table(prs, ctx, 19, "Markkinan ja patentin väliin tarvitaan viisi läpinäkyvää suodatinta", "Arvosilta",
                ["Taso", "Suodatin", "Näyttö"],
                [["1. Kokonaismarkkina", "Tuote- ja mittarirajaus", "Osittainen"], ["2. Oikeusalue", "Voimassa oleva operative claim", "Puuttuu globaalisti"], ["3. Relevantit tuotteet", "Claim chart", f"Rajattu näyttö: {v['de_country']}"], ["4. Relevantti myynti", "Tuote × maa × aika × net sales", "Puuttuu"], ["5. Kassavirta", "Rojalti/sovinto − kulut − verot − viive", "Puuttuu"]],
                "Vain alimman tason kassavirta voi palvella velkaa; ylempi markkina ei sellaisenaan voi.", "WIPO valuation; Evidence Register", [2.1,4.1,4.6] )
    slide_claim(prs, ctx, 20, "Mahdolliset asiakkaat ovat hypoteeseja, eivät vielä näyttöä", "Asiakkaat",
                "Segmentointi johdetaan kaupallistamisreiteistä: valmistajat, teknologiatoimittajat sekä IP-rahoittajat ja ostajat.",
                ["Valmistajalle arvo voi olla lisenssi, toimintarauha tai sovinto.", "Teknologiatoimittajalle arvo voi olla integroitava toiminto tai oikeusasema.", "Rahoittajalle arvo on kontrolloitava, realisoitava kassavirta ja downside-suoja."],
                "Validointi", "Dokumentoi 10–15 haastattelua, päätöskriteerit, vastalauseet, budjetti ja seuraava askel. Julkinen aineisto ei vielä sisällä näitä.", "WIPO licensing; IP Finance" )
    slide_table(prs, ctx, 21, "Kilpailukartta on rakennettava neljälle rinnakkaiselle tasolle", "Kilpailu",
                ["Taso", "Analyysi", "Tuotos"],
                [["Tuotteet", "Toiminto, arkkitehtuuri, hinta, markkina", "Vertailumatriisi"], ["Patentit", "Claims, prior art, status, FTO", "Patenttilandscape"], ["Kiertoratkaisut", "Voiko suorituskyvyn toteuttaa eri tavalla?", "Design-around-arvio"], ["Kaupallinen vaihtoehto", "Lisenssi, hankinta, oma kehitys, riitely", "Buy/build/license/litigate-malli"]],
                "Tekninen etu ei ole uskottava ennen kuin vaihtoehdot on kuvattu ja lähteistetty.", "Evidence Register; EPO" , [2.0,5.3,3.5] )
    slide_process(prs, ctx, 22, "Tuotevalidointi rakentuu katkeamattomaksi todisteketjuksi", "Tekninen validointi",
                  [("Näyte", "Osta oikeasta maasta ja ajankohdasta; dokumentoi myyjä, tuoteversio ja sarjatiedot."), ("Hallussapitoketju", "Tallenna vastaanotto, säilytys, avaaminen, kuvaus ja tiedostojen hashit."), ("Teardown ja testi", "Tee riippumaton mittausprotokolla, raakadata ja toistettavuustesti."), ("Claim chart", "Mapita jokainen vaatimuspiirre todisteeseen ja asiamiehen johtopäätökseen.")],
                  "Markkinamyynti voidaan kohdistaa vasta, kun tuoteidentiteetti ja claim-kytkentä ovat hallittuja.", f"{v['ep_publication']}; {v['de_infringement']['reference']}" )
    slide_process(prs, ctx, 23, "Kaupallistaminen aloitetaan kovista porteista", "Kaupallistaminen",
                  [("Vahvista oikeus", "Ei yhteydenottoa ilman maakohtaista omistusta, voimassaoloa, rasite- ja määräpäivätietoa."), ("Vahvista tuote", "Ei väitettä ilman tuotenäytettä, claim chartia ja paikallisen teon dokumentointia."), ("Pisteytä kohde", "Arvioi myynti, näyttö, vastapuoli, toimivalta, kustannus ja perittävyys."), ("Testaa pilotti", "Neuvottele lisenssi-, sovinto-, luovutus- tai rahoitusvaihtoehdot kontrolloidusti.")],
                  "Yhteydenottokirje ei automaattisesti katkaise vanhentumista tai todista tiedoksiantoa.", "WIPO dispute resolution; EPO IPscore" )
    slide_table(prs, ctx, 24, "Kaupallistamisreitit eroavat kassavirran, kontrollin ja riskin suhteen", "Kaupallistaminen",
                ["Reitti", "Kassavirta", "Keskeinen riski"],
                [["Lisenssi", "Upfront + jatkuva rojalti", "Rojaltipohjan auditointi"], ["Sovinto", "Kertamaksu / vaiheistus", "Ei automaattista ennakkotapausta"], ["Luovutus", "Kauppahinta", "Luovutetaan tuleva upside"], ["IP-rahoitus", "Velka tai revenue share", "Vakuusarvo ja kontrollit"], ["Prosessirahoitus", "Kuluihin sidottu pääoma", "Korkea kustannus ja lopputulosriski"]],
                "Valitse reitti vasta, kun oikeus, vastapuoli, kassavirta ja downside on todennettu.", "WIPO licensing; IP Finance; TPLF mapping", [2.1,4.1,5.2] )
    slide_table(prs, ctx, 25, "Taloudellinen malli alkaa todennettavista syötteistä", "Taloudellinen malli",
                ["Syöte", "Lähde", "Status"],
                [["Maakohtainen relevantti myynti", "Viranomainen / vastapuolen disclosure", "Puuttuu"], ["Claim-mapped osuus", "Asiamiehen tuotedossier", "Puuttuu"], ["Rojalti tai vahinko", "Sopimus / oikeudellinen analyysi", "Puuttuu"], ["Ajoitus ja perittävyys", "Prosessi- ja vastapuolianalyysi", "Puuttuu"], ["Kulut ja verot", "Budjetti / veroasiantuntija", "Puuttuu"], ["Diskontto ja downside", "Riippumaton arvonmääritys", "Puuttuu"]],
                "Nykyisestä julkisesta paketista ei voi johtaa pankkikelpoista NPV:tä.", "Evidence Register; WIPO valuation", [3.3,5.6,1.7] )
    slide_table(prs, ctx, 26, "Herkkyydet on sidottava todellisiin riskeihin", "Herkkyydet",
                ["Ajuri", "Downside", "Base", "Upside"],
                [["Oikeusalueet", "Vain vahvistettu maa", "Priorisoidut maat", "Laajempi varmennettu peitto"], ["Claim-osuvuus", "Yksi tuote", "Validoitu portfolio", "Laaja mutta dokumentoitu osuus"], ["Rojalti / korvaus", "Asiantuntijan alaraja", "Vertailuehdot", "Vain todennettu yläraja"], ["Ajoitus", "Valitus ja pitkä perintä", "Sopimuspolku", "Upfront-rakenne"], ["Kulut", "Täysi riitelybudjetti", "Rajattu pilotti", "Vastapuolen kattamat kulut"]],
                "Taulukkoon ei syötetä prosentteja ennen lähdetodisteita; tämä dia määrittää skenaarioiden rakenteen.", "WIPO valuation", [2.7,3.1,2.0,2.0] )
    slide_table(prs, ctx, 27, "Viisi kysymystä määrittää pankin seuraavan päätöksen", "Tutkijan kysymykset",
                ["#", "Kysymys", "Vaadittu vastaus"],
                [["1", "Mitä tarkalleen omistetaan ja missä se on voimassa?", "Allekirjoitettu oikeusmatriisi"], ["2", "Mikä tuote täyttää mitkä vaatimuspiirteet?", "Claim chart + riippumaton testi"], ["3", "Mikä on todennettu relevantti myynti?", "Tuote–maa–aika-net sales"], ["4", "Mistä ja milloin velanhoitokassa syntyy?", "Sopimukset, maksut, ennuste"], ["5", "Mitä vakuus realisoi downside-tilanteessa?", "Riippumaton arvo + toteutuspolku"]],
                "Jos yksikin vastaus jää olennaisesti auki, rahoitus on rakennettava ehdolliseksi ja vaiheistetuksi.", "Evidence Register", [0.7,5.2,6.5] )
    slide_table(prs, ctx, 28, "Puuttuva aineisto priorisoituu vaikutuksen, ei helppouden mukaan", "Aineistopyynnöt",
                ["Prioriteetti", "Aineisto", "Päätösvaikutus"],
                [["1", "Omistus, siirrot, rasitteet, vuosimaksut", "Olemassa oleva ja kontrolloitava oikeus"], ["2", "Tuotenäytteet, testit ja claim chartit", "Tekninen osuvuus"], ["3", "Relevantti myynti ja vastapuolidata", "Rojaltipohja"], ["4", "Sopimukset, saatavat ja maksut", "Velanhoitokyky"], ["5", "Tilinpäätös, velat, budjetti ja rahoitustarve", "Yhtiö- ja luottoriski"], ["6", "Riippumaton IP-/yritysarvonmääritys", "Vakuus ja hinta"]],
                "Julkinen paketti näyttää aukot; luottamuksellinen näyttö kuuluu pääsyrajattuun datahuoneeseen.", "Evidence Register", [1.2,5.0,5.4] )
    slide_process(prs, ctx, 29, "90 päivän ohjelma muuttaa aukot päätöskelpoisiksi kontrolleiksi", "Seuraavat vaiheet",
                  [("0–30 päivää", "Oikeusmatriisi, määräpäivät, omistus- ja rasiteselvitys sekä yhtiöaineiston indeksi."), ("31–60 päivää", "Tuotedossierit, claim chartit, priorisoitu markkinamyynti ja asiakasvalidointi."), ("61–90 päivää", "Kassavirtamalli, riippumaton arvonmääritys, downside ja term sheet -vaihtoehdot."), ("Päätösportti", "Hyväksy, rajaa, vaiheista tai hylkää rahoitus näkyvillä ehdoilla ja stop-kriteereillä.")],
                  "Paketti päivitetään samasta julkisesta lähdedatasta jokaisessa sivustojulkaisussa.", "Evidence Register; changelog" )
    closing_slide(prs, ctx, 30, "Pankkikelpoinen tarina on todennusketju — ei suurin mahdollinen markkinaluku",
                  [f"Vahva lähtökohta: dokumentoitu patenttiydin ja {v['de_country']}n virallista oikeusnäyttöä.", "Ratkaiseva avoin työ: kansalliset oikeudet, claim-mapped sales, kassavirta ja downside-arvo.", "Suositus: ehdollinen 90 päivän diligence, selkeät päätösportit ja pääsyrajattu datahuone.", "Tämä julkinen paketti on tarkistettava evidenssikartta, ei Pixan Oy:n virallinen kanta tai rahoitussuositus."], "Kaikki lähteet Evidence Registerissä" )
    if len(prs.slides) != 30:
        raise AssertionError(f"Large deck must have 30 slides, got {len(prs.slides)}")
    validate_finnish_deck_fact_tokens(prs, ctx, "large")
    save_presentation(prs, path)


def save_presentation(prs: Presentation, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    normalize_ooxml(path)


def normalize_ooxml(path: Path) -> None:
    """Rewrite an OOXML zip with stable order, metadata and timestamps."""
    with zipfile.ZipFile(path, "r") as source:
        entries = []
        for name in sorted(source.namelist()):
            payload = source.read(name)
            if name == "docProps/core.xml":
                # openpyxl replaces the modified property with wall-clock time
                # during save.  Canonicalise it to the already fixed created
                # timestamp so identical source data always yields identical
                # files and manifest hashes.
                core = payload.decode("utf-8")
                created = re.search(r"<dcterms:created\b[^>]*>([^<]+)</dcterms:created>", core)
                if created:
                    core = re.sub(
                        r"(<dcterms:modified\b[^>]*>)[^<]*(</dcterms:modified>)",
                        rf"\g<1>{created.group(1)}\g<2>",
                        core,
                    )
                    payload = core.encode("utf-8")
            entries.append((name, payload))
    fd, temp_name = tempfile.mkstemp(prefix=path.stem + "-", suffix=path.suffix, dir=path.parent)
    os.close(fd)
    temp_path = Path(temp_name)
    try:
        # Store canonical members without DEFLATE.  ZIP compression output can
        # vary with the system zlib version (macOS builder vs GitHub's Linux
        # runner), whereas ZIP_STORED makes the committed artifact reproducible
        # byte-for-byte across those environments.
        with zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_STORED) as target:
            for name, payload in entries:
                info = zipfile.ZipInfo(name, (1980, 1, 1, 0, 0, 0))
                info.compress_type = zipfile.ZIP_STORED
                info.create_system = 3
                info.external_attr = 0o600 << 16
                target.writestr(info, payload)
        temp_path.replace(path)
        path.chmod(0o644)
    finally:
        temp_path.unlink(missing_ok=True)


def build_workbook(ctx: dict[str, Any], rows: list[dict[str, str]], path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Evidence Register"
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = f"A1:I{len(rows) + 1}"

    for col, header in enumerate(REGISTER_HEADERS, start=1):
        cell = ws.cell(1, col, header)
        cell.fill = PatternFill("solid", fgColor=NAVY)
        cell.font = Font(name="Aptos", size=11, bold=True, color=WHITE)
        cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
    ws.row_dimensions[1].height = 34

    status_fill = {
        "Vahvistettu": "D9EDE7",
        "Tuettu": "DCEAF5",
        "Oletus": "FFF0CB",
        "Puuttuu": "F8DADA",
    }
    thin = Side(style="thin", color="D9E2E7")
    for ridx, item in enumerate(rows, start=2):
        for cidx, header in enumerate(REGISTER_HEADERS, start=1):
            cell = ws.cell(ridx, cidx, item[header])
            cell.font = Font(name="Aptos", size=10, color="000000")
            cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
            cell.border = Border(bottom=thin)
        ws.cell(ridx, 8).fill = PatternFill("solid", fgColor=status_fill[item["Luottamustaso"]])
        ws.cell(ridx, 8).font = Font(name="Aptos", size=10, bold=True, color=INK)
        ws.row_dimensions[ridx].height = 78

    widths = {"A": 44, "B": 22, "C": 54, "D": 46, "E": 15, "F": 42, "G": 42, "H": 16, "I": 54}
    for col, width in widths.items():
        ws.column_dimensions[col].width = width
    ws.sheet_view.zoomScale = 65
    ws.page_setup.orientation = "landscape"
    ws.page_setup.paperSize = ws.PAPERSIZE_A3
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.print_title_rows = "1:1"
    ws.page_margins.left = 0.2
    ws.page_margins.right = 0.2
    ws.page_margins.top = 0.35
    ws.page_margins.bottom = 0.35
    table = Table(displayName="EvidenceRegister", ref=f"A1:I{len(rows) + 1}")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=False)
    ws.add_table(table)
    validation = DataValidation(type="list", formula1='"Vahvistettu,Tuettu,Oletus,Puuttuu"', allow_blank=False)
    ws.add_data_validation(validation)
    validation.add(f"H2:H{len(rows) + 1}")

    summary = wb.create_sheet("Yhteenveto")
    summary.sheet_view.showGridLines = False
    summary.merge_cells("A1:H2")
    summary["A1"] = "Pixan · julkisen evidenssipaketin yhteenveto"
    summary["A1"].fill = PatternFill("solid", fgColor=NAVY)
    summary["A1"].font = Font(name="Aptos Display", size=24, bold=True, color=WHITE)
    summary["A1"].alignment = Alignment(vertical="center")
    meta = [
        ("Versio", ctx["release"]["version"]),
        ("Päivitetty", ctx["as_of"]),
        ("Rajaus", "Julkinen riippumaton kooste; ei Pixan Oy:n virallinen kanta eikä arvo-, laina-, sijoitus- tai oikeudellinen lausunto."),
        ("Rivit", len(rows)),
    ]
    for idx, (label, value) in enumerate(meta, start=4):
        summary.cell(idx, 1, label).font = Font(name="Aptos", size=11, bold=True, color=MUTED)
        summary.cell(idx, 2, value).font = Font(name="Aptos", size=11, color=INK)
        summary.cell(idx, 2).alignment = Alignment(wrap_text=True)
    summary["A10"] = "Näytön jakauma"
    summary["A10"].font = Font(name="Aptos", size=14, bold=True, color=NAVY)
    counts = Counter(item["Luottamustaso"] for item in rows)
    for idx, status in enumerate(("Vahvistettu", "Tuettu", "Oletus", "Puuttuu"), start=11):
        summary.cell(idx, 1, status)
        summary.cell(idx, 2, counts[status])
        summary.cell(idx, 1).fill = PatternFill("solid", fgColor=status_fill[status])
        summary.cell(idx, 1).font = Font(name="Aptos", size=11, bold=True, color=INK)
        summary.cell(idx, 2).font = Font(name="Aptos", size=11, color=INK)
    summary["A17"] = "Kolme vahvinta rahoitusperustetta"
    summary["A17"].font = Font(name="Aptos", size=14, bold=True, color=NAVY)
    strongest = [
        "EPO pysytti patentin muutettuna ja B2-julkaisu on virallisesti jäljitettävissä.",
        "Saksasta on viralliset mitätöinti- ja loukkausratkaisut tarkoin näkyvin rajauksin.",
        "Julkinen markkina-aineisto erottaa viralliset havainnot, proxyt, mallit ja puutteet toisistaan.",
    ]
    for idx, text in enumerate(strongest, start=18):
        summary.cell(idx, 1, idx - 17)
        summary.cell(idx, 2, text)
    summary["A23"] = "Pankkikelpoisuuden neljä korjausta"
    summary["A23"].font = Font(name="Aptos", size=14, bold=True, color=NAVY)
    fixes = [
        "Asiamiehen allekirjoittama oikeus-, omistus-, rasite- ja maksumatriisi.",
        "Priorisoitujen tuotteiden riippumattomat testit ja claim chartit.",
        "Toteutunut tai sopimuspohjainen kassavirta sekä auditoidut taloustiedot.",
        "Riippumaton arvonmääritys ja downside-vakuusanalyysi.",
    ]
    for idx, text in enumerate(fixes, start=24):
        summary.cell(idx, 1, idx - 23)
        summary.cell(idx, 2, text)
    summary.column_dimensions["A"].width = 28
    summary.column_dimensions["B"].width = 95
    summary.sheet_view.zoomScale = 90
    summary.page_setup.orientation = "landscape"
    summary.page_setup.fitToWidth = 1
    summary.page_setup.fitToHeight = 1
    summary.sheet_properties.pageSetUpPr.fitToPage = True
    for row in range(1, 30):
        summary.row_dimensions[row].height = 23
    summary.row_dimensions[1].height = 32
    summary.row_dimensions[2].height = 32

    questions = wb.create_sheet("Tutkijan kysymykset")
    questions.sheet_view.showGridLines = False
    q_headers = ["Prioriteetti", "Todennäköinen kysymys", "Vaadittu näyttö", "Nykytila"]
    q_rows = [
        (1, "Mitä tarkalleen omistetaan ja missä oikeus on voimassa?", "Maakohtainen oikeusmatriisi", "Puuttuu kattavasti"),
        (2, "Mikä tuote täyttää mitkä vaatimuspiirteet?", "Riippumaton testi ja claim chart", "Rajattu Saksan näyttö"),
        (3, "Mikä on todennettu relevantti myynti?", "Tuote–maa–aika-net sales", "Puuttuu"),
        (4, "Mistä ja milloin velanhoitokassa syntyy?", "Sopimukset, maksut ja ennuste", "Puuttuu"),
        (5, "Mitä vakuus realisoi downside-tilanteessa?", "Riippumaton arvio ja toteutuspolku", "Puuttuu"),
    ]
    for cidx, value in enumerate(q_headers, start=1):
        questions.cell(1, cidx, value).fill = PatternFill("solid", fgColor=NAVY)
        questions.cell(1, cidx).font = Font(name="Aptos", size=11, bold=True, color=WHITE)
    for ridx, values in enumerate(q_rows, start=2):
        for cidx, value in enumerate(values, start=1):
            questions.cell(ridx, cidx, value)
            questions.cell(ridx, cidx).alignment = Alignment(wrap_text=True, vertical="top")
            questions.cell(ridx, cidx).font = Font(name="Aptos", size=11, color=INK)
        questions.row_dimensions[ridx].height = 52
    for col, width in zip("ABCD", (14, 52, 50, 28)):
        questions.column_dimensions[col].width = width
    questions.freeze_panes = "A2"
    questions.sheet_view.zoomScale = 90
    questions.page_setup.orientation = "landscape"
    questions.page_setup.fitToWidth = 1
    questions.page_setup.fitToHeight = 1
    questions.sheet_properties.pageSetUpPr.fitToPage = True

    sources_ws = wb.create_sheet("Lähteet")
    sources_ws.sheet_view.showGridLines = False
    sources_ws.append(["Lähdetunnus", "Julkaisija", "Lähdeluokka", "URL", "Haettu / data-as-of"])
    all_sources = list(ctx["market"]["sources"]) + list(ctx["patent_history"]["sources"])
    for source in sorted(all_sources, key=lambda item: item["sourceId"]):
        sources_ws.append([
            source["sourceId"],
            source.get("publisher", ""),
            source.get("sourceKind", source.get("evidenceTier", "")),
            source.get("pageUrl") or source.get("url") or "",
            source.get("retrievedAt", ctx["as_of"]),
        ])
    for cell in sources_ws[1]:
        cell.fill = PatternFill("solid", fgColor=NAVY)
        cell.font = Font(name="Aptos", size=11, bold=True, color=WHITE)
    for row in sources_ws.iter_rows(min_row=2):
        for cell in row:
            cell.font = Font(name="Aptos", size=10, color=INK)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
    for col, width in zip("ABCDE", (30, 38, 25, 80, 20)):
        sources_ws.column_dimensions[col].width = width
    sources_ws.freeze_panes = "A2"
    sources_ws.auto_filter.ref = f"A1:E{sources_ws.max_row}"
    sources_ws.sheet_view.zoomScale = 75
    sources_ws.page_setup.orientation = "landscape"
    sources_ws.page_setup.paperSize = sources_ws.PAPERSIZE_A3
    sources_ws.page_setup.fitToWidth = 1
    sources_ws.page_setup.fitToHeight = 0
    sources_ws.sheet_properties.pageSetUpPr.fitToPage = True
    sources_ws.print_title_rows = "1:1"

    wb.properties.title = "Pixan Evidence Register"
    wb.properties.subject = "Julkinen pankki- ja teknologia-arvioinnin evidenssirekisteri"
    wb.properties.creator = "Pixan Global Market Evidence Atlas"
    stamp = parse_iso_date(ctx["as_of"]).replace(tzinfo=None)
    wb.properties.created = stamp
    wb.properties.modified = stamp
    wb.calculation.fullCalcOnLoad = True
    wb.calculation.forceFullCalc = True
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(path)
    normalize_ooxml(path)


def write_csv(rows: list[dict[str, str]], path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=REGISTER_HEADERS, lineterminator="\n")
        writer.writeheader()
        writer.writerows(rows)


def read_english_register_source() -> list[list[str]]:
    payload = read_json(EN_REGISTER_SOURCE)
    if set(payload) != {"headers", "rows"} or payload.get("headers") != EN_REGISTER_HEADERS:
        raise ValueError("English Evidence Register source has an unexpected schema")
    rows = payload.get("rows")
    if not isinstance(rows, list) or len(rows) != 45:
        raise ValueError("English Evidence Register source must contain exactly 45 rows")
    normalised: list[list[str]] = []
    for index, row in enumerate(rows, start=1):
        if not isinstance(row, list) or len(row) != len(EN_REGISTER_HEADERS):
            raise ValueError(f"English Evidence Register row {index} must contain nine cells")
        values = [str(value) for value in row]
        if not all(value.strip() for value in values):
            raise ValueError(f"English Evidence Register row {index} contains an empty cell")
        if values[7] not in EN_ALLOWED_STATUSES:
            raise ValueError(f"English Evidence Register row {index} has an invalid confidence status")
        normalised.append(values)
    if {row[7] for row in normalised} != EN_ALLOWED_STATUSES:
        raise ValueError("English Evidence Register must visibly use all four classifications")
    public_text_scan(value for row in normalised for value in row)
    return normalised


def write_english_csv(rows: list[list[str]], path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle, lineterminator="\n")
        writer.writerow(EN_REGISTER_HEADERS)
        writer.writerows(rows)


EN_MEDIUM_SECTION_TITLES = [
    "financing thesis",
    "problem",
    "patented solution",
    "patent and ip status",
    "technical differentiation",
    "market size and scope",
    "customers and purchase rationale",
    "competition and alternatives",
    "validation and current evidence",
    "commercialisation model",
    "financial model and sensitivities",
    "risks, controls and next steps",
]


def inspect_english_pptx(path: Path, expected_slides: int, *, require_medium_titles: bool = False) -> None:
    prs = Presentation(path)
    if len(prs.slides) != expected_slides:
        raise AssertionError(f"{path.name}: expected {expected_slides} slides, got {len(prs.slides)}")
    slide_text = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))
        for slide in prs.slides
    ]
    if any(not text.strip() for text in slide_text):
        raise AssertionError(f"{path.name}: every slide must contain readable text")
    combined = "\n".join(slide_text)
    public_text_scan([combined])
    if "independent public evidence" not in combined.casefold():
        raise AssertionError(f"{path.name}: English public-boundary disclosure missing")
    if require_medium_titles:
        for index, title in enumerate(EN_MEDIUM_SECTION_TITLES):
            if title not in " ".join(slide_text[index].casefold().split()):
                raise AssertionError(f"{path.name}: slide {index + 1} lacks section title {title!r}")


def inspect_english_xlsx(path: Path, expected_rows: list[list[str]]) -> None:
    workbook = load_workbook(path, read_only=False, data_only=False)
    if "Evidence Register" not in workbook.sheetnames:
        raise AssertionError(f"{path.name}: Evidence Register sheet missing")
    if any(sheet.sheet_state != "visible" for sheet in workbook.worksheets):
        raise AssertionError(f"{path.name}: hidden sheets are not allowed")
    sheet = workbook["Evidence Register"]
    headers = [str(sheet.cell(1, column).value or "") for column in range(1, 10)]
    if headers != EN_REGISTER_HEADERS:
        raise AssertionError(f"{path.name}: English Evidence Register headers changed")
    rows = []
    for values in sheet.iter_rows(min_row=2, max_col=9, values_only=True):
        row = ["" if value is None else str(value) for value in values]
        if any(value.strip() for value in row):
            rows.append(row)
    if rows != expected_rows:
        raise AssertionError(f"{path.name}: workbook rows differ from the reviewed English source")
    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.comment is not None:
                    raise AssertionError(f"{path.name}: comments are not allowed")
                if isinstance(cell.value, str):
                    public_text_scan([cell.value])
                    if cell.value.startswith("=") and "[" in cell.value:
                        raise AssertionError(f"{path.name}: external workbook formulas are not allowed")


def verify_english_release_lock(ctx: dict[str, Any], english_rows: list[list[str]]) -> dict[str, Any]:
    lock = read_json(EN_LOCK_SOURCE)
    expected_top = {
        "schemaVersion", "release", "asOf", "translationInputs",
        "sourceArtifacts", "artifacts", "generatedBy",
    }
    if set(lock) != expected_top or lock.get("schemaVersion") != 1:
        raise ValueError("English package lock has an unexpected schema")
    expected_release = {
        key: ctx["release"][key] for key in ("id", "version", "publishedAt")
    }
    if lock.get("release") != expected_release or lock.get("asOf") != ctx["as_of"]:
        raise ValueError("English package lock is stale for the current public release")

    expected_translation_paths = {
        str(EN_REGISTER_SOURCE.relative_to(ROOT)),
        str(EN_DECK_TRANSLATIONS_SOURCE.relative_to(ROOT)),
    }
    translation_inputs = lock.get("translationInputs")
    if not isinstance(translation_inputs, list):
        raise ValueError("English package lock translationInputs must be an array")
    translation_by_path = {
        item.get("path"): item for item in translation_inputs if isinstance(item, dict)
    }
    if set(translation_by_path) != expected_translation_paths or len(translation_by_path) != len(translation_inputs):
        raise ValueError("English package lock translation input allowlist differs")
    for relative, item in translation_by_path.items():
        if set(item) != {"path", "sha256"} or item.get("sha256") != sha256(ROOT / relative):
            raise ValueError(f"English package translation input hash differs: {relative}")

    source_specs = {
        "short-deck-fi": OUTPUTS["short-deck"],
        "medium-deck-fi": OUTPUTS["medium-deck"],
        "large-deck-fi": OUTPUTS["large-deck"],
        "evidence-register-fi": OUTPUTS["evidence-register"],
    }
    source_items = lock.get("sourceArtifacts")
    source_by_id = {item.get("id"): item for item in source_items or [] if isinstance(item, dict)}
    if set(source_by_id) != set(source_specs) or len(source_by_id) != len(source_items or []):
        raise ValueError("English package lock must bind all four Finnish source artifacts")
    for artifact_id, path in source_specs.items():
        item = source_by_id[artifact_id]
        expected_path = str(path.relative_to(ROOT))
        if set(item) != {"id", "path", "sha256"} or item.get("path") != expected_path or item.get("sha256") != sha256(path):
            raise ValueError(f"English package source artifact lock differs: {artifact_id}")

    artifact_specs = {
        "short-deck-en": (EN_OUTPUTS["short-deck-en"], "pptx", "slideCount", 6),
        "medium-deck-en": (EN_OUTPUTS["medium-deck-en"], "pptx", "slideCount", 12),
        "large-deck-en": (EN_OUTPUTS["large-deck-en"], "pptx", "slideCount", 30),
        "evidence-register-en": (EN_OUTPUTS["evidence-register-en"], "xlsx", "rowCount", len(english_rows)),
    }
    artifacts = lock.get("artifacts")
    artifact_by_id = {item.get("id"): item for item in artifacts or [] if isinstance(item, dict)}
    if set(artifact_by_id) != set(artifact_specs) or len(artifact_by_id) != len(artifacts or []):
        raise ValueError("English package lock must contain exactly four reviewed artifacts")
    for artifact_id, (path, kind, count_key, count) in artifact_specs.items():
        item = artifact_by_id[artifact_id]
        expected_keys = {"id", "kind", "path", "sha256", "bytes", count_key}
        if set(item) != expected_keys or item.get("kind") != kind:
            raise ValueError(f"English package lock schema differs: {artifact_id}")
        if item.get("path") != str(path.relative_to(ROOT)) or not path.is_file():
            raise ValueError(f"English package artifact is missing: {artifact_id}")
        if item.get("sha256") != sha256(path) or item.get("bytes") != path.stat().st_size or item.get(count_key) != count:
            raise ValueError(f"English package artifact lock differs: {artifact_id}")

    inspect_english_pptx(EN_OUTPUTS["short-deck-en"], 6)
    inspect_english_pptx(EN_OUTPUTS["medium-deck-en"], 12, require_medium_titles=True)
    inspect_english_pptx(EN_OUTPUTS["large-deck-en"], 30)
    inspect_english_xlsx(EN_OUTPUTS["evidence-register-en"], english_rows)
    for path in EN_OUTPUTS.values():
        inspect_reviewed_english_ooxml(path)
    return lock


def inspect_pptx(path: Path, expected_slides: int) -> None:
    prs = Presentation(path)
    if len(prs.slides) != expected_slides:
        raise AssertionError(f"{path.name}: expected {expected_slides} slides, got {len(prs.slides)}")
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text_frame"))
    public_text_scan([text])
    if "Julkinen riippumaton" not in text and "Julkinen riippumaton evidenssikooste" not in text:
        raise AssertionError(f"{path.name}: public boundary disclosure missing")


def inspect_xlsx(path: Path, expected_rows: int) -> None:
    wb = load_workbook(path, read_only=False, data_only=False)
    if wb.sheetnames != ["Evidence Register", "Yhteenveto", "Tutkijan kysymykset", "Lähteet"]:
        raise AssertionError(f"Unexpected workbook sheets: {wb.sheetnames}")
    ws = wb["Evidence Register"]
    headers = [ws.cell(1, col).value for col in range(1, 10)]
    if headers != REGISTER_HEADERS:
        raise AssertionError("Evidence Register headers changed")
    if ws.max_row - 1 != expected_rows:
        raise AssertionError(f"Expected {expected_rows} evidence rows, got {ws.max_row - 1}")
    statuses = {ws.cell(row, 8).value for row in range(2, ws.max_row + 1)}
    if not statuses.issubset(ALLOWED_STATUSES) or statuses != ALLOWED_STATUSES:
        raise AssertionError(f"Evidence statuses invalid or incomplete: {statuses}")
    public_text_scan(str(cell.value or "") for sheet in wb.worksheets for row in sheet.iter_rows() for cell in row)


def inspect_ooxml(path: Path) -> None:
    forbidden_parts = ("vbaproject", "oleobject", "externallink", "connections", "comments", "notesmaster", "notesslide")
    with zipfile.ZipFile(path) as package:
        names = [name.casefold() for name in package.namelist()]
        for name in names:
            if any(part in name for part in forbidden_parts):
                raise AssertionError(f"{path.name}: forbidden OOXML part {name}")
        payload = b"\n".join(package.read(name) for name in package.namelist() if name.endswith((".xml", ".rels")))
        decoded = payload.decode("utf-8", errors="ignore")
        public_text_scan([decoded])


def inspect_reviewed_english_ooxml(path: Path) -> None:
    forbidden_parts = ("vbaproject", "oleobject", "externallink", "connections", "comments", "embedding")
    with zipfile.ZipFile(path) as package:
        names = package.namelist()
        if not names or "[Content_Types].xml" not in names or len(names) != len(set(names)):
            raise AssertionError(f"{path.name}: invalid OOXML package")
        if any(name.startswith("/") or ".." in Path(name).parts for name in names):
            raise AssertionError(f"{path.name}: unsafe OOXML member path")
        payloads: list[str] = []
        for name in names:
            lowered = f"/{name}".casefold()
            if any(part in lowered for part in forbidden_parts):
                raise AssertionError(f"{path.name}: forbidden OOXML part {name}")
            if not name.endswith((".xml", ".rels")):
                continue
            payload = package.read(name).decode("utf-8", errors="replace")
            payloads.append(payload)
            if ("/notesslides/" in lowered or "/notesmasters/" in lowered) and name.endswith(".xml"):
                note_root = ET.fromstring(payload)
                note_text = [
                    str(element.text or "").strip()
                    for element in note_root.iter()
                    if element.tag.endswith("}t") and str(element.text or "").strip()
                ]
                if note_text:
                    raise AssertionError(f"{path.name}: notes part contains text ({name})")
        public_text_scan(payloads)


def artifact_entry(
    artifact_id: str,
    path: Path,
    *,
    kind: str,
    language: str,
    title_fi: str,
    title_en: str,
    slide_count: int | None = None,
    row_count: int | None = None,
) -> dict[str, Any]:
    entry: dict[str, Any] = {
        "id": artifact_id,
        "kind": kind,
        "language": language,
        "titleFi": title_fi,
        "titleEn": title_en,
        "fileName": path.name,
        "path": f"downloads/{path.name}",
        "sha256": sha256(path),
        "bytes": path.stat().st_size,
    }
    if slide_count is not None:
        entry["slideCount"] = slide_count
    if row_count is not None:
        entry["rowCount"] = row_count
    return entry


def write_manifest(
    ctx: dict[str, Any],
    rows: list[dict[str, str]],
    english_rows: list[list[str]],
) -> None:
    manifest_inputs = (*INPUT_FILES, EN_REGISTER_SOURCE, EN_DECK_TRANSLATIONS_SOURCE, EN_LOCK_SOURCE)
    manifest = {
        "schemaVersion": 2,
        "generatedFromPublicDataOnly": True,
        "release": {
            "id": ctx["release"]["id"],
            "version": ctx["release"]["version"],
            "publishedAt": ctx["release"]["publishedAt"],
        },
        "asOf": ctx["as_of"],
        "languages": ["en", "fi"],
        "publicBoundary": {
            "en": "Independent public evidence summary. Not Pixan Oy's official position; not an audit, valuation, legal opinion, investment recommendation or lending recommendation.",
            "fi": "Riippumaton julkinen evidenssikooste. Ei Pixan Oy:n virallinen kanta; ei tilintarkastus, arvonmääritys, oikeudellinen lausunto, sijoitussuositus tai lainasuositus.",
        },
        "inputs": [{"path": str(path.relative_to(ROOT)), "sha256": sha256(path)} for path in manifest_inputs],
        "artifacts": [
            artifact_entry("short-deck-en", EN_OUTPUTS["short-deck-en"], kind="pptx", language="en", title_fi="Suppea pankkidekki (englanti)", title_en="Concise bank deck (English)", slide_count=6),
            artifact_entry("medium-deck-en", EN_OUTPUTS["medium-deck-en"], kind="pptx", language="en", title_fi="Keskikokoinen pankkidekki (englanti)", title_en="Core bank deck (English)", slide_count=12),
            artifact_entry("large-deck-en", EN_OUTPUTS["large-deck-en"], kind="pptx", language="en", title_fi="Laaja pankkidekki (englanti)", title_en="Extended bank deck (English)", slide_count=30),
            artifact_entry("evidence-register-en", EN_OUTPUTS["evidence-register-en"], kind="xlsx", language="en", title_fi="Evidence Register (englanti)", title_en="Evidence Register (English)", row_count=len(english_rows)),
            artifact_entry("short-deck-fi", OUTPUTS["short-deck"], kind="pptx", language="fi", title_fi="Suppea pankkidekki (suomi)", title_en="Concise bank deck (Finnish)", slide_count=6),
            artifact_entry("medium-deck-fi", OUTPUTS["medium-deck"], kind="pptx", language="fi", title_fi="Keskikokoinen pankkidekki (suomi)", title_en="Core bank deck (Finnish)", slide_count=12),
            artifact_entry("large-deck-fi", OUTPUTS["large-deck"], kind="pptx", language="fi", title_fi="Laaja pankkidekki (suomi)", title_en="Extended bank deck (Finnish)", slide_count=30),
            artifact_entry("evidence-register-fi", OUTPUTS["evidence-register"], kind="xlsx", language="fi", title_fi="Evidence Register (suomi)", title_en="Evidence Register (Finnish)", row_count=len(rows)),
        ],
    }
    public_text_scan([json.dumps(manifest, ensure_ascii=False)])
    MANIFEST_OUTPUT.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def build_all() -> dict[str, Any]:
    ctx = build_context()
    regression_check_finnish_deck_fact_binding(ctx)
    rows = evidence_rows(ctx)
    english_rows = read_english_register_source()
    parity_errors = validate_register_parity(
        [[row[header] for header in REGISTER_HEADERS] for row in rows],
        english_rows,
    )
    if parity_errors:
        raise ValueError("Bilingual Evidence Register parity failed:\n- " + "\n- ".join(parity_errors))
    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    build_short_deck(ctx, OUTPUTS["short-deck"])
    build_medium_deck(ctx, OUTPUTS["medium-deck"])
    build_large_deck(ctx, OUTPUTS["large-deck"])
    build_workbook(ctx, rows, OUTPUTS["evidence-register"])
    write_csv(rows, CSV_OUTPUT)
    write_english_csv(english_rows, EN_CSV_OUTPUT)

    inspect_pptx(OUTPUTS["short-deck"], 6)
    inspect_pptx(OUTPUTS["medium-deck"], 12)
    inspect_pptx(OUTPUTS["large-deck"], 30)
    inspect_xlsx(OUTPUTS["evidence-register"], len(rows))
    for path in OUTPUTS.values():
        inspect_ooxml(path)
    verify_english_release_lock(ctx, english_rows)
    write_manifest(ctx, rows, english_rows)
    return {"version": ctx["release"]["version"], "asOf": ctx["as_of"], "evidenceRows": len(rows)}


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--check-determinism", action="store_true", help="Build twice and verify stable artifact hashes.")
    args = parser.parse_args()
    result = build_all()
    if args.check_determinism:
        first = {
            path.name: sha256(path)
            for path in (*OUTPUTS.values(), *EN_OUTPUTS.values(), CSV_OUTPUT, EN_CSV_OUTPUT, MANIFEST_OUTPUT)
        }
        result = build_all()
        second = {
            path.name: sha256(path)
            for path in (*OUTPUTS.values(), *EN_OUTPUTS.values(), CSV_OUTPUT, EN_CSV_OUTPUT, MANIFEST_OUTPUT)
        }
        if first != second:
            changed = sorted(name for name in first if first[name] != second[name])
            raise AssertionError(f"Non-deterministic outputs: {changed}")
    print(json.dumps(result, ensure_ascii=False, sort_keys=True))


if __name__ == "__main__":
    main()
